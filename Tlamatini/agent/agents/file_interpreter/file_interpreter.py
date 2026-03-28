# File-Interpreter Agent - Document parsing and text/image extraction
# Hybrid deterministic/non-deterministic agent
# Action: Read files -> Extract text (and optionally images) -> Log results
#         -> (summarized mode) Query LLM -> Start downstream agents

import os
import sys

# FIX: Disable Intel Fortran runtime Ctrl+C handler
os.environ['FOR_DISABLE_CONSOLE_CTRL_HANDLER'] = '1'

import glob
import time
import yaml
import json
import hashlib
import logging
import subprocess
import urllib.request
import urllib.error
from typing import Dict

# Set working directory to script location
try:
    script_dir = os.path.dirname(os.path.abspath(__file__))
    os.chdir(script_dir)
except Exception as e:
    sys.stderr.write(f"Critical Error: Failed to set working directory: {e}\n")

# Use directory name for log file
CURRENT_DIR_NAME = os.path.basename(os.path.dirname(os.path.abspath(__file__)))
LOG_FILE_PATH = f"{CURRENT_DIR_NAME}.log"

# Reanimation detection: AGENT_REANIMATED=1 means resume from pause
_IS_REANIMATED = os.environ.get('AGENT_REANIMATED') == '1'
if not _IS_REANIMATED:
    open(LOG_FILE_PATH, 'w').close()
logging.basicConfig(
    filename=LOG_FILE_PATH,
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    encoding='utf-8'
)

# Also log to console
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.INFO)
console_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
logging.getLogger().addHandler(console_handler)


def load_config(path: str = "config.yaml") -> Dict:
    try:
        with open(path, "r", encoding="utf-8") as f:
            return yaml.safe_load(f)
    except FileNotFoundError:
        logging.error(f"❌ Error: {path} not found.")
        sys.exit(1)
    except Exception as e:
        logging.error(f"❌ Error parsing {path}: {e}")
        sys.exit(1)


def get_python_command() -> list:
    if not getattr(sys, 'frozen', False):
        return [sys.executable]
    python_home = get_user_python_home()
    if python_home:
        python_exe = os.path.join(python_home, 'python.exe' if sys.platform.startswith('win') else 'python3')
        if os.path.exists(python_exe):
            return [python_exe]
    if sys.platform.startswith('win'):
        bundled_python = os.path.join(os.path.dirname(sys.executable), 'python.exe')
        if os.path.exists(bundled_python):
            return [bundled_python]
        return ['python']
    return ['python3']


def get_user_python_home() -> str:
    if not sys.platform.startswith('win'):
        return os.environ.get('PYTHON_HOME', '')
    try:
        import winreg
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, r'Environment') as key:
            value, _ = winreg.QueryValueEx(key, 'PYTHON_HOME')
            return str(value) if value else ''
    except (FileNotFoundError, OSError):
        return ''


def get_agent_env() -> dict:
    env = os.environ.copy()
    if sys.platform.startswith('win'):
        try:
            import ctypes
            if hasattr(ctypes.windll.kernel32, 'SetDllDirectoryW'):
                ctypes.windll.kernel32.SetDllDirectoryW(None)
        except Exception:
            pass
    if getattr(sys, 'frozen', False) and hasattr(sys, '_MEIPASS'):
        meipass = getattr(sys, '_MEIPASS')
        if meipass:
            path_parts = env.get('PATH', '').split(os.pathsep)
            path_parts = [p for p in path_parts if os.path.normpath(p) != os.path.normpath(meipass)]
            env['PATH'] = os.pathsep.join(path_parts)
    python_home = get_user_python_home()
    if not python_home:
        return env
    env['PYTHON_HOME'] = python_home
    scripts_dir = os.path.join(python_home, 'Scripts')
    current_path = env.get('PATH', '')
    env['PATH'] = f"{python_home};{scripts_dir};{current_path}"
    return env


def get_pool_path() -> str:
    current_dir = os.path.dirname(os.path.abspath(__file__))
    parent = os.path.dirname(current_dir)
    grandparent = os.path.dirname(parent)
    if os.path.basename(grandparent) == 'pools':
        return parent
    if os.path.basename(parent) == 'pools':
        return parent
    return os.path.join(os.path.dirname(current_dir), 'pools')


def get_agent_directory(agent_name: str) -> str:
    return os.path.join(get_pool_path(), agent_name)


def get_agent_script_path(agent_name: str) -> str:
    agent_dir = get_agent_directory(agent_name)
    if os.path.exists(os.path.join(agent_dir, f"{agent_name}.py")):
        return os.path.join(agent_dir, f"{agent_name}.py")
    parts = agent_name.rsplit('_', 1)
    if len(parts) == 2 and parts[1].isdigit():
        base = parts[0]
        if os.path.exists(os.path.join(agent_dir, f"{base}.py")):
            return os.path.join(agent_dir, f"{base}.py")
    return os.path.join(agent_dir, f"{agent_name}.py")


def is_agent_running(agent_name: str) -> bool:
    agent_dir = get_agent_directory(agent_name)
    pid_path = os.path.join(agent_dir, "agent.pid")
    if not os.path.exists(pid_path):
        return False
    try:
        with open(pid_path, "r") as f:
            pid = int(f.read().strip())
    except (ValueError, OSError):
        return False
    try:
        import psutil
        if not psutil.pid_exists(pid):
            return False
        proc = psutil.Process(pid)
        if proc.status() == psutil.STATUS_ZOMBIE:
            return False
        return True
    except Exception:
        try:
            os.kill(pid, 0)
            return True
        except OSError:
            return False


def wait_for_agents_to_stop(agent_names: list):
    if not agent_names:
        return
    waited = 0.0
    poll_interval = 0.5
    while True:
        still_running = [name for name in agent_names if is_agent_running(name)]
        if not still_running:
            return
        if waited >= 10.0:
            logging.error(
                f"❌ WAITING FOR AGENTS TO STOP: {still_running} still running "
                f"after {int(waited)}s. Will keep waiting..."
            )
            waited = 0.0
        time.sleep(poll_interval)
        waited += poll_interval


def start_agent(agent_name: str) -> bool:
    agent_dir = get_agent_directory(agent_name)
    script_path = get_agent_script_path(agent_name)
    if not os.path.exists(script_path):
        logging.error(f"❌ Agent script not found: {script_path}")
        return False
    try:
        cmd = get_python_command() + [script_path]
        logging.info(f"   Command: {cmd}")
        process = subprocess.Popen(
            cmd,
            cwd=agent_dir,
            env=get_agent_env(),
            creationflags=subprocess.CREATE_NO_WINDOW if sys.platform == 'win32' else 0
        )
        try:
            pid_path = os.path.join(agent_dir, "agent.pid")
            with open(pid_path, "w") as f:
                f.write(str(process.pid))
        except Exception as pid_err:
            logging.error(f"⚠️ Failed to write PID file for target {agent_name}: {pid_err}")
        logging.info(f"✅ Started agent '{agent_name}' with PID: {process.pid}")
        return True
    except Exception as e:
        logging.error(f"❌ Failed to start agent '{agent_name}': {e}")
        return False


# PID Management
PID_FILE = "agent.pid"


def write_pid_file():
    try:
        with open(PID_FILE, "w") as f:
            f.write(str(os.getpid()))
    except Exception as e:
        logging.error(f"❌ Failed to write PID file: {e}")


def remove_pid_file():
    for _attempt in range(5):
        try:
            if os.path.exists(PID_FILE):
                os.remove(PID_FILE)
            return
        except PermissionError:
            time.sleep(0.1)
        except Exception as e:
            logging.error(f"❌ Failed to remove PID file: {e}")
            return


# ============================================================
# LLM Query (for 'summarized' mode)
# ============================================================

def query_ollama(host: str, model: str, system_prompt: str, context: str) -> str:
    url = f"{host.rstrip('/')}/api/generate"
    full_prompt = f"{system_prompt}\n\n--- BEGIN DOCUMENT CONTENT ---\n{context}\n--- END DOCUMENT CONTENT ---"
    payload = json.dumps({
        "model": model,
        "prompt": full_prompt,
        "stream": False
    }).encode("utf-8")
    req = urllib.request.Request(
        url,
        data=payload,
        headers={"Content-Type": "application/json"},
        method="POST"
    )
    try:
        with urllib.request.urlopen(req, timeout=300) as resp:
            body = json.loads(resp.read().decode("utf-8"))
            return body.get("response", "")
    except urllib.error.HTTPError as e:
        error_body = e.read().decode("utf-8", errors="replace") if e.fp else ""
        raise RuntimeError(f"Ollama HTTP {e.code}: {error_body}") from e
    except urllib.error.URLError as e:
        raise RuntimeError(f"Cannot reach Ollama at {host}: {e.reason}") from e


# ============================================================
# File Resolution
# ============================================================

def resolve_files(path_filenames: str, recursive: bool = False) -> list:
    """
    Resolve path_filenames into a list of actual file paths.
    Supports wildcards (e.g. 'C:\\temp\\*.docx') or single file paths.
    When recursive=True, injects '**/' to scan subdirectories.
    """
    path_filenames = path_filenames.strip()
    if not path_filenames:
        logging.error("❌ path_filenames is empty. Please configure a file path or wildcard pattern.")
        return []

    # Check if it contains wildcard characters
    has_wildcard = any(c in path_filenames for c in ['*', '?'])

    if has_wildcard:
        pattern = path_filenames
        # When recursive, inject **/ before the filename portion if not already present
        if recursive and '**' not in pattern:
            parent = os.path.dirname(pattern)
            filename_part = os.path.basename(pattern)
            pattern = os.path.join(parent, '**', filename_part) if parent else os.path.join('**', filename_part)
            logging.info(f"🔄 Recursive mode: expanded pattern to '{pattern}'")
        # Use glob to expand wildcards
        matched_files = glob.glob(pattern, recursive=recursive)
        if not matched_files:
            logging.warning(f"⚠️ No files matched the pattern: {pattern}")
            return []
        # Filter out directories, keep only files
        matched_files = [f for f in matched_files if os.path.isfile(f)]
        logging.info(f"📂 Pattern '{pattern}' matched {len(matched_files)} file(s)")
        return sorted(matched_files)
    else:
        # Single file path or bare directory
        if os.path.isdir(path_filenames):
            if recursive:
                # Recursive scan of entire directory tree
                pattern = os.path.join(path_filenames, '**', '*')
                matched_files = glob.glob(pattern, recursive=True)
                matched_files = [f for f in matched_files if os.path.isfile(f)]
                logging.info(f"📂 Recursive scan of '{path_filenames}' found {len(matched_files)} file(s)")
                return sorted(matched_files)
            logging.error(
                f"❌ path_filenames '{path_filenames}' is a directory without wildcards. "
                f"Please add a wildcard pattern, e.g.: {path_filenames}\\*.docx"
            )
            return []
        if not os.path.isfile(path_filenames):
            logging.error(f"❌ File not found or not readable: {path_filenames}")
            return []
        return [path_filenames]


def parse_exclusions(filetype_exclusions: str) -> tuple:
    """
    Parse a comma-separated exclusions string into (excluded_extensions, excluded_filenames).
    Entries with a dot and no other path chars are treated as extensions (e.g. "exe" -> ".exe").
    Entries that look like filenames (e.g. "main.cpp", ".profile") go into excluded_filenames.
    """
    excluded_extensions = set()
    excluded_filenames = set()
    if not filetype_exclusions or not filetype_exclusions.strip():
        return excluded_extensions, excluded_filenames
    for entry in filetype_exclusions.split(','):
        entry = entry.strip()
        if not entry:
            continue
        # If it contains a dot and has chars after it, treat as filename (e.g. "main.cpp", ".profile")
        if '.' in entry and not entry.startswith('.'):
            excluded_filenames.add(entry.lower())
        elif entry.startswith('.') and len(entry) > 1 and '.' not in entry[1:]:
            # Dotfile like ".profile" — treat as filename
            excluded_filenames.add(entry.lower())
        else:
            # Bare extension like "exe", "msi" — normalize to ".exe"
            ext = entry.lower() if entry.startswith('.') else f".{entry.lower()}"
            excluded_extensions.add(ext)
    return excluded_extensions, excluded_filenames


def apply_exclusions(files: list, excluded_extensions: set, excluded_filenames: set) -> list:
    """Filter out files matching excluded extensions or filenames."""
    if not excluded_extensions and not excluded_filenames:
        return files
    original_count = len(files)
    filtered = []
    for f in files:
        basename = os.path.basename(f).lower()
        ext = os.path.splitext(f)[1].lower()
        if ext in excluded_extensions or basename in excluded_filenames:
            continue
        filtered.append(f)
    excluded_count = original_count - len(filtered)
    if excluded_count > 0:
        logging.info(f"🚫 Excluded {excluded_count} file(s) by filetype_exclusions filter")
    return filtered


# ============================================================
# Unique image filename helper
# ============================================================

def unique_image_name(images_dir: str, original_name: str, file_path: str, img_index: int) -> str:
    """
    Generate a collision-free image filename by hashing the source file path + index.
    """
    base, ext = os.path.splitext(original_name)
    if not ext:
        ext = '.png'
    # Hash: source file + original name + index for uniqueness
    hash_input = f"{file_path}:{original_name}:{img_index}"
    short_hash = hashlib.md5(hash_input.encode()).hexdigest()[:8]
    candidate = f"{base}_{short_hash}{ext}"
    return candidate


# ============================================================
# File Parsers
# ============================================================

def extract_txt(file_path: str):
    """Extract text from plain text files."""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return [{'type': 'text', 'content': f.read()}]


def extract_pdf(file_path: str, extract_images: bool = False):
    """Extract text (and optionally images) from PDF files."""
    elements = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                elements.append({'type': 'text', 'content': text})
            if extract_images:
                img_list = page.get_images(full=True)
                for img_index, img_info in enumerate(img_list):
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image.get("ext", "png")
                    image_name = f"page{page_num + 1}_img{img_index + 1}.{image_ext}"
                    elements.append({'type': 'image', 'name': image_name, 'data': image_bytes})
        doc.close()
        return elements
    except ImportError:
        pass

    # Fallback: PyPDF2 / pypdf (text only)
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                elements.append({'type': 'text', 'content': text})
        return elements
    except ImportError:
        pass

    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(file_path)
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                elements.append({'type': 'text', 'content': text})
        return elements
    except ImportError:
        pass

    logging.error("❌ No PDF library available (install PyMuPDF, pypdf, or PyPDF2)")
    return [{'type': 'text', 'content': f'[ERROR: No PDF library available to read {file_path}]'}]


def extract_docx(file_path: str, extract_images: bool = False):
    """Extract text (and optionally images) from Word DOCX files."""
    elements = []
    try:
        from docx import Document
        doc = Document(file_path)

        if extract_images:
            # Extract images from relationships
            img_index = 0
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    img_index += 1
                    image_part = rel.target_part
                    image_ext = os.path.splitext(image_part.partname)[-1] or '.png'
                    image_name = f"docx_img{img_index}{image_ext}"
                    elements.append({'type': 'image', 'name': image_name, 'data': image_part.blob})

        for para in doc.paragraphs:
            if para.text.strip():
                elements.append({'type': 'text', 'content': para.text})

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = '\t'.join(cell.text for cell in row.cells)
                if row_text.strip():
                    elements.append({'type': 'text', 'content': row_text})

        return elements
    except ImportError:
        logging.error("❌ python-docx not installed. Cannot read DOCX files.")
        return [{'type': 'text', 'content': f'[ERROR: python-docx not installed to read {file_path}]'}]


def extract_pptx(file_path: str, extract_images: bool = False):
    """Extract text (and optionally images) from PowerPoint PPTX files."""
    elements = []
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        img_index = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = '\t'.join(cell.text for cell in row.cells)
                        if row_text.strip():
                            slide_texts.append(row_text)

                if extract_images and shape.shape_type == 13:  # Picture
                    img_index += 1
                    image = shape.image
                    image_ext = image.content_type.split('/')[-1] if image.content_type else 'png'
                    image_name = f"slide{slide_num}_img{img_index}.{image_ext}"
                    elements.append({'type': 'image', 'name': image_name, 'data': image.blob})

            if slide_texts:
                elements.append({'type': 'text', 'content': f"[Slide {slide_num}]\n" + '\n'.join(slide_texts)})

        return elements
    except ImportError:
        logging.error("❌ python-pptx not installed. Cannot read PPTX files.")
        return [{'type': 'text', 'content': f'[ERROR: python-pptx not installed to read {file_path}]'}]


def extract_xlsx(file_path: str):
    """Extract text from Excel XLSX files."""
    elements = []
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_lines = [f"[Sheet: {sheet_name}]"]
            for row in ws.iter_rows(values_only=True):
                row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    sheet_lines.append(row_text)
            if len(sheet_lines) > 1:
                elements.append({'type': 'text', 'content': '\n'.join(sheet_lines)})
        wb.close()
        return elements
    except ImportError:
        pass

    # Fallback: xlrd for .xls
    try:
        import xlrd
        wb = xlrd.open_workbook(file_path)
        for sheet in wb.sheets():
            sheet_lines = [f"[Sheet: {sheet.name}]"]
            for row_idx in range(sheet.nrows):
                row_text = '\t'.join(str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols))
                if row_text.strip():
                    sheet_lines.append(row_text)
            if len(sheet_lines) > 1:
                elements.append({'type': 'text', 'content': '\n'.join(sheet_lines)})
        return elements
    except ImportError:
        pass

    logging.error("❌ No Excel library available (install openpyxl or xlrd)")
    return [{'type': 'text', 'content': f'[ERROR: No Excel library available to read {file_path}]'}]


def extract_tex(file_path: str):
    """Extract text from TeX/LaTeX files (plain text reading)."""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return [{'type': 'text', 'content': f.read()}]


def extract_csv(file_path: str):
    """Extract text from CSV files."""
    import csv
    elements = []
    with open(file_path, 'r', encoding='utf-8', errors='replace', newline='') as f:
        reader = csv.reader(f)
        lines = []
        for row in reader:
            lines.append('\t'.join(row))
        if lines:
            elements.append({'type': 'text', 'content': '\n'.join(lines)})
    return elements


def extract_rtf(file_path: str):
    """Extract text from RTF files."""
    try:
        from striprtf.striprtf import rtf_to_text
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            rtf_content = f.read()
        text = rtf_to_text(rtf_content)
        return [{'type': 'text', 'content': text}]
    except ImportError:
        # Fallback: read raw
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return [{'type': 'text', 'content': f.read()}]


def extract_html(file_path: str):
    """Extract text from HTML files."""
    try:
        from bs4 import BeautifulSoup
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        return [{'type': 'text', 'content': soup.get_text(separator='\n')}]
    except ImportError:
        # Fallback: read raw
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return [{'type': 'text', 'content': f.read()}]


def extract_xml(file_path: str):
    """Extract text from XML files."""
    import xml.etree.ElementTree as ET
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        texts = []

        def walk(elem):
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                texts.append(elem.tail.strip())
            for child in elem:
                walk(child)

        walk(root)
        return [{'type': 'text', 'content': '\n'.join(texts)}]
    except ET.ParseError:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return [{'type': 'text', 'content': f.read()}]


def extract_json_file(file_path: str):
    """Extract text from JSON files."""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        content = f.read()
    try:
        data = json.loads(content)
        return [{'type': 'text', 'content': json.dumps(data, indent=2, ensure_ascii=False)}]
    except json.JSONDecodeError:
        return [{'type': 'text', 'content': content}]


def extract_yaml_file(file_path: str):
    """Extract text from YAML files."""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return [{'type': 'text', 'content': f.read()}]


def extract_markdown(file_path: str):
    """Extract text from Markdown files."""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return [{'type': 'text', 'content': f.read()}]


def extract_odt(file_path: str):
    """Extract text from OpenDocument Text files."""
    try:
        from odf import text as odf_text
        from odf.opendocument import load
        doc = load(file_path)
        all_text = []
        for element in doc.getElementsByType(odf_text.P):
            content = ''
            for node in element.childNodes:
                if hasattr(node, 'data'):
                    content += node.data
                elif hasattr(node, '__str__'):
                    content += str(node)
            if content.strip():
                all_text.append(content)
        return [{'type': 'text', 'content': '\n'.join(all_text)}]
    except ImportError:
        logging.error("❌ odfpy not installed. Cannot read ODT files.")
        return [{'type': 'text', 'content': f'[ERROR: odfpy not installed to read {file_path}]'}]


def extract_epub(file_path: str):
    """Extract text from EPUB files."""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(file_path)
        all_text = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text = soup.get_text(separator='\n')
                if text.strip():
                    all_text.append(text)
        return [{'type': 'text', 'content': '\n'.join(all_text)}]
    except ImportError:
        logging.error("❌ ebooklib/beautifulsoup4 not installed. Cannot read EPUB files.")
        return [{'type': 'text', 'content': f'[ERROR: ebooklib not installed to read {file_path}]'}]


# Extension-to-extractor mapping
EXTRACTORS = {
    '.txt': extract_txt,
    '.log': extract_txt,
    '.md': extract_markdown,
    '.csv': extract_csv,
    '.tsv': extract_csv,
    '.json': extract_json_file,
    '.yaml': extract_yaml_file,
    '.yml': extract_yaml_file,
    '.xml': extract_xml,
    '.html': extract_html,
    '.htm': extract_html,
    '.rtf': extract_rtf,
    '.tex': extract_tex,
    '.latex': extract_tex,
    '.bib': extract_tex,
    '.odt': extract_odt,
    '.epub': extract_epub,
    '.ini': extract_txt,
    '.cfg': extract_txt,
    '.conf': extract_txt,
    '.properties': extract_txt,
    '.py': extract_txt,
    '.js': extract_txt,
    '.java': extract_txt,
    '.c': extract_txt,
    '.cpp': extract_txt,
    '.h': extract_txt,
    '.cs': extract_txt,
    '.rb': extract_txt,
    '.php': extract_txt,
    '.sql': extract_txt,
    '.sh': extract_txt,
    '.bat': extract_txt,
    '.ps1': extract_txt,
    '.r': extract_txt,
}

# Extractors that support extract_images parameter
EXTRACTORS_WITH_IMAGES = {
    '.pdf': extract_pdf,
    '.docx': extract_docx,
    '.pptx': extract_pptx,
}

# Extractors for data-only formats (no image support)
EXTRACTORS_DATA_ONLY = {
    '.xlsx': extract_xlsx,
    '.xls': extract_xlsx,
}


def get_extractor(file_path: str, extract_images: bool = False):
    """Get the appropriate extractor function for a file."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext in EXTRACTORS_WITH_IMAGES:
        return lambda fp: EXTRACTORS_WITH_IMAGES[ext](fp, extract_images=extract_images)
    if ext in EXTRACTORS_DATA_ONLY:
        return EXTRACTORS_DATA_ONLY[ext]
    if ext in EXTRACTORS:
        return EXTRACTORS[ext]

    # Default: try to read as text
    logging.warning(f"⚠️ Unknown file extension '{ext}' for {file_path}, attempting plain text read")
    return extract_txt


# ============================================================
# Core Processing
# ============================================================

def process_file_fast(file_path: str):
    """Fast mode: extract text only, no images."""
    extractor = get_extractor(file_path, extract_images=False)
    elements = extractor(file_path)
    text_parts = [e['content'] for e in elements if e['type'] == 'text']
    raw_text = '\n'.join(text_parts)

    logging.info(
        f"INI_FILE: [{file_path}] (fast)\n"
        f"{{\n"
        f"{raw_text}\n"
        f"}}\n"
        f"END_FILE"
    )


def process_file_complete(file_path: str, images_dir: str):
    """Complete mode: extract text + images in order."""
    extractor = get_extractor(file_path, extract_images=True)
    elements = extractor(file_path)

    os.makedirs(images_dir, exist_ok=True)

    output_parts = []
    img_index = 0
    for element in elements:
        if element['type'] == 'text':
            output_parts.append(element['content'])
        elif element['type'] == 'image':
            img_index += 1
            img_name = unique_image_name(images_dir, element['name'], file_path, img_index)
            img_path = os.path.join(images_dir, img_name)
            try:
                with open(img_path, 'wb') as img_f:
                    img_f.write(element['data'])
                output_parts.append(img_name)
                logging.info(f"   📷 Image saved: {img_path}")
            except Exception as e:
                logging.error(f"   ❌ Failed to save image {img_name}: {e}")
                output_parts.append(f"[ERROR saving image: {img_name}]")

    content = '\n'.join(output_parts)
    logging.info(
        f"INI_FILE: [{file_path}] (complete)\n"
        f"{{\n"
        f"{content}\n"
        f"}}\n"
        f"END_FILE"
    )


def process_file_summarized(file_path: str, host: str, model: str):
    """Summarized mode: extract text, send to LLM for summarization."""
    extractor = get_extractor(file_path, extract_images=False)
    elements = extractor(file_path)
    text_parts = [e['content'] for e in elements if e['type'] == 'text']
    raw_text = '\n'.join(text_parts)

    if not raw_text.strip():
        logging.warning(f"⚠️ No text extracted from {file_path}, skipping summarization")
        logging.info(
            f"INI_FILE: [{file_path}] (summarized)\n"
            f"{{\n"
            f"[No text content found in file]\n"
            f"}}\n"
            f"END_FILE"
        )
        return

    system_prompt = (
        "You are a document summarizer. You will receive the full text content of a document. "
        "Your task is to produce a clear, comprehensive summary of the document content. "
        "The summary must be no more than 1024 words. "
        "Focus on the key points, main ideas, and important details. "
        "Produce ONLY the summary text, nothing else."
    )

    try:
        logging.info(f"   🤖 Sending {len(raw_text)} chars to LLM for summarization...")
        summary = query_ollama(host, model, system_prompt, raw_text)
    except RuntimeError as e:
        logging.error(f"   ❌ LLM summarization failed for {file_path}: {e}")
        summary = f"[ERROR: LLM summarization failed: {e}]"

    logging.info(
        f"INI_FILE: [{file_path}] (summarized)\n"
        f"{{\n"
        f"{summary}\n"
        f"}}\n"
        f"END_FILE"
    )


# ============================================================
# Main
# ============================================================

def main():
    config = load_config()
    write_pid_file()
    if _IS_REANIMATED:
        logging.info(f"🔄 {CURRENT_DIR_NAME} REANIMATED (resuming from pause)")
        logging.info("=" * 60)

    try:
        path_filenames = config.get('path_filenames', '')
        reading_type = config.get('reading_type', 'fast').strip().lower()
        recursive = config.get('recursive', False)
        filetype_exclusions = config.get('filetype_exclusions', '')
        source_agents = config.get('source_agents', [])
        target_agents = config.get('target_agents', [])
        llm_config = config.get('llm', {})
        host = llm_config.get('host', 'http://localhost:11434')
        model = llm_config.get('model', 'llama3.1:8b')

        logging.info("📄 FILE-INTERPRETER AGENT STARTED")
        logging.info(f"📂 Path/pattern: {path_filenames}")
        logging.info(f"📋 Reading type: {reading_type}")
        logging.info(f"🔄 Recursive: {recursive}")
        if filetype_exclusions:
            logging.info(f"🚫 Exclusions: {filetype_exclusions}")
        logging.info(f"📥 Source agents: {source_agents}")
        logging.info(f"🎯 Target agents: {target_agents}")
        if reading_type == 'summarized':
            logging.info(f"🤖 LLM: {model} @ {host}")
        logging.info("=" * 60)

        # Validate reading_type
        if reading_type not in ('fast', 'complete', 'summarized'):
            logging.error(f"❌ Invalid reading_type: '{reading_type}'. Must be 'fast', 'complete', or 'summarized'.")
            return

        # Resolve files
        files = resolve_files(path_filenames, recursive=recursive)
        excl_exts, excl_names = parse_exclusions(filetype_exclusions)
        files = apply_exclusions(files, excl_exts, excl_names)
        if not files:
            logging.error("❌ No files to process. Agent stopping.")
            return

        logging.info(f"📄 Processing {len(files)} file(s)...")

        # Set up images directory for 'complete' mode
        images_dir = os.path.join(script_dir, 'images')

        # Process each file
        for i, file_path in enumerate(files, 1):
            logging.info(f"--- [{i}/{len(files)}] Processing: {file_path}")
            try:
                if reading_type == 'fast':
                    process_file_fast(file_path)
                elif reading_type == 'complete':
                    process_file_complete(file_path, images_dir)
                elif reading_type == 'summarized':
                    process_file_summarized(file_path, host, model)
            except Exception as e:
                logging.error(f"❌ Error processing {file_path}: {e}")
                logging.info(
                    f"INI_FILE: [{file_path}] ({reading_type})\n"
                    f"{{\n"
                    f"[ERROR: {e}]\n"
                    f"}}\n"
                    f"END_FILE"
                )

        logging.info(f"✅ All {len(files)} file(s) processed.")

        # Trigger downstream agents
        total_triggered = 0
        if target_agents:
            wait_for_agents_to_stop(target_agents)
            logging.info(f"🚀 Triggering {len(target_agents)} downstream agents...")
            for target in target_agents:
                if start_agent(target):
                    total_triggered += 1

        logging.info(f"🏁 File-Interpreter agent finished. Triggered {total_triggered}/{len(target_agents)} agents.")

    except Exception as e:
        logging.error(f"❌ File-Interpreter agent error: {e}")
    finally:
        time.sleep(0.4)
        remove_pid_file()

    sys.exit(0)


if __name__ == "__main__":
    main()
