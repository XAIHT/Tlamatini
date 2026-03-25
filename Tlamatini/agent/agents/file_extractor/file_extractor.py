# File-Extractor Agent - Extracts text content from files, triggers downstream agents, then stops
# Action: Triggered by upstream -> Read/extract file(s) -> Log content -> Trigger downstream (always) -> Exit

import os
import sys

# FIX: Disable Intel Fortran runtime Ctrl+C handler
os.environ['FOR_DISABLE_CONSOLE_CTRL_HANDLER'] = '1'

import csv
import glob
import json
import re
import time
import yaml
import logging
import subprocess
import xml.etree.ElementTree as ET

# Set working directory to script location
try:
    script_dir = os.path.dirname(os.path.abspath(__file__))
    os.chdir(script_dir)
except Exception as e:
    sys.stderr.write(f"Critical Error: Failed to set working directory: {e}\n")

# Use directory name for log file
CURRENT_DIR_NAME = os.path.basename(os.path.dirname(os.path.abspath(__file__)))
LOG_FILE_PATH = f"{CURRENT_DIR_NAME}.log"
logging.basicConfig(
    filename=LOG_FILE_PATH,
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    encoding='utf-8'
)

# Also log to console
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.INFO)
console_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
logging.getLogger().addHandler(console_handler)


def load_config(path: str = "config.yaml") -> dict:
    try:
        with open(path, "r", encoding="utf-8") as f:
            return yaml.safe_load(f)
    except FileNotFoundError:
        logging.error(f"❌ Error: {path} not found.")
        sys.exit(1)
    except Exception as e:
        logging.error(f"❌ Error parsing {path}: {e}")
        sys.exit(1)


def get_python_command() -> list:
    """
    Get the command to run a Python script.
    - In Dev: Use current sys.executable (handles venvs).
    - In Frozen (Windows): Check for bundled python.exe, else fallback to 'python'.
    - In Frozen (Unix): Fallback to 'python3'.
    """
    if not getattr(sys, 'frozen', False):
        return [sys.executable]

    # Prefer PYTHON_HOME from USER environment variables
    python_home = get_user_python_home()
    if python_home:
        python_exe = os.path.join(python_home, 'python.exe' if sys.platform.startswith('win') else 'python3')
        if os.path.exists(python_exe):
            return [python_exe]

    if sys.platform.startswith('win'):
        bundled_python = os.path.join(os.path.dirname(sys.executable), 'python.exe')
        if os.path.exists(bundled_python):
            return [bundled_python]
        return ['python']

    return ['python3']


def get_user_python_home() -> str:
    """Read PYTHON_HOME exclusively from USER environment variables (Windows registry)."""
    if not sys.platform.startswith('win'):
        return os.environ.get('PYTHON_HOME', '')
    try:
        import winreg
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, r'Environment') as key:
            value, _ = winreg.QueryValueEx(key, 'PYTHON_HOME')
            return str(value) if value else ''
    except (FileNotFoundError, OSError):
        return ''


def get_agent_env() -> dict:
    """Build environment for child processes with PYTHON_HOME from USER env vars on PATH."""
    env = os.environ.copy()

    # Reset PyInstaller's DLL search path alteration on Windows
    if sys.platform.startswith('win'):
        try:
            import ctypes
            if hasattr(ctypes.windll.kernel32, 'SetDllDirectoryW'):
                ctypes.windll.kernel32.SetDllDirectoryW(None)
        except Exception:
            pass

    # Remove PyInstaller's _MEIPASS from PATH to prevent DLL conflicts in child processes
    if getattr(sys, 'frozen', False) and hasattr(sys, '_MEIPASS'):
        meipass = getattr(sys, '_MEIPASS')
        if meipass:
            path_parts = env.get('PATH', '').split(os.pathsep)
            path_parts = [p for p in path_parts if os.path.normpath(p) != os.path.normpath(meipass)]
            env['PATH'] = os.pathsep.join(path_parts)

    python_home = get_user_python_home()
    if not python_home:
        return env

    env['PYTHON_HOME'] = python_home
    scripts_dir = os.path.join(python_home, 'Scripts')
    current_path = env.get('PATH', '')
    env['PATH'] = f"{python_home};{scripts_dir};{current_path}"
    return env


def get_pool_path() -> str:
    """Get the pool directory path where deployed agents reside."""
    current_dir = os.path.dirname(os.path.abspath(__file__))

    # Check if deployed in session: pools/<session_id>/<agent_dir>
    parent = os.path.dirname(current_dir)
    grandparent = os.path.dirname(parent)

    if os.path.basename(grandparent) == 'pools':
        return parent

    if os.path.basename(parent) == 'pools':
        return parent

    return os.path.join(os.path.dirname(current_dir), 'pools')


def get_agent_directory(agent_name: str) -> str:
    return os.path.join(get_pool_path(), agent_name)


def get_agent_script_path(agent_name: str) -> str:
    agent_dir = get_agent_directory(agent_name)
    if os.path.exists(os.path.join(agent_dir, f"{agent_name}.py")):
        return os.path.join(agent_dir, f"{agent_name}.py")

    parts = agent_name.rsplit('_', 1)
    if len(parts) == 2 and parts[1].isdigit():
        base = parts[0]
        if os.path.exists(os.path.join(agent_dir, f"{base}.py")):
            return os.path.join(agent_dir, f"{base}.py")

    return os.path.join(agent_dir, f"{agent_name}.py")


def is_agent_running(agent_name: str) -> bool:
    """Check if an agent is currently running by verifying its PID file and process."""
    agent_dir = get_agent_directory(agent_name)
    pid_path = os.path.join(agent_dir, "agent.pid")

    if not os.path.exists(pid_path):
        return False

    try:
        with open(pid_path, "r") as f:
            pid = int(f.read().strip())
    except (ValueError, OSError):
        return False

    try:
        import psutil
        if not psutil.pid_exists(pid):
            return False
        proc = psutil.Process(pid)
        if proc.status() == psutil.STATUS_ZOMBIE:
            return False
        return True
    except Exception:
        try:
            os.kill(pid, 0)
            return True
        except OSError:
            return False


def wait_for_agents_to_stop(agent_names: list):
    """
    Wait until ALL specified agents have stopped running.
    Logs ERROR every 10 seconds while waiting. Never proceeds until all have stopped.
    """
    if not agent_names:
        return

    waited = 0.0
    poll_interval = 0.5

    while True:
        still_running = [name for name in agent_names if is_agent_running(name)]
        if not still_running:
            return

        if waited >= 10.0:
            logging.error(
                f"❌ WAITING FOR AGENTS TO STOP: {still_running} still running "
                f"after {int(waited)}s. Will keep waiting..."
            )
            waited = 0.0

        time.sleep(poll_interval)
        waited += poll_interval


def start_agent(agent_name: str) -> bool:
    agent_dir = get_agent_directory(agent_name)
    script_path = get_agent_script_path(agent_name)

    if not os.path.exists(script_path):
        logging.error(f"❌ Agent script not found: {script_path}")
        return False

    try:
        cmd = get_python_command() + [script_path]
        logging.info(f"   Command: {cmd}")

        process = subprocess.Popen(
            cmd,
            cwd=agent_dir,
            env=get_agent_env(),
            creationflags=subprocess.CREATE_NO_WINDOW if sys.platform == 'win32' else 0
        )

        try:
            pid_path = os.path.join(agent_dir, "agent.pid")
            with open(pid_path, "w") as f:
                f.write(str(process.pid))
        except Exception as pid_err:
            logging.error(f"⚠️ Failed to write PID file for target {agent_name}: {pid_err}")

        logging.info(f"✅ Started agent '{agent_name}' with PID: {process.pid}")
        return True
    except Exception as e:
        logging.error(f"❌ Failed to start agent '{agent_name}': {e}")
        return False


# PID Management
PID_FILE = "agent.pid"


def write_pid_file():
    try:
        with open(PID_FILE, "w") as f:
            f.write(str(os.getpid()))
    except Exception as e:
        logging.error(f"❌ Failed to write PID file: {e}")


def remove_pid_file():
    for attempt in range(5):
        try:
            if os.path.exists(PID_FILE):
                os.remove(PID_FILE)
            return
        except PermissionError:
            time.sleep(0.1)
        except Exception as e:
            logging.error(f"❌ Failed to remove PID file: {e}")
            return


# ============================================================
# File Resolution (wildcard support via glob)
# ============================================================

def resolve_files(path_filenames: str, recursive: bool = False) -> list:
    """
    Resolve path_filenames into a list of actual file paths.
    Supports wildcards (e.g. 'C:\\temp\\*.docx') or single file paths.
    When recursive=True, injects '**/' to scan subdirectories.
    """
    path_filenames = path_filenames.strip()
    if not path_filenames:
        logging.error("❌ path_filenames is empty. Please configure a file path or wildcard pattern.")
        return []

    has_wildcard = any(c in path_filenames for c in ['*', '?'])

    if has_wildcard:
        pattern = path_filenames
        if recursive and '**' not in pattern:
            parent = os.path.dirname(pattern)
            filename_part = os.path.basename(pattern)
            pattern = os.path.join(parent, '**', filename_part) if parent else os.path.join('**', filename_part)
            logging.info(f"🔄 Recursive mode: expanded pattern to '{pattern}'")
        matched_files = glob.glob(pattern, recursive=recursive)
        if not matched_files:
            logging.warning(f"⚠️ No files matched the pattern: {pattern}")
            return []
        matched_files = [f for f in matched_files if os.path.isfile(f)]
        logging.info(f"📂 Pattern '{pattern}' matched {len(matched_files)} file(s)")
        return sorted(matched_files)
    else:
        if os.path.isdir(path_filenames):
            if recursive:
                pattern = os.path.join(path_filenames, '**', '*')
                matched_files = glob.glob(pattern, recursive=True)
                matched_files = [f for f in matched_files if os.path.isfile(f)]
                logging.info(f"📂 Recursive scan of '{path_filenames}' found {len(matched_files)} file(s)")
                return sorted(matched_files)
            logging.error(
                f"❌ path_filenames '{path_filenames}' is a directory without wildcards. "
                f"Please add a wildcard pattern, e.g.: {path_filenames}\\*.*"
            )
            return []
        if not os.path.isfile(path_filenames):
            logging.error(f"❌ File not found or not readable: {path_filenames}")
            return []
        return [path_filenames]


def parse_exclusions(filetype_exclusions: str) -> tuple:
    """
    Parse a comma-separated exclusions string into (excluded_extensions, excluded_filenames).
    Entries with a dot and no other path chars are treated as extensions (e.g. "exe" -> ".exe").
    Entries that look like filenames (e.g. "main.cpp", ".profile") go into excluded_filenames.
    """
    excluded_extensions = set()
    excluded_filenames = set()
    if not filetype_exclusions or not filetype_exclusions.strip():
        return excluded_extensions, excluded_filenames
    for entry in filetype_exclusions.split(','):
        entry = entry.strip()
        if not entry:
            continue
        if '.' in entry and not entry.startswith('.'):
            excluded_filenames.add(entry.lower())
        elif entry.startswith('.') and len(entry) > 1 and '.' not in entry[1:]:
            excluded_filenames.add(entry.lower())
        else:
            ext = entry.lower() if entry.startswith('.') else f".{entry.lower()}"
            excluded_extensions.add(ext)
    return excluded_extensions, excluded_filenames


def apply_exclusions(files: list, excluded_extensions: set, excluded_filenames: set) -> list:
    """Filter out files matching excluded extensions or filenames."""
    if not excluded_extensions and not excluded_filenames:
        return files
    original_count = len(files)
    filtered = []
    for f in files:
        basename = os.path.basename(f).lower()
        ext = os.path.splitext(f)[1].lower()
        if ext in excluded_extensions or basename in excluded_filenames:
            continue
        filtered.append(f)
    excluded_count = original_count - len(filtered)
    if excluded_count > 0:
        logging.info(f"🚫 Excluded {excluded_count} file(s) by filetype_exclusions filter")
    return filtered


# ============================================================
# Text Extractors (text-only, no images)
# ============================================================

def extract_txt(file_path: str) -> str:
    """Extract text from plain text files."""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return f.read()


def extract_csv(file_path: str) -> str:
    """Extract text from CSV/TSV files."""
    lines = []
    with open(file_path, 'r', encoding='utf-8', errors='replace', newline='') as f:
        reader = csv.reader(f)
        for row in reader:
            lines.append('\t'.join(row))
    return '\n'.join(lines)


def extract_json_file(file_path: str) -> str:
    """Extract text from JSON files."""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        content = f.read()
    try:
        data = json.loads(content)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return content


def extract_yaml_file(file_path: str) -> str:
    """Extract text from YAML files."""
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return f.read()


def extract_xml(file_path: str) -> str:
    """Extract text from XML files."""
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        texts = []

        def walk(elem):
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                texts.append(elem.tail.strip())
            for child in elem:
                walk(child)

        walk(root)
        return '\n'.join(texts)
    except ET.ParseError:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()


def extract_html(file_path: str) -> str:
    """Extract text from HTML files."""
    try:
        from bs4 import BeautifulSoup
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        return soup.get_text(separator='\n')
    except ImportError:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()


def extract_rtf(file_path: str) -> str:
    """Extract text from RTF files."""
    try:
        from striprtf.striprtf import rtf_to_text
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            rtf_content = f.read()
        return rtf_to_text(rtf_content)
    except ImportError:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()


def extract_pdf(file_path: str) -> str:
    """Extract text from PDF files (text only, no images)."""
    texts = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                texts.append(text)
        doc.close()
        return '\n'.join(texts)
    except ImportError:
        pass

    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                texts.append(text)
        return '\n'.join(texts)
    except ImportError:
        pass

    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(file_path)
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                texts.append(text)
        return '\n'.join(texts)
    except ImportError:
        pass

    logging.error("❌ No PDF library available (install PyMuPDF, pypdf, or PyPDF2)")
    return f'[ERROR: No PDF library available to read {file_path}]'


def extract_docx(file_path: str) -> str:
    """Extract text from Word DOCX files (text only, no images)."""
    try:
        from docx import Document
        doc = Document(file_path)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = '\t'.join(cell.text for cell in row.cells)
                if row_text.strip():
                    texts.append(row_text)
        return '\n'.join(texts)
    except ImportError:
        logging.error("❌ python-docx not installed. Cannot read DOCX files.")
        return f'[ERROR: python-docx not installed to read {file_path}]'


def extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint PPTX files (text only, no images)."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = '\t'.join(cell.text for cell in row.cells)
                        if row_text.strip():
                            slide_texts.append(row_text)
            if slide_texts:
                texts.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_texts))
        return '\n'.join(texts)
    except ImportError:
        logging.error("❌ python-pptx not installed. Cannot read PPTX files.")
        return f'[ERROR: python-pptx not installed to read {file_path}]'


def extract_xlsx(file_path: str) -> str:
    """Extract text from Excel XLSX/XLS files."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_lines = [f"[Sheet: {sheet_name}]"]
            for row in ws.iter_rows(values_only=True):
                row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    sheet_lines.append(row_text)
            if len(sheet_lines) > 1:
                texts.append('\n'.join(sheet_lines))
        wb.close()
        return '\n'.join(texts)
    except ImportError:
        pass

    try:
        import xlrd
        wb = xlrd.open_workbook(file_path)
        texts = []
        for sheet in wb.sheets():
            sheet_lines = [f"[Sheet: {sheet.name}]"]
            for row_idx in range(sheet.nrows):
                row_text = '\t'.join(str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols))
                if row_text.strip():
                    sheet_lines.append(row_text)
            if len(sheet_lines) > 1:
                texts.append('\n'.join(sheet_lines))
        return '\n'.join(texts)
    except ImportError:
        pass

    logging.error("❌ No Excel library available (install openpyxl or xlrd)")
    return f'[ERROR: No Excel library available to read {file_path}]'


def extract_odt(file_path: str) -> str:
    """Extract text from OpenDocument Text files."""
    try:
        from odf import text as odf_text
        from odf.opendocument import load
        doc = load(file_path)
        all_text = []
        for element in doc.getElementsByType(odf_text.P):
            content = ''
            for node in element.childNodes:
                if hasattr(node, 'data'):
                    content += node.data
                elif hasattr(node, '__str__'):
                    content += str(node)
            if content.strip():
                all_text.append(content)
        return '\n'.join(all_text)
    except ImportError:
        logging.error("❌ odfpy not installed. Cannot read ODT files.")
        return f'[ERROR: odfpy not installed to read {file_path}]'


def extract_epub(file_path: str) -> str:
    """Extract text from EPUB files."""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(file_path)
        all_text = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text = soup.get_text(separator='\n')
                if text.strip():
                    all_text.append(text)
        return '\n'.join(all_text)
    except ImportError:
        logging.error("❌ ebooklib/beautifulsoup4 not installed. Cannot read EPUB files.")
        return f'[ERROR: ebooklib not installed to read {file_path}]'


def extract_strings(file_path: str, min_length: int = 4) -> str:
    """
    Extract printable strings from binary files, similar to the Linux 'strings' command.
    Finds sequences of printable ASCII characters of at least min_length.
    """
    printable_pattern = re.compile(
        rb'[\x20-\x7E\t\r\n]{' + str(min_length).encode() + rb',}'
    )
    try:
        with open(file_path, 'rb') as f:
            data = f.read()
        matches = printable_pattern.findall(data)
        return '\n'.join(m.decode('ascii', errors='replace') for m in matches)
    except Exception as e:
        logging.error(f"❌ Failed to extract strings from {file_path}: {e}")
        return f'[ERROR: Failed to extract strings from {file_path}: {e}]'


# Extension-to-extractor mapping (text-only, matching File-Interpreter)
EXTRACTORS = {
    '.txt': extract_txt,
    '.log': extract_txt,
    '.md': extract_txt,
    '.csv': extract_csv,
    '.tsv': extract_csv,
    '.json': extract_json_file,
    '.yaml': extract_yaml_file,
    '.yml': extract_yaml_file,
    '.xml': extract_xml,
    '.html': extract_html,
    '.htm': extract_html,
    '.rtf': extract_rtf,
    '.tex': extract_txt,
    '.latex': extract_txt,
    '.bib': extract_txt,
    '.odt': extract_odt,
    '.epub': extract_epub,
    '.ini': extract_txt,
    '.cfg': extract_txt,
    '.conf': extract_txt,
    '.properties': extract_txt,
    '.py': extract_txt,
    '.js': extract_txt,
    '.java': extract_txt,
    '.c': extract_txt,
    '.cpp': extract_txt,
    '.h': extract_txt,
    '.cs': extract_txt,
    '.rb': extract_txt,
    '.php': extract_txt,
    '.sql': extract_txt,
    '.sh': extract_txt,
    '.bat': extract_txt,
    '.ps1': extract_txt,
    '.r': extract_txt,
    '.pdf': extract_pdf,
    '.docx': extract_docx,
    '.pptx': extract_pptx,
    '.xlsx': extract_xlsx,
    '.xls': extract_xlsx,
}


def extract_file(file_path: str) -> str:
    """Extract text from a file using the appropriate extractor, or strings for unknown types."""
    ext = os.path.splitext(file_path)[1].lower()
    extractor = EXTRACTORS.get(ext)
    if extractor:
        return extractor(file_path)
    logging.info(f"⚠️ Unknown file type '{ext}' for {file_path} — using strings extraction")
    return extract_strings(file_path)


def main():
    config = load_config()

    # Write PID file immediately
    write_pid_file()

    try:
        path_filenames = config.get('path_filenames', '')
        recursive = config.get('recursive', False)
        filetype_exclusions = config.get('filetype_exclusions', '')
        target_agents = config.get('target_agents', [])

        logging.info("📦 FILE-EXTRACTOR AGENT STARTED")
        logging.info(f"📄 Path/pattern: {path_filenames}")
        logging.info(f"🔄 Recursive: {recursive}")
        if filetype_exclusions:
            logging.info(f"🚫 Exclusions: {filetype_exclusions}")
        logging.info(f"🎯 Targets: {target_agents}")

        # Resolve files (supports wildcards)
        files = resolve_files(path_filenames, recursive=recursive)
        excl_exts, excl_names = parse_exclusions(filetype_exclusions)
        files = apply_exclusions(files, excl_exts, excl_names)
        files_extracted = 0
        files_failed = 0

        for file_path in files:
            try:
                logging.info(f"📖 Extracting: {file_path}")
                raw_text = extract_file(file_path)
                logging.info(
                    f"INI_FILE: [{file_path}] (extracted)\n"
                    f"{{\n"
                    f"{raw_text}\n"
                    f"}}\n"
                    f"END_FILE"
                )
                files_extracted += 1
            except Exception as e:
                logging.error(f"❌ Failed to extract {file_path}: {e}")
                files_failed += 1

        logging.info(
            f"📊 Extraction complete: {files_extracted} succeeded, "
            f"{files_failed} failed out of {len(files)} file(s)"
        )

        # Trigger downstream agents independently of extraction result
        total_triggered = 0
        if target_agents:
            wait_for_agents_to_stop(target_agents)
            logging.info(f"🚀 Triggering {len(target_agents)} downstream agents...")
            for target in target_agents:
                if start_agent(target):
                    total_triggered += 1

        logging.info(
            f"🏁 File-Extractor agent finished. "
            f"Extracted {files_extracted}/{len(files)} files. "
            f"Triggered {total_triggered}/{len(target_agents)} agents."
        )

    finally:
        # Keep LED green briefly for visual feedback
        time.sleep(0.4)
        remove_pid_file()

    sys.exit(0)


if __name__ == "__main__":
    main()
