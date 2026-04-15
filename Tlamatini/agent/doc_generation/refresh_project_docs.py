from __future__ import annotations

import json
import subprocess
import sys
import textwrap
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


SCRIPT_PATH = Path(__file__).resolve()
DOC_DIR = SCRIPT_PATH.parent
PROJECT_DIR = SCRIPT_PATH.parents[2]
REPO_ROOT = SCRIPT_PATH.parents[3]

if str(DOC_DIR) not in sys.path:
    sys.path.insert(0, str(DOC_DIR))

from mardown_to_pdf import markdown_text_to_pdf


PDF_OUTPUT = REPO_ROOT / "tlamatini_app_summary.pdf"
PPT_OUTPUT = REPO_ROOT / "Tlamatini_eXtended_Artificial_Intelligence_Humanly_Tempered.pptx"
BUILD_DIR = REPO_ROOT / "build" / "documentation_refresh"
MARKDOWN_OUTPUT = BUILD_DIR / "tlamatini_app_summary_source.md"
JSON_OUTPUT = BUILD_DIR / "documentation_refresh_context.json"

DOCUMENT_PATHS = [
    "tlamatini_app_summary.pdf",
    "Tlamatini_eXtended_Artificial_Intelligence_Humanly_Tempered.pptx",
]

PDF_CSS = """
@page { size: A4; margin: 16mm; }
body {
  font-family: Helvetica, Arial, sans-serif;
  font-size: 10.5pt;
  line-height: 1.35;
  color: #172031;
}
h1 {
  font-size: 24pt;
  color: #0f172a;
  border-bottom: 2px solid #0f766e;
  padding-bottom: 6px;
  margin: 0 0 10px;
}
h2 {
  font-size: 16pt;
  color: #115e59;
  margin: 18px 0 8px;
}
h3 {
  font-size: 12.5pt;
  color: #1f2937;
  margin: 14px 0 6px;
}
p {
  margin: 5px 0;
}
ul, ol {
  margin: 6px 0 6px 18px;
}
li {
  margin: 3px 0;
}
code {
  font-family: Courier, monospace;
  font-size: 9pt;
  background-color: #f1f5f9;
}
pre {
  font-family: Courier, monospace;
  font-size: 9pt;
  border: 1px solid #cbd5e1;
  background-color: #f8fafc;
  padding: 8px;
  white-space: pre-wrap;
}
table {
  width: 100%;
  border-collapse: collapse;
  margin: 8px 0 12px;
}
th, td {
  border: 1px solid #cbd5e1;
  padding: 6px 7px;
  vertical-align: top;
}
th {
  background-color: #dcfce7;
}
blockquote {
  border-left: 4px solid #0f766e;
  padding-left: 10px;
  color: #334155;
  margin-left: 0;
}
.small {
  font-size: 9pt;
  color: #475569;
}
"""


PALETTE = {
    "ink": RGBColor(15, 23, 42),
    "slate": RGBColor(51, 65, 85),
    "teal": RGBColor(15, 118, 110),
    "teal_light": RGBColor(204, 251, 241),
    "gold": RGBColor(217, 119, 6),
    "gold_light": RGBColor(254, 243, 199),
    "red": RGBColor(190, 24, 93),
    "red_light": RGBColor(252, 231, 243),
    "sky": RGBColor(14, 165, 233),
    "sky_light": RGBColor(224, 242, 254),
    "paper": RGBColor(248, 250, 252),
    "panel": RGBColor(255, 255, 255),
    "line": RGBColor(203, 213, 225),
    "green": RGBColor(22, 163, 74),
    "green_light": RGBColor(220, 252, 231),
}


@dataclass
class CommitInfo:
    full_hash: str
    short_hash: str
    committed_at: str
    subject: str

    @property
    def pretty_date(self) -> str:
        return iso_to_human(self.committed_at)


@dataclass
class FileDelta:
    path: str
    insertions: int
    deletions: int

    @property
    def churn(self) -> int:
        return self.insertions + self.deletions


@dataclass
class ChangeTheme:
    title: str
    summary: str
    bullets: list[str]
    files: list[str]


THEMES = [
    ChangeTheme(
        title="Multi-Turn loop hardening",
        summary=(
            "The checked Multi-Turn path is now more stable, more explicit, "
            "and more cleanly separated from the original one-shot chat path."
        ),
        bullets=[
            "April 12 fixes targeted the breaking loop and the frozen-route path so the checked executor no longer stalls as easily.",
            "April 13 refined the stopping scheme so Multi-Turn sessions can wind down without leaving the runtime in a confusing state.",
            "The planner plus capability registry continue to shape request-scoped context prefetch, execution, monitoring, and answer stages.",
            "Unchecked chat still keeps the legacy validation and broad-tool behavior, which preserves backward compatibility.",
        ],
        files=[
            "Tlamatini/agent/mcp_agent.py",
            "Tlamatini/agent/global_execution_planner.py",
            "Tlamatini/agent/capability_registry.py",
            "Tlamatini/agent/rag/interface.py",
            "Tlamatini/agent/rag/chains/unified.py",
        ],
    ),
    ChangeTheme(
        title="Wrapped runtime permanence and observability",
        summary=(
            "Wrapped chat-agent launches now preserve each attempt, expose more "
            "runtime logging, and behave more predictably across source and frozen builds."
        ),
        bullets=[
            "Each wrapped launch gets a unique sequenced directory under `agent/agents/pools/_chat_runs_/` so failed attempts are not overwritten.",
            "Path discovery now distinguishes source mode from frozen mode and logs absolute resolution decisions for easier debugging.",
            "Visible console windows are suppressed when checked Multi-Turn launches detached helper processes.",
            "Django logging now exposes INFO-level runtime, planner, and tool-selection messages in the console.",
        ],
        files=[
            "Tlamatini/agent/chat_agent_runtime.py",
            "Tlamatini/agent/tools.py",
            "Tlamatini/tlamatini/settings.py",
            "Tlamatini/agent/agents/notifier/config.yaml",
            "Tlamatini/agent/agents/file_extractor/config.yaml",
        ],
    ),
    ChangeTheme(
        title="Flow creation from successful answers",
        summary=(
            "A successful Multi-Turn answer can now become a starter workflow file "
            "instead of ending as plain chat text."
        ),
        bullets=[
            "Backend tool calls now carry success state plus ACP agent display-name mapping for flow generation.",
            "The frontend only renders `Create Flow` when the reply used Multi-Turn, at least one eligible tool succeeded, and the answer was classified as successful.",
            "Generated `.flw` files are built as a simple `Starter -> Agent(s) -> Ender` chain, with tool arguments mapped into initial node configs.",
            "Management tools like run-status, run-log, and stop helpers are intentionally excluded from generated flow nodes.",
        ],
        files=[
            "Tlamatini/agent/mcp_agent.py",
            "Tlamatini/agent/services/response_parser.py",
            "Tlamatini/agent/static/agent/js/agent_page_chat.js",
            "Tlamatini/agent/static/agent/css/agent_page.css",
            "multi-turn-flow01.flw",
        ],
    ),
    ChangeTheme(
        title="Answer Analizer based flow gating",
        summary=(
            "The new answer classifier replaces brittle keyword heuristics with a "
            "second narrow LLM check that decides whether the task actually succeeded."
        ),
        bullets=[
            "A dedicated `answer_analizer.py` module sends the final answer through a strict `SUCCESS` or `FAILURE` classification prompt.",
            "Partial but useful completion is treated as `SUCCESS`, while refusals and unrecoverable failures are treated as `FAILURE`.",
            "On internal classifier failure the UI defaults to showing the flow button, which favors recoverability over silent hiding.",
            "The WebSocket payload now includes `answer_success` so the frontend can gate the button without inventing its own parser.",
        ],
        files=[
            "Tlamatini/agent/services/answer_analizer.py",
            "Tlamatini/agent/services/response_parser.py",
            "Tlamatini/agent/consumers.py",
            "Tlamatini/agent/static/agent/js/agent_page_chat.js",
        ],
    ),
    ChangeTheme(
        title="Documentation expansion plus Gatewayer deep dive",
        summary=(
            "The operator docs are materially richer now, especially around "
            "Multi-Turn answer-to-flow generation and Gatewayer's ingress model."
        ),
        bullets=[
            "README now includes a dedicated `Flow Creation from Multi-Turn Answers` section that walks through backend logging, frontend gating, and `.flw` output.",
            "A new `gatewayer_explanation.md` file documents HTTP ingress, folder-drop mode, authentication, crash recovery, and observability.",
            "The Gatewayer docs now explain where the implementation is complete versus where config keys are declared but not yet enforced.",
            "The refreshed docs also clarify that Gatewayer is a workflow ingress boundary, not a chat-router equivalent.",
        ],
        files=[
            "README.md",
            "Tlamatini/agent/agents/gatewayer/gatewayer_explanation.md",
            "Tlamatini/agent/agents/flowcreator/agentic_skill.md",
        ],
    ),
    ChangeTheme(
        title="Dependency and workflow support updates",
        summary=(
            "The Python environment gained spreadsheet support and the docs now better "
            "reflect current operator guidance."
        ),
        bullets=[
            "`openpyxl` was added to `requirements.txt`, expanding the bundled Python toolchain for spreadsheet-oriented tasks.",
            "The README updates on April 13 and April 14 align the public documentation with the latest Multi-Turn behavior.",
            "A temporary example `.flw` file was added during feature work and then removed, leaving the flow-generation feature in product code rather than shipping a scratch artifact.",
        ],
        files=[
            "requirements.txt",
            "README.md",
            "multi-turn-flow01.flw",
        ],
    ),
]


COMMIT_FOCUS = {
    "174e6a0": "Fix Multi-Turn loop stability.",
    "d66d749": "Fix frozen-mode routes for Multi-Turn.",
    "2a5965d": "Tighten Multi-Turn stopping behavior.",
    "6316d91": "Preserve wrapped runtime history and expose more logs.",
    "325cc99": "Expand operator-facing docs.",
    "892670c": "Add openpyxl dependency.",
    "78d40ff": "Finalize Python recommendation documentation/testing.",
    "156287b": "Turn successful answers into downloadable `.flw` workflows.",
    "060db0a": "Remove temporary pivot flow artifact.",
    "8b40669": "Add answer-success classification and README coverage.",
}


def git(*args: str) -> str:
    result = subprocess.run(
        ["git", *args],
        cwd=REPO_ROOT,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        check=True,
    )
    return result.stdout.strip()


def iso_to_human(iso_value: str) -> str:
    dt = datetime.fromisoformat(iso_value.replace("Z", "+00:00"))
    return dt.strftime("%B %d, %Y %I:%M %p UTC%z")


def plain_date(iso_value: str) -> str:
    dt = datetime.fromisoformat(iso_value.replace("Z", "+00:00"))
    return dt.strftime("%Y-%m-%d")


def count_requirements_lines(path: Path) -> int:
    count = 0
    for line in path.read_text(encoding="utf-8").splitlines():
        stripped = line.strip()
        if stripped and not stripped.startswith("#"):
            count += 1
    return count


def load_commits(start_ref: str) -> list[CommitInfo]:
    raw = git("log", "--reverse", "--format=%H%x1f%h%x1f%cI%x1f%s", f"{start_ref}..HEAD")
    commits: list[CommitInfo] = []
    if not raw:
        return commits
    for line in raw.splitlines():
        full_hash, short_hash, committed_at, subject = line.split("\x1f", 3)
        commits.append(
            CommitInfo(
                full_hash=full_hash,
                short_hash=short_hash,
                committed_at=committed_at,
                subject=subject,
            )
        )
    return commits


def load_diff_stats(start_ref: str) -> list[FileDelta]:
    raw = git("diff", "--numstat", f"{start_ref}..HEAD")
    stats: list[FileDelta] = []
    for line in raw.splitlines():
        if not line.strip():
            continue
        inserted, deleted, path = line.split("\t", 2)
        if inserted == "-" or deleted == "-":
            continue
        stats.append(
            FileDelta(
                path=path,
                insertions=int(inserted),
                deletions=int(deleted),
            )
        )
    stats.sort(key=lambda item: item.churn, reverse=True)
    return stats


def recent_commit_window(paths: Iterable[str]) -> tuple[str, str]:
    baseline = git("log", "-n", "1", "--format=%H", "--", *paths)
    baseline_date = git("show", "-s", "--format=%cI", baseline)
    return baseline, baseline_date


def collect_context() -> dict:
    baseline_commit, baseline_date = recent_commit_window(DOCUMENT_PATHS)
    head_full = git("rev-parse", "HEAD")
    head_short = git("rev-parse", "--short", "HEAD")
    head_date = git("show", "-s", "--format=%cI", "HEAD")
    head_subject = git("show", "-s", "--format=%s", "HEAD")

    commits = load_commits(baseline_commit)
    diff_stats = load_diff_stats(baseline_commit)
    total_insertions = sum(item.insertions for item in diff_stats)
    total_deletions = sum(item.deletions for item in diff_stats)

    tracked_files = len([line for line in git("ls-files").splitlines() if line.strip()])
    workflow_agents = len(
        [
            entry
            for entry in (PROJECT_DIR / "agent" / "agents").iterdir()
            if entry.is_dir() and entry.name != "pools" and (entry / "config.yaml").exists()
        ]
    )
    js_modules = len(list((PROJECT_DIR / "agent" / "static" / "agent" / "js").glob("*.js")))
    css_files = len(list((PROJECT_DIR / "agent" / "static" / "agent" / "css").glob("*.css")))
    html_templates = len(list((PROJECT_DIR / "agent" / "templates" / "agent").glob("*.html")))
    requirements_count = count_requirements_lines(REPO_ROOT / "requirements.txt")

    current_time = datetime.now().astimezone().isoformat(timespec="seconds")

    return {
        "generated_at": current_time,
        "baseline_commit": baseline_commit,
        "baseline_short": baseline_commit[:7],
        "baseline_date": baseline_date,
        "head_commit": head_full,
        "head_short": head_short,
        "head_date": head_date,
        "head_subject": head_subject,
        "tracked_files": tracked_files,
        "workflow_agents": workflow_agents,
        "js_modules": js_modules,
        "css_files": css_files,
        "html_templates": html_templates,
        "requirements_count": requirements_count,
        "recent_commits": commits,
        "recent_commit_count": len(commits),
        "diff_stats": diff_stats,
        "changed_files_count": len(diff_stats),
        "insertions": total_insertions,
        "deletions": total_deletions,
        "top_changed_files": diff_stats[:10],
        "themes": THEMES,
    }


def build_markdown(context: dict) -> str:
    commit_rows = []
    for commit in context["recent_commits"]:
        focus = COMMIT_FOCUS.get(commit.short_hash, commit.subject)
        commit_rows.append(
            f"| {plain_date(commit.committed_at)} | `{commit.short_hash}` | {focus} |"
        )

    top_file_rows = []
    for delta in context["top_changed_files"]:
        top_file_rows.append(
            f"| `{delta.path}` | +{delta.insertions} | -{delta.deletions} | {delta.churn} |"
        )

    theme_chunks = []
    for theme in context["themes"]:
        bullets = "\n".join(f"- {item}" for item in theme.bullets)
        files = ", ".join(f"`{item}`" for item in theme.files)
        theme_chunks.append(
            textwrap.dedent(
                f"""
                ## {theme.title}

                {theme.summary}

                {bullets}

                Key files: {files}
                """
            ).strip()
        )

    exact_subjects = "\n".join(
        f"1. `{commit.short_hash}` - {commit.subject}" for commit in context["recent_commits"]
    )

    markdown = textwrap.dedent(
        f"""
        # Tlamatini App Summary

        Repository refresh generated on {iso_to_human(context["generated_at"])}.

        The previous committed documentation refresh for `tlamatini_app_summary.pdf` and `Tlamatini_eXtended_Artificial_Intelligence_Humanly_Tempered.pptx` landed on {iso_to_human(context["baseline_date"])} in commit `{context["baseline_short"]}`. This update analyzes everything committed after that point through current `HEAD` `{context["head_short"]}` ({context["head_subject"]}) at {iso_to_human(context["head_date"])}.

        ## Executive Summary

        Tlamatini's post-April-11 work is concentrated in six areas: Multi-Turn loop hardening, wrapped-runtime permanence, flow generation from successful Multi-Turn answers, LLM-based answer-success gating, expanded Gatewayer/operator documentation, and a small dependency refresh via `openpyxl`.

        The most important user-facing change is that a successful checked Multi-Turn answer can now produce a starter `.flw` workflow from the exact tools that ran. The most important infrastructure change is that wrapped chat-agent runs now persist as sequenced runtime directories with stronger logging and better source-vs-frozen path handling.

        ## Refresh Window At A Glance

        | Measure | Value |
        |---|---|
        | Previous docs commit | `{context["baseline_short"]}` |
        | Previous docs timestamp | {iso_to_human(context["baseline_date"])} |
        | Current HEAD | `{context["head_short"]}` |
        | Current HEAD timestamp | {iso_to_human(context["head_date"])} |
        | Commits in this window | {context["recent_commit_count"]} |
        | Files changed in this window | {context["changed_files_count"]} |
        | Total insertions | +{context["insertions"]} |
        | Total deletions | -{context["deletions"]} |
        | Current tracked files | {context["tracked_files"]} |
        | Workflow agents | {context["workflow_agents"]} |
        | Frontend JS modules | {context["js_modules"]} |
        | Frontend CSS files | {context["css_files"]} |
        | Chat / ACP HTML templates | {context["html_templates"]} |
        | Python requirements | {context["requirements_count"]} |

        ## Commit Timeline Since The Last Documentation Refresh

        | Date | Commit | Focus |
        |---|---|---|
        __COMMIT_ROWS__

        __THEME_CHUNKS__

        ## Most Changed Files Since The Previous Docs Refresh

        | Path | Insertions | Deletions | Churn |
        |---|---|---|---|
        __TOP_FILE_ROWS__

        ## Current Operator Takeaways

        - Use unchecked chat for fast one-shot answers and checked Multi-Turn for multi-step work where tool selection, monitoring, or wrapped agents matter.
        - Treat the new `Create Flow` button as a starting point, not as a finished production workflow. Validate the `.flw` in ACP before running it.
        - When debugging wrapped agents, inspect the sequenced runtime folders under `agent/agents/pools/_chat_runs_/` and the new INFO-level logs from the planner, tool layer, and chat runtime.
        - Gatewayer is now documented well enough to act as a serious ingress boundary, but the docs also clearly identify config fields that are declared today yet not fully enforced in code.
        - `openpyxl` is now present in the bundled environment, which broadens what generated Python or workflow scripts can do with `.xlsx` files.

        ## Documentation Sync Notes

        README now covers the end-to-end flow-creation pipeline from checked Multi-Turn answers, the answer-success classifier, the wrapped runtime lifecycle, and a much deeper Gatewayer section. The refreshed PPT deck that accompanies this PDF turns those same changes into a concise operator-facing briefing.

        ## Appendix: Exact Subjects In The Window

        __EXACT_SUBJECTS__

        <p class="small">
        Generated from the current repository state, git history after the last committed document refresh, and the checked-in README / source files that now define the latest Multi-Turn, runtime, and Gatewayer behavior.
        </p>
        """
    ).strip()

    markdown = markdown.replace("__COMMIT_ROWS__", "\n".join(commit_rows))
    markdown = markdown.replace("__THEME_CHUNKS__", "\n\n".join(theme_chunks))
    markdown = markdown.replace("__TOP_FILE_ROWS__", "\n".join(top_file_rows))
    markdown = markdown.replace("__EXACT_SUBJECTS__", exact_subjects)
    return markdown + "\n"


def set_shape_fill(shape, color: RGBColor) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_background(
    slide,
    audit: list[tuple[int, int, int, int, str]],
    accent_color: RGBColor,
    slide_width: int,
    slide_height: int,
) -> None:

    background = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        slide_width,
        slide_height,
    )
    set_shape_fill(background, PALETTE["paper"])
    background.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(0.28),
        slide_height,
    )
    set_shape_fill(accent, accent_color)
    accent.line.fill.background()

    ribbon = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.28),
        0,
        slide_width - Inches(0.28),
        Inches(0.22),
    )
    set_shape_fill(ribbon, accent_color)
    ribbon.fill.transparency = 0.82
    ribbon.line.fill.background()


def add_textbox(
    slide,
    audit: list[tuple[int, int, int, int, str]],
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 18,
    color: RGBColor | None = None,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    name: str = "textbox",
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or PALETTE["ink"]
    audit.append((box.left, box.top, box.width, box.height, name))


def add_title(
    slide,
    audit: list[tuple[int, int, int, int, str]],
    title: str,
    subtitle: str,
    accent_color: RGBColor,
) -> None:
    add_textbox(
        slide,
        audit,
        left=0.68,
        top=0.48,
        width=11.8,
        height=0.4,
        text=subtitle,
        font_size=12,
        color=accent_color,
        bold=True,
        name="subtitle",
    )
    add_textbox(
        slide,
        audit,
        left=0.68,
        top=0.88,
        width=11.8,
        height=0.7,
        text=title,
        font_size=24,
        color=PALETTE["ink"],
        bold=True,
        name="title",
    )


def add_panel(
    slide,
    audit: list[tuple[int, int, int, int, str]],
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list[str],
    accent_color: RGBColor,
    accent_fill: RGBColor,
    name: str,
    body_size: int = 16,
) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PALETTE["panel"]
    panel.line.color.rgb = PALETTE["line"]
    audit.append((panel.left, panel.top, panel.width, panel.height, name))

    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.18),
    )
    set_shape_fill(stripe, accent_color)
    stripe.line.fill.background()

    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.22),
        Inches(top + 0.28),
        Inches(min(width - 0.44, 2.85)),
        Inches(0.44),
    )
    set_shape_fill(chip, accent_fill)
    chip.line.fill.background()

    add_textbox(
        slide,
        audit,
        left=left + 0.32,
        top=top + 0.34,
        width=min(width - 0.64, 2.65),
        height=0.28,
        text=title,
        font_size=13,
        color=accent_color,
        bold=True,
        name=f"{name}-header",
    )

    body_box = slide.shapes.add_textbox(
        Inches(left + 0.22),
        Inches(top + 0.82),
        Inches(width - 0.44),
        Inches(height - 1.0),
    )
    body = body_box.text_frame
    body.clear()
    body.word_wrap = True
    body.vertical_anchor = MSO_ANCHOR.TOP
    body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    body.margin_left = Pt(4)
    body.margin_right = Pt(4)
    body.margin_top = Pt(2)
    body.margin_bottom = Pt(2)

    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(body_size)
        paragraph.font.color.rgb = PALETTE["slate"]
        paragraph.space_after = Pt(6)
        paragraph.bullet = True

    try:
        body.fit_text(max_size=body_size)
    except Exception:
        body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    audit.append((body_box.left, body_box.top, body_box.width, body_box.height, f"{name}-body"))


def add_stat_card(
    slide,
    audit: list[tuple[int, int, int, int, str]],
    left: float,
    top: float,
    width: float,
    title: str,
    value: str,
    accent_color: RGBColor,
    fill_color: RGBColor,
    name: str,
) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(1.1),
    )
    set_shape_fill(card, fill_color)
    card.line.fill.background()
    audit.append((card.left, card.top, card.width, card.height, name))

    add_textbox(
        slide,
        audit,
        left=left + 0.18,
        top=top + 0.14,
        width=width - 0.36,
        height=0.3,
        text=title,
        font_size=11,
        color=accent_color,
        bold=True,
        name=f"{name}-title",
    )
    add_textbox(
        slide,
        audit,
        left=left + 0.18,
        top=top + 0.42,
        width=width - 0.36,
        height=0.45,
        text=value,
        font_size=20,
        color=PALETTE["ink"],
        bold=True,
        name=f"{name}-value",
    )


def add_flow_strip(
    slide,
    audit: list[tuple[int, int, int, int, str]],
    left: float,
    top: float,
    labels: list[str],
) -> None:
    current_left = left
    box_width = 2.05
    arrow_width = 0.5
    for index, label in enumerate(labels):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(current_left),
            Inches(top),
            Inches(box_width),
            Inches(0.82),
        )
        set_shape_fill(card, PALETTE["sky_light"] if index % 2 == 0 else PALETTE["green_light"])
        card.line.fill.background()
        audit.append((card.left, card.top, card.width, card.height, f"flow-{index}"))
        add_textbox(
            slide,
            audit,
            left=current_left + 0.12,
            top=top + 0.17,
            width=box_width - 0.24,
            height=0.34,
            text=label,
            font_size=12,
            color=PALETTE["ink"],
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"flow-text-{index}",
        )
        current_left += box_width
        if index < len(labels) - 1:
            add_textbox(
                slide,
                audit,
                left=current_left,
                top=top + 0.14,
                width=arrow_width,
                height=0.34,
                text="->",
                font_size=16,
                color=PALETTE["teal"],
                bold=True,
                align=PP_ALIGN.CENTER,
                name=f"arrow-{index}",
            )
            current_left += arrow_width


def rectangles_overlap(
    first: tuple[int, int, int, int, str],
    second: tuple[int, int, int, int, str],
) -> bool:
    l1, t1, w1, h1, _ = first
    l2, t2, w2, h2, _ = second
    r1 = l1 + w1
    b1 = t1 + h1
    r2 = l2 + w2
    b2 = t2 + h2
    return l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2


def layout_audit(audit: list[tuple[int, int, int, int, str]], prs: Presentation, slide_number: int) -> None:
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    for left, top, width, height, name in audit:
        if left < 0 or top < 0 or left + width > slide_width or top + height > slide_height:
            raise RuntimeError(
                f"Slide {slide_number}: shape '{name}' is outside the slide bounds."
            )
    content_items = [item for item in audit if not item[4].startswith("flow-text") and not item[4].startswith("arrow")]
    for index, first in enumerate(content_items):
        for second in content_items[index + 1 :]:
            if first[4].startswith("subtitle") or second[4].startswith("subtitle"):
                continue
            if first[4].startswith("title") and second[4].startswith("title"):
                continue
            if rectangles_overlap(first, second):
                if first[4].split("-")[0] == second[4].split("-")[0]:
                    continue
                raise RuntimeError(
                    f"Slide {slide_number}: '{first[4]}' overlaps '{second[4]}'."
                )


def build_presentation(context: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide(title: str, subtitle: str, accent: RGBColor):
        slide = prs.slides.add_slide(blank)
        audit: list[tuple[int, int, int, int, str]] = []
        add_background(slide, audit, accent, prs.slide_width, prs.slide_height)
        add_title(slide, audit, title, subtitle, accent)
        return slide, audit

    slide, audit = new_slide(
        "Tlamatini Documentation Refresh",
        "Update window: previous docs commit April 11, 2026 -> current HEAD April 14, 2026",
        PALETTE["teal"],
    )
    add_textbox(
        slide,
        audit,
        left=0.72,
        top=1.72,
        width=6.45,
        height=1.15,
        text=(
            "This rebuilt deck summarizes the exact repository changes that landed "
            "after the previous PDF/PPT refresh and turns them into a concise operator briefing."
        ),
        font_size=19,
        color=PALETTE["slate"],
        name="hero-copy",
    )
    add_stat_card(slide, audit, 7.55, 1.6, 1.62, "Commits", str(context["recent_commit_count"]), PALETTE["teal"], PALETTE["teal_light"], "hero-stat-1")
    add_stat_card(slide, audit, 9.35, 1.6, 1.62, "Files", str(context["changed_files_count"]), PALETTE["gold"], PALETTE["gold_light"], "hero-stat-2")
    add_stat_card(slide, audit, 11.15, 1.6, 1.62, "Agents", str(context["workflow_agents"]), PALETTE["red"], PALETTE["red_light"], "hero-stat-3")
    add_panel(
        slide,
        audit,
        left=0.72,
        top=3.05,
        width=5.95,
        height=3.1,
        title="What changed most",
        bullets=[
            "Checked Multi-Turn got loop, route, stop, and runtime-lifecycle hardening.",
            "Successful answers can now become starter `.flw` workflows.",
            "A new answer classifier gates flow generation more safely than regex checks.",
            "Gatewayer documentation is now deep enough for serious operator use.",
        ],
        accent_color=PALETTE["teal"],
        accent_fill=PALETTE["teal_light"],
        name="hero-left",
        body_size=16,
    )
    add_panel(
        slide,
        audit,
        left=6.95,
        top=3.05,
        width=5.65,
        height=3.1,
        title="Refresh window facts",
        bullets=[
            f"Previous docs commit: {context['baseline_short']}",
            f"Current HEAD: {context['head_short']}",
            f"Insertions / deletions: +{context['insertions']} / -{context['deletions']}",
            f"Current tracked files: {context['tracked_files']}",
            f"Frontend modules: {context['js_modules']} JS, {context['css_files']} CSS, {context['html_templates']} HTML",
        ],
        accent_color=PALETTE["gold"],
        accent_fill=PALETTE["gold_light"],
        name="hero-right",
        body_size=15,
    )
    layout_audit(audit, prs, 1)

    slide, audit = new_slide("Commit Window", "The 10 commits analyzed in this refresh", PALETTE["gold"])
    left_commits = context["recent_commits"][:5]
    right_commits = context["recent_commits"][5:]
    add_panel(
        slide,
        audit,
        left=0.72,
        top=1.58,
        width=5.9,
        height=4.95,
        title="April 12 -> April 13",
        bullets=[
            f"{plain_date(commit.committed_at)} | {commit.short_hash} | {COMMIT_FOCUS.get(commit.short_hash, commit.subject)}"
            for commit in left_commits
        ],
        accent_color=PALETTE["gold"],
        accent_fill=PALETTE["gold_light"],
        name="timeline-left",
        body_size=14,
    )
    add_panel(
        slide,
        audit,
        left=6.75,
        top=1.58,
        width=5.9,
        height=4.95,
        title="April 13 -> April 14",
        bullets=[
            f"{plain_date(commit.committed_at)} | {commit.short_hash} | {COMMIT_FOCUS.get(commit.short_hash, commit.subject)}"
            for commit in right_commits
        ],
        accent_color=PALETTE["teal"],
        accent_fill=PALETTE["teal_light"],
        name="timeline-right",
        body_size=14,
    )
    layout_audit(audit, prs, 2)

    slide, audit = new_slide("Current Repository Snapshot", "Where the project stands after the refresh window", PALETTE["sky"])
    add_stat_card(slide, audit, 0.78, 1.65, 2.15, "Tracked files", str(context["tracked_files"]), PALETTE["sky"], PALETTE["sky_light"], "snapshot-1")
    add_stat_card(slide, audit, 3.12, 1.65, 2.15, "Workflow agents", str(context["workflow_agents"]), PALETTE["teal"], PALETTE["teal_light"], "snapshot-2")
    add_stat_card(slide, audit, 5.46, 1.65, 2.15, "JS modules", str(context["js_modules"]), PALETTE["gold"], PALETTE["gold_light"], "snapshot-3")
    add_stat_card(slide, audit, 7.8, 1.65, 2.15, "HTML templates", str(context["html_templates"]), PALETTE["red"], PALETTE["red_light"], "snapshot-4")
    add_stat_card(slide, audit, 10.14, 1.65, 2.15, "Requirements", str(context["requirements_count"]), PALETTE["green"], PALETTE["green_light"], "snapshot-5")
    add_panel(
        slide,
        audit,
        left=0.78,
        top=3.05,
        width=6.02,
        height=2.95,
        title="Why the snapshot matters",
        bullets=[
            "The repo is still centered on Django + Channels + an in-repo `agent` application, but the hot path has moved further toward agentic orchestration.",
            "The most changed files in this window are README, `agent_page_chat.js`, `mcp_agent.py`, `chat_agent_runtime.py`, and the new Gatewayer explanation.",
            "The refresh updates both summary artifacts so they match the current HEAD instead of the earlier April 11 picture.",
        ],
        accent_color=PALETTE["sky"],
        accent_fill=PALETTE["sky_light"],
        name="snapshot-left",
        body_size=16,
    )
    add_panel(
        slide,
        audit,
        left=7.0,
        top=3.05,
        width=5.55,
        height=2.95,
        title="Largest themes by churn",
        bullets=[
            "README expansion and sync work.",
            "Frontend chat logic for flow creation and safe rendering.",
            "Multi-Turn backend orchestration in `mcp_agent.py`.",
            "Runtime persistence and logging in `chat_agent_runtime.py`.",
        ],
        accent_color=PALETTE["teal"],
        accent_fill=PALETTE["teal_light"],
        name="snapshot-right",
        body_size=16,
    )
    layout_audit(audit, prs, 3)

    slide, audit = new_slide("Multi-Turn Hardening", "Stability, compatibility, and request-scoped execution", PALETTE["teal"])
    theme = THEMES[0]
    add_panel(
        slide,
        audit,
        left=0.78,
        top=1.6,
        width=6.0,
        height=4.95,
        title=theme.title,
        bullets=theme.bullets,
        accent_color=PALETTE["teal"],
        accent_fill=PALETTE["teal_light"],
        name="mt-left",
        body_size=16,
    )
    add_panel(
        slide,
        audit,
        left=6.98,
        top=1.6,
        width=5.55,
        height=4.95,
        title="Files that moved this work",
        bullets=[
            "`mcp_agent.py` keeps the explicit tool loop and planner injection.",
            "`global_execution_planner.py` keeps the request DAG and selected stages.",
            "`capability_registry.py` continues to shape which tools/context enter the turn.",
            "`rag/interface.py` and `rag/chains/unified.py` keep the opt-in split between checked and unchecked modes.",
        ],
        accent_color=PALETTE["gold"],
        accent_fill=PALETTE["gold_light"],
        name="mt-right",
        body_size=15,
    )
    layout_audit(audit, prs, 4)

    slide, audit = new_slide("Wrapped Runtime Persistence", "Every wrapped run now leaves a more inspectable trail", PALETTE["gold"])
    theme = THEMES[1]
    add_panel(
        slide,
        audit,
        left=0.78,
        top=1.6,
        width=6.1,
        height=4.95,
        title=theme.title,
        bullets=theme.bullets,
        accent_color=PALETTE["gold"],
        accent_fill=PALETTE["gold_light"],
        name="runtime-left",
        body_size=16,
    )
    add_panel(
        slide,
        audit,
        left=7.0,
        top=1.6,
        width=5.53,
        height=4.95,
        title="Operator value",
        bullets=[
            "You can inspect failed runs without losing later successful ones.",
            "Detached launches are less noisy in checked Multi-Turn mode.",
            "Frozen-build path problems are easier to diagnose from logs alone.",
            "Runtime, planner, and tool selection logs now appear in the Django console at INFO level.",
        ],
        accent_color=PALETTE["red"],
        accent_fill=PALETTE["red_light"],
        name="runtime-right",
        body_size=16,
    )
    layout_audit(audit, prs, 5)

    slide, audit = new_slide("Flow Creation From Answers", "Successful Multi-Turn executions can now seed ACP workflows", PALETTE["sky"])
    theme = THEMES[2]
    add_panel(
        slide,
        audit,
        left=0.78,
        top=1.58,
        width=12.0,
        height=2.3,
        title=theme.title,
        bullets=theme.bullets,
        accent_color=PALETTE["sky"],
        accent_fill=PALETTE["sky_light"],
        name="flow-top",
        body_size=16,
    )
    add_flow_strip(
        slide,
        audit,
        left=0.95,
        top=4.25,
        labels=[
            "Checked Multi-Turn",
            "tool_calls_log",
            "answer_success",
            "Create Flow",
            "download .flw",
        ],
    )
    add_panel(
        slide,
        audit,
        left=1.55,
        top=5.15,
        width=10.5,
        height=1.2,
        title="Key guardrails",
        bullets=[
            "Only successful flow-eligible tools become nodes, and management-only runtime tools stay out of the generated workflow.",
            "The output is a starter chain, not a production-validated flow; ACP validation is still the next step.",
        ],
        accent_color=PALETTE["teal"],
        accent_fill=PALETTE["teal_light"],
        name="flow-bottom",
        body_size=15,
    )
    layout_audit(audit, prs, 6)

    slide, audit = new_slide("Answer Analizer", "Safer gating for the new Create Flow action", PALETTE["red"])
    theme = THEMES[3]
    add_panel(
        slide,
        audit,
        left=0.78,
        top=1.58,
        width=6.0,
        height=4.95,
        title=theme.title,
        bullets=theme.bullets,
        accent_color=PALETTE["red"],
        accent_fill=PALETTE["red_light"],
        name="analizer-left",
        body_size=16,
    )
    add_panel(
        slide,
        audit,
        left=6.98,
        top=1.58,
        width=5.55,
        height=4.95,
        title="Why this matters",
        bullets=[
            "The UI no longer guesses success by scraping arbitrary wording from the answer.",
            "Backend classification keeps the decision close to the actual execution result and final answer text.",
            "If classification itself fails, the system chooses the less frustrating fallback: show the button rather than hide it.",
        ],
        accent_color=PALETTE["gold"],
        accent_fill=PALETTE["gold_light"],
        name="analizer-right",
        body_size=16,
    )
    layout_audit(audit, prs, 7)

    slide, audit = new_slide("Gatewayer And Docs Expansion", "Operator-facing coverage is much deeper than before", PALETTE["teal"])
    theme = THEMES[4]
    add_panel(
        slide,
        audit,
        left=0.78,
        top=1.58,
        width=6.08,
        height=4.95,
        title=theme.title,
        bullets=theme.bullets,
        accent_color=PALETTE["teal"],
        accent_fill=PALETTE["teal_light"],
        name="docs-left",
        body_size=16,
    )
    add_panel(
        slide,
        audit,
        left=7.0,
        top=1.58,
        width=5.53,
        height=4.95,
        title="Gatewayer themes now documented",
        bullets=[
            "HTTP ingress plus optional folder-drop ingestion.",
            "Bearer and HMAC auth plus allowlist notes.",
            "Crash recovery via persisted queue and dedup files.",
            "Stable log markers for Monitor Log and Summarizer workflows.",
            "Clear caveats about config fields that exist before full enforcement.",
        ],
        accent_color=PALETTE["sky"],
        accent_fill=PALETTE["sky_light"],
        name="docs-right",
        body_size=15,
    )
    layout_audit(audit, prs, 8)

    slide, audit = new_slide("Dependency And File Impact", "What changed around the edges of the feature work", PALETTE["gold"])
    theme = THEMES[5]
    add_panel(
        slide,
        audit,
        left=0.78,
        top=1.58,
        width=5.9,
        height=4.95,
        title=theme.title,
        bullets=theme.bullets,
        accent_color=PALETTE["gold"],
        accent_fill=PALETTE["gold_light"],
        name="impact-left",
        body_size=16,
    )
    top_five = context["top_changed_files"][:5]
    add_panel(
        slide,
        audit,
        left=6.92,
        top=1.58,
        width=5.61,
        height=4.95,
        title="Top files by churn",
        bullets=[
            f"{item.path} | +{item.insertions} / -{item.deletions}"
            for item in top_five
        ],
        accent_color=PALETTE["red"],
        accent_fill=PALETTE["red_light"],
        name="impact-right",
        body_size=15,
    )
    layout_audit(audit, prs, 9)

    slide, audit = new_slide("Recommendations", "How to use the refreshed capabilities safely", PALETTE["green"])
    add_panel(
        slide,
        audit,
        left=0.92,
        top=1.75,
        width=11.7,
        height=4.75,
        title="Recommended next-step behavior",
        bullets=[
            "Keep unchecked chat as the default for direct analysis; switch to checked Multi-Turn when you want tool orchestration, monitoring, or flow seeding.",
            "Inspect generated `.flw` output in ACP and run validation before executing it in a real workflow.",
            "Use the new runtime directories and INFO logs whenever a wrapped agent behaves differently between source mode and a frozen desktop build.",
            "Treat Gatewayer as an ingress control point: pair it with explicit auth, monitoring, and downstream validation rather than assuming every declared config key is already enforced.",
            "Leverage `openpyxl` when generated Python or workflow scripts need to read or write `.xlsx` content.",
        ],
        accent_color=PALETTE["green"],
        accent_fill=PALETTE["green_light"],
        name="recommendations",
        body_size=17,
    )
    layout_audit(audit, prs, 10)

    PPT_OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_OUTPUT)


def write_outputs() -> None:
    context = collect_context()
    BUILD_DIR.mkdir(parents=True, exist_ok=True)
    markdown = build_markdown(context)
    MARKDOWN_OUTPUT.write_text(markdown, encoding="utf-8")
    serializable_context = {
        key: value
        for key, value in context.items()
        if key not in {"recent_commits", "diff_stats", "top_changed_files", "themes"}
    }
    JSON_OUTPUT.write_text(
        json.dumps(
            {
                **serializable_context,
                "recent_commits": [commit.__dict__ for commit in context["recent_commits"]],
                "diff_stats": [delta.__dict__ for delta in context["diff_stats"]],
                "top_changed_files": [delta.__dict__ for delta in context["top_changed_files"]],
                "themes": [theme.__dict__ for theme in context["themes"]],
            },
            indent=2,
        ),
        encoding="utf-8",
    )
    markdown_text_to_pdf(markdown, PDF_OUTPUT, base_dir=REPO_ROOT, css_text=PDF_CSS)
    build_presentation(context)
    print(f"Updated PDF: {PDF_OUTPUT}")
    print(f"Updated PPTX: {PPT_OUTPUT}")
    print(f"Saved source markdown: {MARKDOWN_OUTPUT}")


if __name__ == "__main__":
    write_outputs()
