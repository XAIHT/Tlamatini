# ═══════════════════════════════════════════════════════════════════
#   ✦  T L A M A T I N I  ✦   —   "one who knows"
#
#   Created by  Angela López Mendoza   ·   @angelahack1
#   Developer · Architect · Creator of Tlamatini
#
#   Every line of this file was written by Angela López Mendoza.
# ═══════════════════════════════════════════════════════════════════
#   Tlamatini Author Banner — do not remove (releases scrub the name automatically)
from __future__ import annotations

import ast
import io
import importlib.util
import json
import os
import re
import subprocess
import tokenize
from dataclasses import dataclass
from datetime import datetime, timedelta
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.lib.utils import ImageReader
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    Preformatted,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


SCRIPT_PATH = Path(__file__).resolve()
DOC_DIR = SCRIPT_PATH.parent
PROJECT_DIR = SCRIPT_PATH.parents[2]
REPO_ROOT = SCRIPT_PATH.parents[3]
BUILD_DIR = REPO_ROOT / "build" / "documentation_refresh"
REFERENCE_MEDIA_DIR = BUILD_DIR / "reference_ppt_media"

PDF_OUTPUT = REPO_ROOT / "tlamatini_app_summary.pdf"
PPT_OUTPUT = REPO_ROOT / "Tlamatini_eXtended_Artificial_Intelligence_Humanly_Tempered.pptx"
CONTEXT_OUTPUT = BUILD_DIR / "complete_project_dossier_context.json"
TREE_OUTPUT = BUILD_DIR / "complete_tracked_file_tree.txt"

BINARY_EXTENSIONS = {
    ".7z",
    ".dll",
    ".exe",
    ".flw",
    ".ico",
    ".jar",
    ".jpg",
    ".jpeg",
    ".mp3",
    ".mp4",
    ".pdf",
    ".png",
    ".pptx",
    ".pyc",
    ".wav",
    ".zip",
}

LANGUAGE_BY_EXTENSION = {
    ".bat": ("Batch", "Batch"),
    ".css": ("CSS", "CSS"),
    ".html": ("HTML", "HTML"),
    ".js": ("JavaScript", "JS"),
    ".json": ("JSON", "JSON"),
    ".md": ("Markdown", "MD"),
    ".mjs": ("JavaScript module", "MJS"),
    ".pmt": ("Prompt template", "Prompt"),
    ".proto": ("Protocol Buffers", "Proto"),
    ".ps1": ("PowerShell", "PowerShell"),
    ".py": ("Python", "Python"),
    ".txt": ("Text", "Text"),
    ".yaml": ("YAML", "YAML"),
    ".yml": ("YAML", "YAML"),
}

THEME = {
    "obsidian": RGBColor(8, 13, 13),
    "void": RGBColor(13, 18, 19),
    "stone": RGBColor(32, 35, 34),
    "panel": RGBColor(22, 26, 25),
    "panel2": RGBColor(34, 39, 36),
    "white": RGBColor(232, 232, 222),
    "muted": RGBColor(180, 186, 174),
    "copper": RGBColor(196, 128, 82),
    "copper2": RGBColor(142, 89, 52),
    "jade": RGBColor(72, 191, 143),
    "jade2": RGBColor(33, 127, 98),
    "amber": RGBColor(224, 171, 93),
    "line": RGBColor(89, 95, 86),
}

SLIDE_W = 13.333
SLIDE_H = 7.5
RECENT_GIT_WINDOW_DAYS = 3
RECENT_GIT_WINDOW_LABEL = "last 3 days"
RECENT_GIT_WINDOW_TITLE = "Recent Git Window"
RECENT_GIT_HIGHLIGHT_TITLE = "recent highlights"
RECENT_GIT_APPENDIX_SUBTITLE = "all commits from the last 3 days according to git"


@dataclass
class LineStats:
    language: str
    short: str
    files: int = 0
    total_lines: int = 0
    effective_lines: int = 0


@dataclass
class FileStats:
    path: str
    language: str
    total_lines: int
    effective_lines: int


@dataclass
class CommitInfo:
    short_hash: str
    committed_at: str
    subject: str


@dataclass
class CommitBaseline:
    full_hash: str
    short_hash: str
    committed_at: str
    subject: str


def git(*args: str) -> str:
    result = subprocess.run(
        ["git", *args],
        cwd=REPO_ROOT,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        check=True,
    )
    return result.stdout.strip()


def local_stamp() -> str:
    return datetime.now().astimezone().strftime("%B %d, %Y %I:%M %p UTC%z")


def iso_date(iso_value: str) -> str:
    return datetime.fromisoformat(iso_value.replace("Z", "+00:00")).strftime("%Y-%m-%d")


def discover_reference_deck() -> Path | None:
    desktop = Path.home() / "OneDrive" / "Desktop"
    if not desktop.exists():
        desktop = Path.home() / "Desktop"
    matches = sorted(desktop.glob("TLAMATINI_ El Saber*.pptx"))
    return matches[0] if matches else None


def extract_reference_media() -> list[Path]:
    deck = discover_reference_deck()
    if deck is None or not deck.exists():
        return []
    REFERENCE_MEDIA_DIR.mkdir(parents=True, exist_ok=True)
    extracted: list[Path] = []
    try:
        with ZipFile(deck) as archive:
            for name in archive.namelist():
                if not name.startswith("ppt/media/"):
                    continue
                suffix = Path(name).suffix.lower()
                if suffix not in {".png", ".jpg", ".jpeg"}:
                    continue
                target = REFERENCE_MEDIA_DIR / Path(name).name
                target.write_bytes(archive.read(name))
                extracted.append(target)
    except Exception:
        return []
    return extracted


def git_tracked_paths() -> list[str]:
    return [line for line in git("ls-files").splitlines() if line.strip()]


def git_untracked_paths() -> list[str]:
    return [line for line in git("ls-files", "--others", "--exclude-standard").splitlines() if line.strip()]


def inventory_paths() -> list[str]:
    return sorted(set(git_tracked_paths()) | set(git_untracked_paths()))


def has_esphomer_assets() -> bool:
    return (PROJECT_DIR / "agent" / "agents" / "esphomer" / "config.yaml").exists()


def build_tree(paths: list[str]) -> str:
    root: dict[str, dict] = {}
    for raw_path in sorted(paths):
        parts = raw_path.replace("\\", "/").split("/")
        node = root
        for part in parts:
            node = node.setdefault(part, {})

    def walk(node: dict, prefix: str = "") -> list[str]:
        lines: list[str] = []
        entries = sorted(node.items(), key=lambda item: (bool(item[1]), item[0].lower()))
        for index, (name, child) in enumerate(entries):
            last = index == len(entries) - 1
            connector = "`-- " if last else "|-- "
            suffix = "/" if child else ""
            lines.append(f"{prefix}{connector}{name}{suffix}")
            if child:
                extension = "    " if last else "|   "
                lines.extend(walk(child, prefix + extension))
        return lines

    return "Tlamatini/\n" + "\n".join(walk(root))


def detect_language(path: str) -> tuple[str, str]:
    suffix = Path(path).suffix.lower()
    return LANGUAGE_BY_EXTENSION.get(suffix, ("Other text", "Other"))


def looks_binary(path: Path) -> bool:
    if path.suffix.lower() in BINARY_EXTENSIONS:
        return True
    try:
        chunk = path.read_bytes()[:2048]
    except OSError:
        return True
    return b"\x00" in chunk


def read_text(path: Path) -> str | None:
    if looks_binary(path):
        return None
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
        except OSError:
            return None
    return None


def python_docstring_lines(text: str) -> set[int]:
    try:
        tree = ast.parse(text)
    except SyntaxError:
        return set()
    doc_lines: set[int] = set()
    candidates = [tree, *[node for node in ast.walk(tree) if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef))]]
    for node in candidates:
        if not node.body:
            continue
        first = node.body[0]
        if isinstance(first, ast.Expr) and isinstance(first.value, ast.Constant) and isinstance(first.value.value, str):
            start = getattr(first, "lineno", None)
            end = getattr(first, "end_lineno", start)
            if start and end:
                doc_lines.update(range(start, end + 1))
    return doc_lines


def count_python_effective(text: str) -> int:
    doc_lines = python_docstring_lines(text)
    effective_lines: set[int] = set()
    try:
        tokens = tokenize.generate_tokens(io.StringIO(text).readline)
        for token in tokens:
            if token.type in {
                tokenize.COMMENT,
                tokenize.NL,
                tokenize.NEWLINE,
                tokenize.ENCODING,
                tokenize.ENDMARKER,
                tokenize.INDENT,
                tokenize.DEDENT,
            }:
                continue
            line_number = token.start[0]
            if line_number not in doc_lines:
                effective_lines.add(line_number)
    except tokenize.TokenError:
        return count_generic_effective(text, ".py")
    return len(effective_lines)


def remove_block_comments(text: str, suffix: str) -> str:
    if suffix in {".js", ".mjs", ".css", ".proto"}:
        return re.sub(r"/\*.*?\*/", "", text, flags=re.DOTALL)
    if suffix == ".html" or suffix == ".md":
        return re.sub(r"<!--.*?-->", "", text, flags=re.DOTALL)
    if suffix == ".ps1":
        return re.sub(r"<#.*?#>", "", text, flags=re.DOTALL)
    return text


def strip_inline_comment(line: str, suffix: str) -> str:
    stripped = line.strip()
    if suffix in {".yaml", ".yml", ".ps1"}:
        return "" if stripped.startswith("#") else line
    if suffix in {".js", ".mjs", ".css", ".proto"}:
        return "" if stripped.startswith("//") else line
    if suffix == ".bat":
        lowered = stripped.lower()
        return "" if lowered.startswith("rem ") or stripped.startswith("::") else line
    return line


def count_generic_effective(text: str, suffix: str) -> int:
    text = remove_block_comments(text, suffix)
    count = 0
    for raw_line in text.splitlines():
        line = strip_inline_comment(raw_line, suffix)
        if line.strip():
            count += 1
    return count


def line_stats_for_paths(paths: list[str]) -> tuple[list[LineStats], list[FileStats], int, int]:
    by_language: dict[str, LineStats] = {}
    file_rows: list[FileStats] = []
    binary_count = 0
    skipped_count = 0

    for rel_path in paths:
        absolute = REPO_ROOT / rel_path
        text = read_text(absolute)
        if text is None:
            if absolute.suffix.lower() in BINARY_EXTENSIONS:
                binary_count += 1
            else:
                skipped_count += 1
            continue

        language, short = detect_language(rel_path)
        total = len(text.splitlines())
        suffix = absolute.suffix.lower()
        effective = count_python_effective(text) if suffix == ".py" else count_generic_effective(text, suffix)

        stats = by_language.setdefault(language, LineStats(language=language, short=short))
        stats.files += 1
        stats.total_lines += total
        stats.effective_lines += effective
        file_rows.append(FileStats(rel_path, language, total, effective))

    language_rows = sorted(by_language.values(), key=lambda item: item.effective_lines, reverse=True)
    file_rows.sort(key=lambda item: item.effective_lines, reverse=True)
    return language_rows, file_rows, binary_count, skipped_count


def _commit_is_retired_for_dossier(subject: str) -> bool:
    lowered = subject.lower()
    return any(token in lowered for token in ("toast", "toaster", "native_toast", "windows-toast"))


def recent_commits(limit: int = 10) -> list[CommitInfo]:
    raw = git("log", "-n40", "--format=%h%x1f%cI%x1f%s")
    commits: list[CommitInfo] = []
    for line in raw.splitlines():
        short_hash, committed_at, subject = line.split("\x1f", 2)
        if _commit_is_retired_for_dossier(subject):
            continue
        commits.append(CommitInfo(short_hash, committed_at, subject))
        if len(commits) >= limit:
            break
    return commits


def recent_week_commits(days: int = RECENT_GIT_WINDOW_DAYS) -> list[CommitInfo]:
    now = datetime.now().astimezone()
    if RECENT_GIT_WINDOW_LABEL == "today":
        since = now.replace(hour=0, minute=0, second=0, microsecond=0).isoformat()
    else:
        since = (now - timedelta(days=days)).isoformat()
    raw = git("log", f"--since={since}", "--format=%h%x1f%cI%x1f%s")
    commits: list[CommitInfo] = []
    for line in raw.splitlines():
        short_hash, committed_at, subject = line.split("\x1f", 2)
        commits.append(CommitInfo(short_hash, committed_at, subject))
    return commits


def last_visual_doc_commit() -> CommitBaseline | None:
    raw = git(
        "log",
        "-n1",
        "--format=%H%x1f%h%x1f%cI%x1f%s",
        "--",
        PDF_OUTPUT.name,
        PPT_OUTPUT.name,
    )
    if not raw:
        return None
    full_hash, short_hash, committed_at, subject = raw.split("\x1f", 3)
    return CommitBaseline(full_hash, short_hash, committed_at, subject)


def commits_since_visual_docs(baseline: CommitBaseline | None) -> list[CommitInfo]:
    if baseline is None:
        return recent_commits()
    raw = git("log", f"{baseline.full_hash}..HEAD", "--format=%h%x1f%cI%x1f%s")
    commits: list[CommitInfo] = []
    for line in raw.splitlines():
        short_hash, committed_at, subject = line.split("\x1f", 2)
        commits.append(CommitInfo(short_hash, committed_at, subject))
    return commits


def weekly_highlights(commits: list[CommitInfo]) -> list[str]:
    subjects = [commit.subject.lower() for commit in commits]
    highlights: list[str] = []
    has_current_release_wave = any(
        "1.36.0" in subject
        or "1.35.0" in subject
        or "1.33.2" in subject
        or "1.33.0" in subject
        or "1.32.0" in subject
        or "video-analizer" in subject
        or "video-analyzer" in subject
        or "video analyzer" in subject
        or "video_analyzer" in subject
        or "search of prompts" in subject
        or "prompt search" in subject
        or "flw generation" in subject
        or ".flw generation" in subject
        or "retrying behaviour" in subject
        or "self-healing" in subject
        or "self healing" in subject
        or "failure" in subject
        or "failures" in subject
        or "answer classifier" in subject
        or "answer_success" in subject
        or "create flow" in subject
        or "pdcp" in subject
        or "projectdiscovery" in subject
        or "vulnx" in subject
        or "go-deny" in subject
        or "go deny" in subject
        or "zavuerer" in subject
        or "zavu" in subject
        or "annouces skill" in subject
        or "announce skill" in subject
        or "angela" in subject
        or "creator" in subject
        or "author" in subject
        or "private data" in subject
        or "private-data" in subject
        or "public build" in subject
        or "public verify" in subject
        or "3x" in subject
        or "performance" in subject
        or "image-interpreter" in subject
        or "image interpreter" in subject
        or "image_interpreter" in subject
        or "config dialog" in subject
        or "1.39.2" in subject
        or "1.39.3" in subject
        or "1.39.4" in subject
        or "1.39.5" in subject
        or "1.40.0" in subject
        or "1.40.1" in subject
        or "1.41.0" in subject
        or "1.41.2" in subject
        or "1.41.3" in subject
        or "1.41.4" in subject
        or "1.42.0" in subject
        or "stm32er" in subject
        or "structuredcontent" in subject
        or "configurable" in subject
        or "django_port" in subject
        or "port of tlamatini" in subject
        or "hard-cancel" in subject
        or "drag and drop image" in subject
        or "catalog of prompts" in subject
        or "nmapper" in subject
        or "nmap" in subject
        or "pentesting" in subject
        or "pentest" in subject
        or "cyber-sec" in subject
        or "cyber sec" in subject
        or "startup dialog" in subject
        or "catalog of prompts" in subject
        or "flowpills" in subject
        or "unrealer" in subject
        or "scaffold" in subject
        or "smoothness" in subject
        for subject in subjects
    )
    if has_current_release_wave:
        highlights.append(
            "The live Git window resolves cleanly to `v1.42.0`: tag, local `main`, and `origin/main` all point to commit `c58b01ad`. Git/source inspection extends README.md and BookOfTlamatini.md before the generated inventory tables derive the current agent, tool, skill, asset, and effective-line totals."
        )
        highlights.append(
            "The current release wave is broader than a badge bump: `v1.42.0` releases STM32er PlatformIO coverage, stepwise camera-verification demos, and the contiguous prompt-catalog migration; `v1.41.4` External-MCP `structuredContent`, `v1.41.3` catalog grouping, `v1.41.2` Hard Cancel, and `v1.41.0` screenshot paste/drop remain carried foundations."
        )
    if any("structuredcontent" in subject for subject in subjects):
        highlights.append(
            "External MCP stdio and network calls now deliver both human-readable content blocks and machine-readable `structuredContent` to the LLM, preventing valid structured-output servers from looking empty and triggering repeat-call cancellation."
        )
    if any("drag and drop image" in subject or "accept drag" in subject for subject in subjects):
        highlights.append(
            "The screenshot-to-chat path accepts clipboard bitmaps or dropped image files, re-encodes them safely into Tlamatini's Temp directory, inserts each absolute path at the remembered caret, and exposes removable thumbnail chips so Image-Interpreter can consume the same local path immediately."
        )
    if any("hard-cancel" in subject or "hard cancel" in subject for subject in subjects):
        highlights.append(
            "Hard Cancel now mints a monotonically increasing run epoch per user and permanently latches only the cancelled epoch; executor, retry, self-healing, Ask-Execs, status-emitter, and frontend guards stop the old run without poisoning the next request or another user's concurrent work."
        )
    if any("catalog of prompts" in subject or "prompt catalog" in subject for subject in subjects):
        highlights.append(
            "Catalog-of-Prompts migrations classify 106 historical rows into 13 operator categories and physically remove 13 duplicate ACPX demos without renumbering survivors; the primary endpoint is gap-tolerant, the fallback skips gaps, and the UI adds grouped sections plus ranked fuzzy search."
        )
    if any(
        "1.40.1" in subject
        or "django_port" in subject
        or "port of tlamatini" in subject
        or "configurable" in subject
        for subject in subjects
    ):
        highlights.append(
            "The `v1.40.1` configurable-web-port contract retires the hardcoded 8000: `config.json`'s `django_port` is resolved by `manage.py::_resolve_django_port()` and injected into every launch path by `_apply_configured_port()` (frozen double-click, `.flw` association, browser auto-open, source `runserver`, and `startserver`), so a machine where Windows has RESERVED port 8000 — the `WinError 10013` startup death that a frozen install previously could not escape without a rebuild — is fixed by editing one line. Resolution is fail-open to 8000 and an explicit command-line port always wins, both pinned by 24 tests in `agent/test_django_port_config.py`."
        )
    if any("flowpills" in subject or "companion" in subject for subject in subjects):
        highlights.append(
            "The `v1.40.0` companion-app contract adds `agent_manifest.py`, a six-value `HKCU\\Software\\XAIHT\\Tlamatini` discovery key, `_tlamatini_agents_manifest.json` with per-file SHA-256 values, preserved-agent uninstall metadata, launch/install/build integration, and 17 focused tests so FlowPills can locate valid agent templates without importing Tlamatini or scanning drives."
        )
    if any("unrealer" in subject or "scaffold" in subject for subject in subjects):
        highlights.append(
            "The Unrealer wave adds a two-field Catalog-of-Prompts route to a ready-to-build Unreal Engine 5.8 C++ project, including plugin wiring and Visual Studio 2026 guidance, while normalizing `/Content` paths to `/Game` and sending `assign_material` through the plugin's real `slot_index` wire key."
        )
    if any("smoothness" in subject for subject in subjects):
        highlights.append(
            "The `v1.39.5` smoothness pass bounds previously open-ended I/O, preserves partial Nmapper results, makes subprocess decoding UTF-8-safe, hardens background External MCP supervision, keeps deferred deliverables and request-scoped orphan evidence, and improves `.flw` secret redaction without weakening runtime behavior."
        )
    if any(
        "retrying behaviour" in subject
        or "self-healing" in subject
        or "self healing" in subject
        or "failure" in subject
        or "failures" in subject
        for subject in subjects
    ):
        highlights.append(
            "The newest committed reliability wave adds `agent/self_healing.py`: every Multi-Turn model step now runs under a watchdog, switches recovery tactics on transient failures, broadcasts live recovery status, and preserves already-executed agent work instead of discarding it."
        )
    if any("answer classifier" in subject or "answer_success" in subject or "create flow" in subject for subject in subjects):
        highlights.append(
            "Create Flow is no longer gated by the removed whole-answer SUCCESS/FAILURE classifier: the button appears when Multi-Turn has at least one successful agent call, and the generated `.flw` keeps only successful executions."
        )
    if any("pdcp" in subject or "projectdiscovery" in subject or "vulnx" in subject or "go-deny" in subject or "go deny" in subject for subject in subjects):
        highlights.append(
            "The newest committed Discoverer wave strengthens ProjectDiscovery operation: PDCP key setup is documented and auto-injected, legacy `cvemap` CVE searches now run through `vulnx`, and Go-toolchain source-control guardrails keep the private compiler/cache out of the repository."
        )
    if any(
        "video-analizer" in subject
        or "video-analyzer" in subject
        or "video analyzer" in subject
        or "video_analyzer" in subject
        for subject in subjects
    ):
        highlights.append(
            "The newest agent feature is Video-Analyzer: she reads recorded videos, gates obvious no-motion failures deterministically, runs two Ollama vision models in parallel, merges their reports, emits `INI_SECTION_VIDEO_ANALYZER`, and exposes substring-safe `TLM_VERDICT::<TOKEN>` values for Forker-driven robotic loops."
        )
    if any("search of prompts" in subject or "prompt search" in subject for subject in subjects):
        highlights.append(
            "Prompt discovery was upgraded too: the Tools dialog now supports substring, word-start, and fuzzy matching with mode badges, making the growing seeded prompt catalog easier to browse without memorizing exact titles."
        )
    if any("flw generation" in subject or ".flw generation" in subject for subject in subjects):
        highlights.append(
            "Generated `.flw` files now use a serpentine/boustrophedon canvas layout so long Multi-Turn chains stay on screen with continuous wiring instead of collapsing into unreadable off-canvas rows."
        )
    if any("zavuerer" in subject or "zavu" in subject for subject in subjects):
        highlights.append(
            "The newest working-tree feature is Zavuerer, the 83rd workflow-agent type: she sends authorized, opted-in messages through Zavu's unified API for SMS, WhatsApp, Telegram, Email, and Voice, with a wrapped `chat_agent_zavuerer`, Access Keys Wizard support, Parametrizer output fields, canvas connection handling, and demo/catalog migrations."
        )
    if any("image-interpreter" in subject or "image interpreter" in subject or "image_interpreter" in subject for subject in subjects):
        highlights.append(
            "The latest handbook/source delta upgrades Image-Interpreter into a triple-model vision pipeline: `qwen3.5:cloud` and `gemma4:cloud` interpret each image in parallel on dedicated Ollama connections, then `glm-5.2:cloud` merges both reports into one structured `INI_SECTION_IMAGE_INTERPRETER` result."
        )
    if any("config dialog" in subject or "config -> models" in subject or "config models" in subject for subject in subjects):
        highlights.append(
            "The current Config -> Models dialog now mirrors that triple-model reality with three Image-Interpreter fields: interpreter 1, interpreter 2, and image merger, plus safe defaults so older preserved configs are not stranded with empty required values."
        )
    if not has_current_release_wave and any(
        "1.26.5" in subject
        or "mit license" in subject
        or "mcp-doctor" in subject
        or "mcp doctor" in subject
        or "external mcp income" in subject
        or "unrealer parametrization" in subject
        or "game creation prompt" in subject
        or "drift of number of agents" in subject
        for subject in subjects
    ):
        highlights.append(
            "The live Git window now lands on `v1.26.5`: README, BookOfTlamatini, and the version surfaces identify the project as 1.26.5, with the license moved to MIT, the public Discord invite added, the live catalog corrected to 82 workflow agents, and the External MCP/MCP Doctor/Unrealer maintenance wave folded into the documentation."
        )
        highlights.append(
            "The `v1.26.5` code changes are maintenance-heavy but operator-visible: External MCP input handling and catalog cleanup were improved, MCP Doctor now enumerates all active external MCPs instead of stopping at the first, Unrealer parameterization was fixed, and a new Unreal game-creation prompt was seeded for guided demos."
        )
    if not has_current_release_wave and any("1.26.0" in subject or "external mcp" in subject or "esphomer" in subject for subject in subjects):
        highlights.append(
            "The current Git window still carries the `v1.26.0` External MCP baseline forward: a config-driven universal client connects to external MCP servers over stdio, streamable HTTP, SSE, or WebSocket, with MCP Doctor diagnostics, full-surface Multi-Turn tool binding, and Step-by-Step setup layered on top of the earlier ESPHomer firmware lane."
        )
    if any("installation/use steps" in subject or "getting start" in subject or "config images" in subject or "ollama token" in subject for subject in subjects):
        highlights.append(
            "The latest handbook pass on June 19, 2026 shifts the emphasis from release bragging to operator onboarding: BookOfTlamatini now opens with an easy-follow five-step setup path (install Tlamatini, install/sign in to Ollama, pull models, configure the app, start using Multi-Turn only when needed), and it adds first-run configuration screenshots plus a clearer local-vs-remote Ollama-token rule."
        )
    if any("disclaimer" in subject and "agent" in subject for subject in subjects):
        highlights.append(
            "The documentation now carries a clear Agent-directory disclaimer: workflow agents under `Tlamatini/agent/agents/` are plain-Python user-jurisdiction code, so the user who enables, edits, configures, chains, or runs them is responsible for their security boundary, credentials, targets, and downstream effects."
        )
    if has_esphomer_assets():
        highlights.append(
            "The tagged `v1.26.0` release now includes ESPHomer as a fourth firmware lane, bridging Tlamatini to ESPHome so she can author YAML device configs, validate, compile, upload, and observe smart-home firmware from chat or canvas."
        )
    if any("image about" in subject or "video" in subject or "kyber" in subject for subject in subjects):
        highlights.append(
            "The same window refreshes the visible about/presentation media too: `TlamatiniAbout.png` replaces the old JPEG asset, and `TlamatiniAndKyber.mp4` joins the repository as a new shipped visual asset that the dossier inventory now counts."
        )
    if any("v1.23.0" in subject or "data-preserving self-update" in subject or "migrate users' db on self-update" in subject for subject in subjects):
        highlights.append(
            "The current Git window includes the `v1.23.0` release from June 15, 2026: the in-app updater now preserves the user's database and migrates it back into the new build on first launch, so chat history and custom Tool/Mcp/Agent toggles survive a packaged update."
        )
    if any("numpy/opencv" in subject or ("embed" in subject and ("numpy" in subject or "opencv" in subject)) for subject in subjects):
        highlights.append(
            "The same `v1.23.0` wave hardens frozen builds for the media family: numpy and OpenCV are now embedded in both the carried Python and the frozen `_internal`, and `build.py` fails loudly if either import is missing instead of shipping a broken Recorder / Camcorder / AudioPlayer / VideoPlayer / Whisperer path."
        )
    if any("file-reading" in subject or "file-modification" in subject or "tool-order rule" in subject or "quoted args" in subject for subject in subjects):
        highlights.append(
            "The working tree also advances the file-navigation and file-editing operator surface: the new Globber, Grepper, and Editor agents/tools let Tlamatini discover files by pattern, search contents by regex, and make surgical in-place edits without falling back to a shell command."
        )
    if any("blenderer" in subject or "blender" in subject for subject in subjects):
        highlights.append(
            "The same recent release span still includes the Blenderer foundation: Tlamatini reaches a live Blender session through the official Blender MCP add-on socket (`localhost:9876`), both as the wrapped `chat_agent_blenderer` tool and as a visual workflow node."
        )
    if any("self-update" in subject or "check for updates" in subject or "apply_update.ps1" in subject or "start_update" in subject for subject in subjects):
        highlights.append(
            "The in-app self-update path itself is now mature across the current Git window: packaged installs can check GitHub releases, stage a download, hand the locked-file replacement to `apply_update.ps1`, and preserve both operator state and one `agents_backup` generation."
        )
    if has_current_release_wave or any(
        "1.36.0" in subject
        or "1.35.0" in subject
        or "1.33.2" in subject
        or "1.33.0" in subject
        or "1.32.0" in subject
        or "video-analizer" in subject
        or "video-analyzer" in subject
        or "video analyzer" in subject
        or "video_analyzer" in subject
        or "search of prompts" in subject
        or "prompt search" in subject
        or "flw generation" in subject
        or ".flw generation" in subject
        or "retrying behaviour" in subject
        or "self-healing" in subject
        or "self healing" in subject
        or "failure" in subject
        or "failures" in subject
        or "answer classifier" in subject
        or "answer_success" in subject
        or "create flow" in subject
        or "pdcp" in subject
        or "projectdiscovery" in subject
        or "vulnx" in subject
        or "go-deny" in subject
        or "go deny" in subject
        or "zavuerer" in subject
        or "zavu" in subject
        or "annouces skill" in subject
        or "announce skill" in subject
        or "angela" in subject
        or "creator" in subject
        or "private data" in subject
        or "public build" in subject
        or "public verify" in subject
        or "3x" in subject
        or "documentation" in subject
        or "docs" in subject
        or "disclaimer" in subject
        for subject in subjects
    ):
        highlights.append(
            "The latest dossier pass resolves the product at tagged `v1.42.0`, combines README.md and BookOfTlamatini.md with source/Git truth for the released STM32er PlatformIO expansion and carried External-MCP structured output, and retains the complete installation, Ollama, architecture, usage, tree, line inventory, and responsibility context."
        )
    elif not has_current_release_wave and any(
        "1.26.5" in subject
        or "mit license" in subject
        or "discord" in subject
        or "mcp-doctor" in subject
        or "mcp doctor" in subject
        or "drift of number of agents" in subject
        or "documentation" in subject
        or "docs" in subject
        or "disclaimer" in subject
        for subject in subjects
    ):
        highlights.append(
            "The latest documentation pass aligns the handbook and source with the `v1.26.5` line, which matters here because the dossier must reflect the live MIT license, Discord community link, 82-agent / 89-tool inventory, and the Agent-directory responsibility boundary instead of older public badges or stale prose."
        )
    elif any(
        "1.26.1" in subject
        or "disclaimer" in subject
        or ("documentation" in subject and ("1.26.1" in subject or "external mcp" in subject))
        for subject in subjects
    ):
        highlights.append(
            "The latest documentation pass aligns the handbook and source with the `v1.26.5` line, which matters here because some older badges or prose lines still lag behind the live 82-agent / 89-tool inventory and the new Agent-directory responsibility boundary."
        )
    if any("filecreator" in subject or "file creator" in subject or ("truncate" in subject and "file" in subject) for subject in subjects):
        highlights.append(
            "The `v1.19.5` File-Creator hardening pass now writes content byte-for-byte: plain `content` is re-extracted verbatim, and heavy escape/binary payloads can travel through `content_b64`, eliminating the wrong-symbol corruption that broke backslash-dense Java, JSON, and regex files."
        )
    if any("source code as the codebase" in subject or "self modify" in subject or "copy_source_assets" in subject or "source snapshot" in subject for subject in subjects):
        highlights.append(
            "The same release completes the self-modify story: `build.py --self-modify` now generates a rebuildable `TlamatiniSourceCode/` snapshot through `copy_source_assets.py`, so a self-able-modify build carries her own source tree, rebuild instructions, and redacted secrets in one honest package."
        )
    if any("api-keys wizard" in subject or "api keys wizard" in subject or ("api" in subject and "wizard" in subject) for subject in subjects):
        highlights.append(
            "Operator ergonomics improved too: the Config menu now includes an API-Keys Wizard dialog, giving the user a browser-side way to enter and persist provider credentials without hand-editing `config.json`."
        )
    if any("talker" in subject or "whisperer" in subject or "recorder" in subject or "camcorder" in subject or "audio" in subject for subject in subjects):
        highlights.append(
            "The recent media-and-voice wave is still visible in Git: Talker and Whisperer extend Tlamatini from text-only operation into female-voice text-to-speech plus speech-to-text, while Recorder/Camcorder/Shoter/AudioPlayer/VideoPlayer complete the broader media I/O family."
        )
    if any("watchdog" in subject or "hung command" in subject or "hanged" in subject for subject in subjects):
        highlights.append(
            "Runtime resilience advanced as well: the autonomous command watchdog can now detect shell wrappers that are alive but making no CPU or I/O progress, then reap only the wedged console-interpreter subtree without touching healthy long-running work."
        )
    if any("audioplayer" in subject or "videoplayer" in subject or ("playback" in subject and ("audio" in subject or "video" in subject)) for subject in subjects):
        highlights.append(
            "New media-playback pair completes the media-I/O family: AudioPlayer plays an audio file to the speakers (soundfile + sounddevice, volume and a truncate/loop time budget) and VideoPlayer plays a video file with audio on a chosen display (ffpyplayer — its wheel bundles ffmpeg + SDL — plus an OpenCV window, with display/volume/time-budget/window/fullscreen). Both are observational/output, on the canvas and as wrapped chat_agent_audioplayer / chat_agent_videoplayer tools."
        )
    if any("camcorder" in subject or "recorder" in subject for subject in subjects):
        highlights.append(
            "New observational capture pair: Camcorder (webcam photo/video via OpenCV) and Recorder (microphone WAV via sounddevice) — read-only siblings of Shoter, on the canvas and as wrapped chat_agent_camcorder / chat_agent_recorder tools."
        )
    if any("arduiner" in subject or "arduino" in subject for subject in subjects):
        highlights.append(
            "Arduiner added as the third microcontroller agent: a direct arduino-cli bridge that builds and uploads firmware for any fqbn-selected board, with zero-config bootstrap, auto board-core install, and a serial-port safety preflight."
        )
    if any("flow-making" in subject or "flow making" in subject or "flw" in subject for subject in subjects):
        highlights.append(
            "The new in-process flow-making skill turns a plain objective into a canvas-loadable .flw by driving the FlowCreator engine, so operators build runnable flows straight from chat."
        )
    if any("temp/template" in subject or "template generation" in subject for subject in subjects):
        highlights.append(
            "Directory policy: every transient file now stays under <app>/Temp and every scaffolded firmware/engine project tree under <app>/Templates (never C:/Temp or %TEMP%), pinned before Django starts."
        )
    if any("esp32er" in subject or "es32er" in subject or "platformio" in subject for subject in subjects):
        highlights.append(
            "The embedded-firmware surface remains broad in the current tree: ESP32er keeps the direct PlatformIO Core path for scaffold/build/upload/monitor work, STM32er and Arduiner cover the other hardware lanes, and the live working tree now adds ESPHomer as the ESPHome smart-home device bridge."
        )
    if any("asking on the chain of multi-turn" in subject or "ask exec" in subject or "execution interrupted" in subject for subject in subjects):
        highlights.append(
            "Human approval remains part of the modern safety story too: Ask Execs can still stop Multi-Turn before the next state-changing step, wait for a Proceed or Deny decision, and surface denial through the explicit red interruption banner."
        )
    if any("stm32" in subject or "stmer" in subject or "firmware" in subject and "hardware" in subject for subject in subjects):
        highlights.append(
            "The firmware-control branch remains active in the current codebase: STM32er still bridges the STM32 Template Project MCP for scaffold/build/flash/observe/reset flows, guarded by a fail-safe hardware preflight before unsafe mutations."
        )
    if any("herself" in subject or "self" in subject and "modify" in subject or "self knowledge" in subject for subject in subjects):
        highlights.append(
            "Self-awareness is now part of the steady baseline rather than a one-off milestone: Tlamatini carries a first-person self-knowledge map and, in self-modify builds, can inspect the bundled source snapshot that describes how she actually works."
        )
    if any("4096" in subject or "degrees of liberty" in subject or "turn" in subject and "freedom" in subject for subject in subjects):
        highlights.append(
            "Multi-Turn autonomy expanded too: the default iteration ceiling now reaches 4096 turns, giving long operator chains far more room before they hit the loop cap."
        )
    if any("unrealer" in subject or "unreal mcp" in subject or "xaiht" in subject for subject in subjects):
        highlights.append(
            "Unrealer kept growing across the same window: the docs now point at the public `XAIHT/XaihtUnrealEngineMCP` fork and the full 53-command, nine-category Unreal MCP surface it exposes to chat and canvas."
        )
    if any("kalier" in subject or "kali" in subject or "pentest" in subject for subject in subjects):
        highlights.append(
            "A still-relevant platform branch remains the Kali Linux bridge: Kalier behaves as the embedded MCP-Kali-Server client, so the Kali box URL is configured once in `Config -> URLs` and auto-injected into `chat_agent_kalier` runs."
        )
    if any("windower" in subject or "window manager" in subject or "window" in subject and "multi-turn" in subject for subject in subjects):
        highlights.append(
            "Another still-visible platform branch is Windower: a deterministic Win32 window-manager surface for focusing, tiling, resizing, listing, and closing windows from both Multi-Turn chat and the visual canvas."
        )
    if any("playwrighter" in subject or "playwright" in subject for subject in subjects):
        highlights.append(
            "Playwrighter remains part of the broader platform story too: a real-browser automation surface for scripted Playwright flows from both Multi-Turn chat and the visual canvas."
        )
    if any("tkinter" in subject or "unstable" in subject or "native dialog" in subject for subject in subjects):
        highlights.append(
            "The immediately previous stability pass is still visible in Git too: Tkinter was removed from the unstable runtime-facing dialog path in favor of native Windows dialog helpers, reducing UI instability around file and folder picking while keeping the browser/operator flow intact."
        )
    if any("reporting tables" in subject or "widths" in subject for subject in subjects):
        highlights.append(
            "Reporting-table layout was also tightened during the same window, improving readability of generated execution/reporting surfaces without changing the underlying operational data."
        )
    if any(("reviewer" in subject and "state" in subject) or ("reviewer" in subject and "handling" in subject) for subject in subjects):
        highlights.append(
            "A still-relevant reviewer follow-up is the behavioral-accuracy patch: the review prompt distinguishes uncommitted working-tree diffs from committed history and teaches the model Tlamatini’s managed-secret scrub convention, reducing false positives around local credentials in config files."
        )
    if any("reviewer" in subject or "analyzer" in subject or "security audit" in subject or "code review" in subject for subject in subjects):
        highlights.append(
            "A still-visible platform branch is Reviewer and Analyzer: code review plus deterministic security scanning are available from both the canvas and the skill layer."
        )
    if any("number and descriptions of agents" in subject or "markdowns" in subject or "agentic_skill" in subject for subject in subjects):
        highlights.append(
            "Agent-catalog consistency work also remains visible: the live count, the markdown bestiaries, the flow-creator skill catalog, and the sidebar-description source were brought back into alignment around one shared workflow-agent inventory."
        )
    if any("unreal" in subject or "unreal-engine mcp" in subject or "unreal engine enabled" in subject for subject in subjects):
        highlights.append(
            "Unreal MCP support remains part of the broader platform story: the Unrealer agent, its chat-wrapped tool, canvas wiring, seeded prompts, and the direct TCP bridge into a live Unreal Engine 5 editor."
        )
    if any("orphan" in subject or "cleanup" in subject or "sec/perf" in subject for subject in subjects):
        highlights.append(
            "Windows process hygiene also remains visible in recent history: a three-tier reaper, hardened detached spawn sites, ACPX process-tree termination, and user-visible survivor reporting when anything truly refuses to die."
        )
    if any("de-compresser" in subject or "de compresser" in subject for subject in subjects):
        highlights.append(
            "The archive-automation branch remains visible too: De-Compresser adds deterministic archive compression/decompression, Multi-Turn exposure, ACP canvas wiring, and the `py7zr` fallback path."
        )
    if any("version" in subject or "worldwide system" in subject for subject in subjects):
        highlights.append(
            "Release-identity work also remains relevant: the SemVer policy, git-tag sourcing, runtime version surfaces, and build-time embedding across the Windows artefacts now define how Tlamatini reports her version."
        )
    if any("menu db" in subject or "database" in subject or "browse buttons" in subject for subject in subjects):
        highlights.append(
            "Another still-relevant operator-facing branch is the DB dropdown: backup, Set DB staging for the next start-up, startup swap-in/rollback mechanics, and native Browse buttons on both dialogs."
        )
    if any("gpu" in subject or "autoload" in subject or "spining" in subject for subject in subjects):
        highlights.append(
            f"GPU-host behavior changed during the {RECENT_GIT_WINDOW_LABEL}: performance hooks, model-pinning startup behavior, and autoload-at-restart reliability were all touched."
        )
    if any("reconnection" in subject or ("config" in subject and "dialog" in subject) for subject in subjects):
        highlights.append(
            "The config-plane UI grew another safety layer: when model/URL dialog saves materially change live runtime inputs, the chat now prompts the operator to reconnect before trusting the current session state."
        )
    if any("acpx" in subject for subject in subjects):
        highlights.append(
            f"ACPX-related work remained visible during the {RECENT_GIT_WINDOW_LABEL}: runtime, documentation, or operator-surface changes were still part of the active maintenance stream."
        )
    if any("shortcut" in subject or "restrictive" in subject or "policy" in subject for subject in subjects):
        highlights.append(
            "Windows deployment hardening continued with CreateShortcut fixes and improved behavior on restricted-policy machines."
        )
    if any("teletlamatini" in subject for subject in subjects):
        highlights.append(
            "TeleTlamatini grew throughout the window, including agent/runtime updates, config movement, FlowCreator/ACPX integration, and related documentation changes."
        )
    if any("compiler" in subject or "contract" in subject for subject in subjects):
        highlights.append(
            "The current Git window still shows the flow compiler / agent-contract direction influencing how chat-created and ACP-created workflows converge."
        )
    if any("multi-turn" in subject or "multi turn" in subject or "tool quota" in subject for subject in subjects):
        highlights.append(
            "Multi-Turn behavior kept evolving across the period through quota tuning, execution-table persistence, autonomous-action improvements, and broader tool enablement."
        )
    if any("glm-5.2:cloud" in subject or "default in config.json" in subject for subject in subjects):
        highlights.append(
            "The checked-in runtime defaults also moved: the shared config now points at `glm-5.2:cloud`, so the handbook and dossier need to describe the shipped cloud-first baseline honestly instead of assuming only the older local model defaults."
        )
    if any("attention" in subject or "flash" in subject or "notifications" in subject or "notifier" in subject for subject in subjects):
        highlights.append(
            "Operator attention routing also changed: when Ask Execs or a Notifier event needs the user, Tlamatini can now flash her own Windows taskbar presence and write an uppercase attention banner into `tlamatini.log`."
        )
    if any("pythonxer" in subject or "forked windows execution" in subject or "reporting on the log file" in subject or "project skills" in subject for subject in subjects):
        highlights.append(
            "Follow-up implementation work in the same post-STM32 window tightened Pythonxer downstream execution, Windows forked-command launching, execution-log detail, and project-skill loading, so the release story is not only a UI toggle but also a reliability pass around the operator chain."
        )
    if any("scroll" in subject or "icon" in subject or "web page" in subject or "scheme" in subject for subject in subjects):
        highlights.append(
            "The operator surface evolved too: ACP/ACPX visual mechanics, canvas scrolling, icons, and page framing all received polish."
        )
    if any("persistent" in subject or "execution table" in subject or "summarizer" in subject for subject in subjects):
        highlights.append(
            "Execution observability improved through persistent execution tables, summarizer work, and the broader reportability push around agent runs."
        )
    if any("document" in subject or "docs" in subject or "framing" in subject for subject in subjects):
        highlights.append(
            "Documentation itself changed during the recent window, so the regenerated dossier is part of the tracked operator surface rather than an external afterthought."
        )
    if highlights:
        return highlights
    return [f"Git history shows focused maintenance across the operator surface, release mechanics, and runtime behavior during {RECENT_GIT_WINDOW_LABEL}."]


def visual_doc_highlights(commits: list[CommitInfo]) -> list[str]:
    subjects = [commit.subject.lower() for commit in commits]
    highlights: list[str] = []
    has_current_release_wave = any(
        "1.36.0" in subject
        or "1.35.0" in subject
        or "1.33.2" in subject
        or "1.33.0" in subject
        or "1.32.0" in subject
        or "video-analizer" in subject
        or "video-analyzer" in subject
        or "video analyzer" in subject
        or "video_analyzer" in subject
        or "search of prompts" in subject
        or "prompt search" in subject
        or "flw generation" in subject
        or ".flw generation" in subject
        or "retrying behaviour" in subject
        or "self-healing" in subject
        or "self healing" in subject
        or "failure" in subject
        or "failures" in subject
        or "answer classifier" in subject
        or "answer_success" in subject
        or "create flow" in subject
        or "pdcp" in subject
        or "projectdiscovery" in subject
        or "vulnx" in subject
        or "go-deny" in subject
        or "go deny" in subject
        or "zavuerer" in subject
        or "zavu" in subject
        or "annouces skill" in subject
        or "announce skill" in subject
        or "angela" in subject
        or "creator" in subject
        or "author" in subject
        or "private data" in subject
        or "public build" in subject
        or "public verify" in subject
        or "3x" in subject
        or "image-interpreter" in subject
        or "image interpreter" in subject
        or "image_interpreter" in subject
        or "config dialog" in subject
        or "1.39.5" in subject
        or "1.40.0" in subject
        or "1.40.1" in subject
        or "1.41.0" in subject
        or "1.41.2" in subject
        or "1.41.3" in subject
        or "1.41.4" in subject
        or "1.42.0" in subject
        or "stm32er" in subject
        or "structuredcontent" in subject
        or "django_port" in subject
        or "port of tlamatini" in subject
        or "configurable" in subject
        or "hard-cancel" in subject
        or "drag and drop image" in subject
        or "catalog of prompts" in subject
        or "flowpills" in subject
        or "unrealer" in subject
        or "scaffold" in subject
        or "smoothness" in subject
        for subject in subjects
    )
    if has_current_release_wave:
        highlights.append(
            "The `v1.42.0` release at `c58b01ad` promotes the STM32er PlatformIO expansion, camera-verified stepwise demos, and contiguous prompt-catalog migration into the tagged release; local `main` and `origin/main` are aligned on that commit."
        )
        highlights.append(
            "The release delta expands `stm32er.py`, config, registry, contracts, tools, tests, descriptions, migrations, and proposal assets together; the existing local configuration-only worktree changes remain private and are neither reproduced nor modified by dossier generation."
        )
    if any("structuredcontent" in subject for subject in subjects):
        highlights.append(
            "The MCP formatter preserves plain text and error behavior, unwraps a sole `{result: ...}` envelope, serializes structured payloads safely, and caps oversized structured content at a configurable character budget."
        )
    if any("drag and drop image" in subject or "accept drag" in subject for subject in subjects):
        highlights.append(
            "Image ingestion is path-native rather than attachment-native: Pillow flattens/re-encodes images under a 25 MB ceiling, the Temp path is inserted at the caret, chips can remove both preview and text, and layout observers account for the new row so the textarea and Send control remain visible."
        )
    if any("hard-cancel" in subject or "hard cancel" in subject for subject in subjects):
        highlights.append(
            "The cancellation delta closes every resurrection path: the per-user epoch survives legacy-flag clearing, propagates through payload rebuilds, stops executor/retry/self-healing loops, denies blocked Ask-Execs prompts, revokes late tactic emitters, and prevents stale frontend frames from putting the UI back into Cancel state."
        )
    if any("catalog of prompts" in subject or "prompt catalog" in subject for subject in subjects):
        highlights.append(
            "Prompt-catalog migrations group the historical catalog into 13 categories and delete 13 redundant ACPX variants while keeping surviving ids stable; the UI renders category sections, supports numeric/acronym/subsequence search, highlights matches, and restores grouped order when the query clears."
        )
    if any(
        "retrying behaviour" in subject
        or "self-healing" in subject
        or "self healing" in subject
        or "failure" in subject
        or "failures" in subject
        for subject in subjects
    ):
        highlights.append(
            "The newest committed reliability wave adds `agent/self_healing.py`, routes Multi-Turn model calls through `SelfHealingInvoker`, broadcasts live recovery status, trims or retries context when model calls fail, and preserves already-executed tool evidence for a degraded but truthful final answer."
        )
    if any("answer classifier" in subject or "answer_success" in subject or "create flow" in subject for subject in subjects):
        highlights.append(
            "The old whole-answer `answer_success` classifier was removed: Create Flow now appears when Multi-Turn produced at least one successful agent call, generated `.flw` files keep only successful executions, and Exec Report remains tool-evidence rather than a global verdict."
        )
    if any("pdcp" in subject or "projectdiscovery" in subject or "vulnx" in subject or "go-deny" in subject or "go deny" in subject for subject in subjects):
        highlights.append(
            "The latest committed visual-doc delta now includes Discoverer hardening too: PDCP key retrieval/setup, `cvemap` routed to ProjectDiscovery `vulnx`, the latest-CVE demo prompt, the `git_deny_go.py` guard, and the no-Go-toolchain-in-source-control rule."
        )
    if any(
        "video-analizer" in subject
        or "video-analyzer" in subject
        or "video analyzer" in subject
        or "video_analyzer" in subject
        for subject in subjects
    ):
        highlights.append(
            "Video-Analyzer adds a motion-verdict surface for hardware-in-the-loop demos: Camcorder can record a board, Video-Analyzer can judge the motion with deterministic and model-backed checks, and Forker can branch on `TLM_VERDICT::PASS_OK`, `FAIL_NO_MOTION`, `FAIL_WRONG_MOTION`, `UNCLEAR`, or `ANALYSIS_ERROR`."
        )
    if any("search of prompts" in subject or "prompt search" in subject for subject in subjects):
        highlights.append(
            "Prompt sorting/search work now makes the seeded prompt catalog easier to operate from the Tools dialog: substring, word-start, fuzzy scoring, and mode badges reduce hunting in a large prompt list."
        )
    if any("flw generation" in subject or ".flw generation" in subject for subject in subjects):
        highlights.append(
            "The latest `.flw` generator work changes layout, not just metadata: large generated flows use alternating row direction and row capacity rules so visual workflows avoid unreadable overlap and off-screen chains."
        )
    if any("zavuerer" in subject or "zavu" in subject for subject in subjects):
        highlights.append(
            "The fresh `v1.33.2` artifact delta centers on Zavuerer and cleanup: README.md, BookOfTlamatini.md, `agents_descriptions.md`, `agent/agents/zavuerer/`, migrations `0159`-`0164`, Access Keys Wizard wiring, capability hints, model/config defaults, and frontend/runtime cleanup all describe the same unified-messaging release family."
        )
    if any("image-interpreter" in subject or "image interpreter" in subject or "image_interpreter" in subject for subject in subjects):
        highlights.append(
            "Since the last committed PDF/PPTX refresh, Image-Interpreter was rebuilt into a fail-safe triple-model analyzer: two dedicated Ollama vision interpreters run in parallel, a barrier waits for both, and a merger model fuses the reports while preserving raw-output fallbacks when any leg fails."
        )
    if any("config dialog" in subject or "config -> models" in subject or "config models" in subject for subject in subjects):
        highlights.append(
            "The post-refresh UI/backend delta also updates Config -> Models so the browser exposes all three Image-Interpreter model slots and `tools.py` seeds the matching wrapped-tool defaults before per-call overrides."
        )
    if not has_current_release_wave and any(
        "1.26.5" in subject
        or "mit license" in subject
        or "mcp-doctor" in subject
        or "mcp doctor" in subject
        or "external mcp income" in subject
        or "unrealer parametrization" in subject
        or "game creation prompt" in subject
        or "drift of number of agents" in subject
        for subject in subjects
    ):
        highlights.append(
            "Since the last committed PDF/PPTX refresh, the repository advanced to `v1.26.5`: the handbook and README now carry the MIT license identity, Discord invite, corrected 82-agent catalog, and the updated External MCP/MCP Doctor/Unrealer maintenance story."
        )
        highlights.append(
            "The post-baseline implementation delta includes improved External MCP input/catalog handling, MCP Doctor enumeration across all active external MCPs, Unrealer parameterization fixes, seeded Unreal game-demo prompts, Create Superuser wizard prompt migrations, FlowCreation backslash behavior fixes, Step-by-Step polish, and context-bloat/chat-history windowing work."
        )
    if not has_current_release_wave and any("1.26.0" in subject or "external mcp" in subject or "esphomer" in subject for subject in subjects):
        highlights.append(
            "The broader `v1.26.x` line still centers on the External MCPs universal client: four transports, the MCP Doctor agent, full-surface Multi-Turn tool binding, and Step-by-Step setup layered on top of the ESPHomer firmware lane."
        )
    if any("installation/use steps" in subject or "getting start" in subject or "config images" in subject or "ollama token" in subject for subject in subjects):
        highlights.append(
            "The newest post-tag documentation wave is about usability rather than another runtime feature: BookOfTlamatini now begins with a five-step getting-started path, clarifies installer-vs-source choices, adds configuration screenshots, and explains that a localhost Ollama usually needs no token while a remote Ollama may."
        )
    if any("disclaimer" in subject and "agent" in subject for subject in subjects):
        highlights.append(
            "The latest Markdown source now explicitly warns that the plain-Python workflow agents are under user jurisdiction: Tlamatini offers orchestration and guardrails, but user-run agents can touch files, shells, APIs, credentials, external MCPs, hardware, and networks selected by the operator, so resulting breaches or unsafe actions remain the operator's responsibility."
        )
    if any("esphomer" in subject or "esphome" in subject for subject in subjects):
        highlights.append(
            "Since the last dossier baseline, ESPHomer changed across its ESPHome firmware agent, wrapped `chat_agent_esphomer` tool, sample YAML project, migrations, tests, and handbook coverage for smart-home device provisioning."
        )
    if any("image about" in subject or "video" in subject or "kyber" in subject for subject in subjects):
        highlights.append(
            "The same span also refreshes the shipped visual media: the old `TlamatiniAbout.jpg` gives way to `TlamatiniAbout.png`, and `TlamatiniAndKyber.mp4` is now part of the repository asset tree and line/inventory context."
        )
    if any("v1.23.0" in subject or "data-preserving self-update" in subject or "migrate users' db on self-update" in subject for subject in subjects):
        highlights.append(
            "Since the last committed PDF/PPTX refresh, `v1.23.0` made packaged self-update data-preserving: the user's database is staged through `DB/ToLoad/`, restored into the new build, and migrated on next launch so chat history and custom toggles survive the upgrade."
        )
    if any("numpy/opencv" in subject or ("embed" in subject and ("numpy" in subject or "opencv" in subject)) for subject in subjects):
        highlights.append(
            "The same window also embedded numpy and OpenCV into both shipped Python runtimes, closing the frozen-build dependency gap for Recorder, Camcorder, AudioPlayer, VideoPlayer, and Whisperer."
        )
    if any("file-reading" in subject or "file-modification" in subject or "tool-order rule" in subject or "quoted args" in subject for subject in subjects):
        highlights.append(
            "The operator surface also expanded with the file-navigation/file-edit trio: Globber, Grepper, and Editor now exist as workflow agents and wrapped chat tools, giving Tlamatini deterministic file discovery, regex search, and surgical in-place edit steps."
        )
    if any("blenderer" in subject or "blender" in subject for subject in subjects):
        highlights.append(
            "Since the last committed PDF/PPTX refresh, Blenderer entered the platform: a live Blender bridge over the official MCP add-on socket, available both on the canvas and as `chat_agent_blenderer`."
        )
    if any("self-update" in subject or "check for updates" in subject or "apply_update.ps1" in subject or "start_update" in subject for subject in subjects):
        highlights.append(
            "The same refresh window also delivered the in-app self-update path: `self_update.py`, new update endpoints, staged release downloads, and the external `apply_update.ps1` swap helper that preserves operator state during upgrade."
        )
    if not has_current_release_wave and any(
        "1.26.1" in subject
        or "disclaimer" in subject
        or "external mcp" in subject
        or ("documentation" in subject and "1.26.1" in subject)
        for subject in subjects
    ):
        highlights.append(
            "The latest versioning/documentation commits move the source-of-truth product story to the `v1.26.5` line: the External MCPs universal client remains the headline runtime capability, and the new Agent-directory disclaimer now makes the user-jurisdiction boundary explicit beside the current 82-agent / 89-tool runtime surface."
        )
    if any("filecreator" in subject or "file creator" in subject or ("truncate" in subject and "file" in subject) for subject in subjects):
        highlights.append(
            "Since the last committed PDF/PPTX refresh, `v1.19.5` hardened File-Creator with a byte-exact write path: verbatim `content` plus a `content_b64` channel now preserve backslash-heavy and binary payloads without wrong-symbol corruption."
        )
    if any("source code as the codebase" in subject or "self modify" in subject or "copy_source_assets" in subject or "source snapshot" in subject for subject in subjects):
        highlights.append(
            "The same release also completes the self-modify packaging story: `build.py --self-modify` now uses `copy_source_assets.py` to generate a rebuildable `TlamatiniSourceCode/` snapshot with redacted secrets, omitted heavy media, and rebuild instructions carried beside the application."
        )
    if any("api-keys wizard" in subject or "api keys wizard" in subject or ("api" in subject and "wizard" in subject) for subject in subjects):
        highlights.append(
            "A new operator-facing convenience layer landed too: Config now exposes an API-Keys Wizard dialog so cloud-provider credentials can be entered and updated from the browser instead of by manually editing `config.json`."
        )
    if any("watchdog" in subject or "hung command" in subject or "hanged" in subject for subject in subjects):
        highlights.append(
            "Recent runtime hardening added the autonomous command watchdog: a boot-time daemon thread that samples CPU and I/O progress across shell-interpreter subtrees and reaps only the genuinely wedged ones, closing the gap left by timeouts and post-return orphan cleanup."
        )
    if any("audioplayer" in subject or "videoplayer" in subject or ("playback" in subject and ("audio" in subject or "video" in subject)) for subject in subjects):
        highlights.append(
            "Since the last committed PDF/PPTX refresh, the media-PLAYBACK pair landed and completed the media-I/O family: AudioPlayer plays an audio file to the speakers (soundfile + sounddevice — volume in percent and a time-played budget that truncates a longer file or loops a shorter one), and VideoPlayer plays a video file with audio on a chosen display (ffpyplayer, whose wheel bundles ffmpeg + SDL, plus an OpenCV window — display, volume, the same time budget, window size, and fullscreen). Both are observational/output and ship on the canvas and as wrapped chat_agent_audioplayer / chat_agent_videoplayer tools."
        )
    if any("camcorder" in subject or "recorder" in subject for subject in subjects):
        highlights.append(
            "New observational capture pair: Camcorder (webcam photo/video via OpenCV) and Recorder (microphone WAV via sounddevice) — read-only siblings of Shoter, on the canvas and as wrapped chat_agent_camcorder / chat_agent_recorder tools."
        )
    if any("arduiner" in subject or "arduino" in subject for subject in subjects):
        highlights.append(
            "Arduiner added as the third microcontroller agent: a direct arduino-cli bridge that builds and uploads firmware for any fqbn-selected board, with zero-config bootstrap, auto board-core install, and a serial-port safety preflight."
        )
    if any("flow-making" in subject or "flow making" in subject or "flw" in subject for subject in subjects):
        highlights.append(
            "The new in-process flow-making skill turns a plain objective into a canvas-loadable .flw by driving the FlowCreator engine, so operators build runnable flows straight from chat."
        )
    if any("temp/template" in subject or "template generation" in subject for subject in subjects):
        highlights.append(
            "Directory policy: every transient file now stays under <app>/Temp and every scaffolded firmware/engine project tree under <app>/Templates (never C:/Temp or %TEMP%), pinned before Django starts."
        )
    if any("esp32er" in subject or "es32er" in subject or "platformio" in subject for subject in subjects):
        highlights.append(
            "The embedded-firmware branch remains part of the current product story too: ESP32er keeps the PlatformIO path for scaffold/build/upload/monitor work, while the broader handbook now needs to present the firmware stack as STM32er + ESP32er + Arduiner + ESPHomer rather than a smaller trio."
        )
    if any("asking on the chain of multi-turn" in subject or "ask exec" in subject or "execution interrupted" in subject for subject in subjects):
        highlights.append(
            "The human-in-the-loop gate remains current: Ask Execs can still pause before each state-changing Multi-Turn step, wait for Proceed or Deny through `ExecPermissionBroker`, and halt the chain safely with the explicit red interruption banner."
        )
    if any("stm32" in subject or "stmer" in subject or "firmware" in subject and "hardware" in subject for subject in subjects):
        highlights.append(
            "The STM32 branch is still active in the same modern operator surface: STM32er remains the bridge into the STM32 Template Project MCP for scaffold/build/flash/observe/reset workflows guarded by a critical preflight."
        )
    if any("self" in subject and ("herself" in subject or "modify" in subject or "source code" in subject) for subject in subjects):
        highlights.append(
            "The self-knowledge and optional self-modification surface has matured into the current baseline: she can describe her runtime honestly and, when `TlamatiniSourceCode/` is present, inspect the bundled rebuildable source tree that explains how she works."
        )
    if any("4096" in subject or "degrees of liberty" in subject or "turn" in subject and "freedom" in subject for subject in subjects):
        highlights.append(
            "The same window raised the default Multi-Turn iteration ceiling from 256 to 4096, making long autonomous operator chains practical without an early loop-cap failure."
        )
    if any("unrealer" in subject or "unreal mcp" in subject or "xaiht" in subject for subject in subjects):
        highlights.append(
            "Unrealer advanced substantially too: the docs now point at the public `XAIHT/XaihtUnrealEngineMCP` fork and the full 53-command Unreal MCP surface, with new demos and better parameter-path guidance."
        )
    if any("kalier" in subject or "kali" in subject for subject in subjects):
        highlights.append(
            "Kalier also matured during the same span: `v1.7.1` made Tlamatini the embedded MCP-Kali-Server client for chat-side runs, so operators configure the Kali box once in `Config -> URLs` instead of repeating it in every prompt."
        )
    if any("glm-5.2:cloud" in subject or "default in config.json" in subject or "pythonxer" in subject or "forked windows execution" in subject or "project skills" in subject or "reporting on the log file" in subject for subject in subjects):
        highlights.append(
            "The same span also refined the shipped operating baseline: handbook simplification, a `glm-5.2:cloud` checked-in default, stronger execution logging, Pythonxer downstream fixes, Windows forked-process polish, and cleaner project-skill loading."
        )
    if any("attention" in subject or "flash" in subject or "notifications" in subject or "notifier" in subject for subject in subjects):
        highlights.append(
            "Since the last committed PDF/PPTX refresh, operator attention handling moved to a concrete Windows path: browser-side Ask Execs and Notifier events can call `/agent/flash_window/`, which lets Tlamatini flash her own taskbar button and persist an uppercase attention banner in `tlamatini.log`."
        )
    if any("esp32 template project" in subject or "template project" in subject and "esp32" in subject for subject in subjects):
        highlights.append(
            "The same span also documented the `ESP32TemplateProject` reference repository: a plain PlatformIO project that gives ESP32er a known-good, GitHub-ready firmware baseline for build, upload, and serial-monitor verification."
        )
    if any("doc" in subject or "markdown" in subject or "graphical" in subject for subject in subjects):
        highlights.append(
            "The markdown handbooks themselves were revised during that span, so this dossier refresh is carrying forward not only code/runtime changes but also the corrected operator-language and release-story wording."
        )
    if highlights:
        return highlights
    return ["Since the last committed PDF/PPTX refresh, Git shows focused platform evolution across autonomy, operator ergonomics, runtime self-knowledge, and documentation fidelity."]


def _load_module_from_path(module_name: str, path: Path):
    spec = importlib.util.spec_from_file_location(module_name, str(path))
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Could not load module spec from {path}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def resolve_version_info() -> dict[str, str]:
    # Honor an explicit TLAMATINI_VERSION override, exactly like the build
    # scripts do (versioning.py precedence #2). This lets a release be
    # documented before its git tag is cut — set TLAMATINI_VERSION=X.Y.Z when
    # running the generator. With no override it falls through to git tags.
    override = os.environ.get("TLAMATINI_VERSION", "").strip()
    if override:
        return {"version": override, "build": override, "commit": "override", "date": "", "source": "env"}
    version_module_path = PROJECT_DIR / "agent" / "version.py"
    try:
        module = _load_module_from_path("tlamatini_agent_version", version_module_path)
        info = module.get_version_info()
        if isinstance(info, dict):
            return {
                "version": str(info.get("version", "0.0.0+unknown")),
                "build": str(info.get("build", "0.0.0+unknown")),
                "commit": str(info.get("commit", "unknown")),
                "date": str(info.get("date", "")),
                "source": str(info.get("source", "unknown")),
            }
    except Exception:
        pass
    return {
        "version": "0.0.0+unknown",
        "build": "0.0.0+unknown",
        "commit": "unknown",
        "date": "",
        "source": "unknown",
    }


def workflow_agents() -> list[str]:
    agents_root = PROJECT_DIR / "agent" / "agents"
    names = []
    for entry in agents_root.iterdir():
        if entry.is_dir() and (entry / "config.yaml").exists() and entry.name != "pools":
            names.append(entry.name)
    return sorted(names)


def count_requirements() -> int:
    count = 0
    for line in (REPO_ROOT / "requirements.txt").read_text(encoding="utf-8").splitlines():
        stripped = line.strip()
        if stripped and not stripped.startswith("#"):
            count += 1
    return count


def count_agent_description_rows() -> int:
    path = REPO_ROOT / "agents_descriptions.md"
    count = 0
    for line in path.read_text(encoding="utf-8").splitlines():
        stripped = line.strip()
        if stripped.startswith("| **") and "| Agent |" not in stripped and "| Purpose |" not in stripped:
            count += 1
    return count


def count_wrapped_chat_agent_tools() -> int:
    path = PROJECT_DIR / "agent" / "chat_agent_registry.py"
    text = path.read_text(encoding="utf-8")
    return len(re.findall(r'tool_name\s*=\s*"([^"]+)"', text))


def count_skills() -> int:
    skills_root = PROJECT_DIR / "agent" / "skills_pkg"
    return sum(1 for entry in skills_root.iterdir() if entry.is_dir() and (entry / "SKILL.md").exists())


def collect_context() -> dict:
    tracked = git_tracked_paths()
    untracked = git_untracked_paths()
    paths = inventory_paths()
    tree_text = build_tree(paths)
    language_rows, file_rows, binary_count, skipped_count = line_stats_for_paths(paths)
    total_effective = sum(row.effective_lines for row in language_rows)
    total_lines = sum(row.total_lines for row in language_rows)
    agents = workflow_agents()
    wrapped_chat_tools = count_wrapped_chat_agent_tools()
    skills_count = count_skills()
    reference_media = extract_reference_media()
    weekly = recent_week_commits()
    visual_baseline = last_visual_doc_commit()
    visual_commits = commits_since_visual_docs(visual_baseline)
    version_info = resolve_version_info()

    context = {
        "generated_at": local_stamp(),
        "head_short": git("rev-parse", "--short", "HEAD"),
        "head_full": git("rev-parse", "HEAD"),
        "head_subject": git("show", "-s", "--format=%s", "HEAD"),
        "head_date": git("show", "-s", "--format=%cI", "HEAD"),
        "inventory_files": len(paths),
        "tracked_files": len(tracked),
        "untracked_files": len(untracked),
        "tracked_paths": tracked,
        "untracked_paths": untracked,
        "inventory_paths": paths,
        "tree_text": tree_text,
        "language_rows": language_rows,
        "file_rows": file_rows,
        "binary_count": binary_count,
        "skipped_count": skipped_count,
        "total_effective_lines": total_effective,
        "total_lines": total_lines,
        "workflow_agents": agents,
        "workflow_agent_count": len(agents),
        "agent_description_rows": count_agent_description_rows(),
        "wrapped_chat_agent_count": wrapped_chat_tools,
        "core_python_tool_count": 20,
        "acpx_tool_count": 12,
        "total_multi_turn_tools": 20 + wrapped_chat_tools + 12,
        "skills_count": skills_count,
        "requirements_count": count_requirements(),
        "js_modules": len(list((PROJECT_DIR / "agent" / "static" / "agent" / "js").glob("*.js"))),
        "css_files": len(list((PROJECT_DIR / "agent" / "static" / "agent" / "css").glob("*.css"))),
        "html_templates": len(list((PROJECT_DIR / "agent" / "templates" / "agent").glob("*.html"))),
        "migrations": len(list((PROJECT_DIR / "agent" / "migrations").glob("*.py"))) - 1,
        "recent_commits": recent_commits(),
        "weekly_commits": weekly,
        "weekly_highlights": weekly_highlights(weekly),
        "visual_doc_baseline": visual_baseline,
        "visual_doc_commits": visual_commits,
        "visual_doc_highlights": visual_doc_highlights(visual_commits),
        "reference_media": reference_media,
        "version_info": version_info,
    }
    return context


SYSTEM_OVERVIEW = [
    "Tlamatini is a self-hosted AI developer assistant (cloud LLMs by default; the app and RAG run locally) built with Django, Django Channels, LangChain, LangGraph, FAISS/BM25 retrieval, and a large in-repository agent application.",
    "She combines a browser chat surface, a Retrieval-Augmented Generation stack, a Multi-Turn tool executor, MCP-backed context providers, wrapped chat-agent runtimes, and a visual Agentic Control Panel for workflow design.",
    "She is designed for development operations: codebase analysis, file and directory context, deterministic file discovery/search/editing, command execution, Python execution, screenshots, web/search helpers, notifications and attention routing, DevOps tools, authorized cyber-security assessment, local model operation, Windows packaging and uninstall registration, first-person self-knowledge about her own runtime, and embedded-firmware control for STM32, ESP32-class, Arduino-class, and ESPHome smart-home boards.",
]

AGENT_DIRECTORY_DISCLAIMER = [
    "Every agent in `Tlamatini/agent/agents/` is intentionally plain Python so the user can read, audit, edit, restrict, or disable its operating code. This transparency is a user-control mechanism, not a warranty that an agent is secure or suitable for a particular environment.",
    "Agents have no independent authority or jurisdiction. The user alone decides whether, where, how, and with which permissions an agent runs; enabling, configuring, modifying, chaining, or executing it places that execution under the user's control and jurisdiction.",
    "The user is responsible for code/config review, least-privilege secrets and credentials, authorized files and targets, browsers, shells, APIs, external MCPs, machines, hardware, downstream systems, supervision, and compliance with applicable law, policy, license, contract, and authorization.",
    "By running an agent, the user accepts responsibility for its actions and consequences. To the fullest extent permitted by applicable law, security breaches, data exposure or loss, unauthorized actions, credential leaks, unsafe automation, violations, compromise, device damage, financial loss, or other harm arising from use are the responsibility of the user who runs it.",
    "Tlamatini's orchestration, documentation, examples, and guardrails do not authorize third-party access and cannot replace the user's security review, permission controls, monitoring, or legal compliance.",
]

WHAT_IT_DOES = [
    "Answers codebase questions with loaded file or directory context.",
    "Uses hybrid retrieval to extract metadata, split content, rank source chunks, and respect context budgets.",
    "Can discover files by glob pattern, search their contents by regex, and make surgical in-place replacements through the Globber, Grepper, and Editor agent/tool trio.",
    "Can connect to external MCP servers declared in `external_mcps.json`, expose their remote tools to Multi-Turn under the `ext__<server>__<tool>` naming convention, and supervise those connections through status / reconnect / doctor / import / list / call helpers.",
    "Gives operators GUI-first database maintenance through the new DB dropdown for backup and staged database replacement.",
    "Lets operators manage provider secrets from the browser through the Config -> Access Keys Wizard instead of hand-editing `config.json`.",
    "Can drive Blender through the Blenderer agent, using the official Blender MCP add-on socket to inspect scenes, mutate objects and materials, run raw code, and automate renders from chat or the workflow canvas.",
    "Can pause before every state-changing Multi-Turn execution and ask the operator to approve or deny that exact step through the Ask Execs checkbox.",
    "Can raise a Windows attention signal when the browser needs the operator: Ask Execs prompts and Notifier events can flash Tlamatini’s own taskbar presence and leave an uppercase banner in `tlamatini.log`.",
    "Registers packaged installs in Windows `Installed apps` / `Programs and Features` with a real uninstall entry, so the release behaves like a normal installed application instead of only a shortcut bundle.",
    "Can update packaged installs in place through About -> Check for updates: she checks the latest GitHub release, stages the download, swaps locked files externally, and preserves operator state such as `config.json`, the database, and content.",
    "Hardens generated files: File-Creator’s bulk-write path now preserves long content byte-complete even when it contains heavy quoting or semicolon-rich source text.",
    "Warns GPU-host operators before a directory-context load is likely to saturate VRAM and degrade embedding throughput.",
    "Exposes a coherent versioning surface across builds, runtime UI, logs, and an open health-check endpoint.",
    "Carries a first-person self-knowledge map so she can answer more accurately about her own architecture, ports, runtime modes, pages, and capabilities.",
    "Can command Kali Linux offensive-security tooling through MCP-Kali-Server for authorized recon, enumeration, web scanning, and assessment workflows.",
    "Can run local authorized nmap reconnaissance through Nmapper, a use-only bridge that resolves a user-installed nmap, defaults to unprivileged TCP connect scanning, refuses unsafe/missing prerequisites gracefully, and never bundles or redistributes nmap.",
    "Can diagnose an external MCP before the first live connection through the MCP Doctor agent and wrapped `chat_agent_mcp_doctor` tool, checking transport, runtime requirements, PATH availability, placeholder secrets, and the next operator step.",
    "Can scaffold, author, build, flash, reset, and observe STM32 firmware through STM32er: the template-MCP path remains for STM32F407, while released `v1.42.0` adds a PlatformIO path for supported boards from Blue Pill/F1 through mainstream F/G/L/H7/U5/WB families, with fail-safe preflight before hardware mutation.",
    "Can scaffold, author, build, upload, and monitor ESP32-class firmware through ESP32er and PlatformIO Core, with zero-config bootstrap and a serial-aware preflight before hardware mutation.",
    "Can author YAML-based smart-home firmware through ESPHomer and ESPHome, including zero-config bootstrap, device-config generation, validation, compile, USB/OTA upload, and bounded log observation for ESP32 / ESP8266 / RP2040 / BK72xx devices.",
    "Can play media on the operator's machine: an audio file to the speakers through AudioPlayer (soundfile + sounddevice — volume in percent and a time-played budget that truncates a longer file or loops a shorter one), or a video file with audio on a chosen display through VideoPlayer (ffpyplayer, whose wheel bundles ffmpeg + SDL, plus an OpenCV window — display, volume, the same truncate/loop time budget, window size, and fullscreen); both are observational/output and ship on the canvas and as wrapped chat tools.",
    "Can SPEAK and LISTEN: Talker (text-to-speech) renders input_text to a 24 kHz WAV through an Ollama neural TTS model (default Orpheus-3b-FT) and is female-voice-only by design, while Whisperer (speech-to-text) records the microphone itself or transcribes a file via faster-whisper locally (NVIDIA-GPU auto-detect with an always-present CPU fallback) or a cloud Whisper API; both light a zero-latency console REC indicator driven by the live audio stream and are observational/output, on the canvas and as wrapped chat tools.",
    "Can manage real desktop windows by title: focus them, tile them, resize them, list them, and close them deterministically through Win32 calls.",
    "Can drive a real Playwright browser through scripted interactive steps for logins, forms, assertions, downloads, extraction, and end-to-end UI checks.",
    "Can drive a live Unreal Engine 5 editor through the Unreal MCP plugin, from either Multi-Turn chat or the visual workflow canvas.",
    "Actively reaps orphaned Windows console-host and pool-child processes so long Multi-Turn or ACPX sessions do not leave misleading Tlamatini-icon ghosts in Task Manager.",
    "Can autonomously kill only genuinely hung shell-wrapper subtrees through a boot-time command watchdog that measures CPU and I/O progress instead of guessing from elapsed time.",
    "Runs checked Multi-Turn requests through request-scoped planning, capability selection, tool calls, observations, monitoring, and final synthesis.",
    "Can optionally inspect and modify her own bundled source tree in self-modify builds after verifying that `TlamatiniSourceCode/` is actually present.",
    "Launches wrapped copies of selected workflow agents in isolated runtime folders without mutating templates.",
    "Lets users design, validate, save, pause, resume, and stop visual workflows through the Agentic Control Panel.",
    "Turns successful Multi-Turn tool executions into starter `.flw` workflows that can be inspected and validated in ACP.",
    "Packages the project into a distributable Windows release with installer and uninstaller tooling.",
]

HOW_IT_WORKS = [
    "Browser UI sends chat and workflow requests through Django views and Channels WebSockets.",
    "RAG chains load selected file/directory context, retrieve relevant chunks, and build answer prompts.",
    "DB-menu actions validate directories or SQLite files in the browser, then call Django views that either copy the live database out or stage a replacement into `DB/ToLoad/db.sqlite3`.",
    "Config -> Access Keys Wizard reads masked provider-key status from the backend and persists only the edited secrets, keeping the browser flow honest without dumping live values back to the page.",
    "The file-navigation and file-edit trio sits above raw shell execution: Globber enumerates matching files, Grepper scans content with regex while pruning noisy/binary trees, and Editor performs byte-exact in-place replacements without rewriting an entire file.",
    "External MCP connectivity is catalog-driven: `external_mcps.json` stores Claude-style `mcpServers` entries, `external_mcp_manager.py` negotiates stdio / streamable-HTTP / SSE / WebSocket transports, and wrapped remote tools are surfaced into the planner as `ext__<server>__<tool>` only after the connection is healthy.",
    "The MCP Doctor path is intentionally safer than a live connect: it reads the configured server entry, validates transport shape, runtime commands, PATH/toolchain presence, placeholder secrets, and docs/source URLs, then returns an onboarding diagnosis without consuming the server's real tool surface.",
    "Blenderer opens the official Blender MCP add-on TCP socket (default `localhost:9876`), sends one action payload or raw code-execution request, and returns the structured result through the same wrapped-tool / canvas contract used by the rest of the agent catalog.",
    "When Ask Execs is enabled, the synchronous Multi-Turn executor stops before each state-changing tool call, emits an `exec_permission_request`, and waits on `ExecPermissionBroker` until the browser sends Proceed or Deny.",
    "When the browser surfaces an Ask Execs prompt or a Notifier event, JavaScript can POST to `/agent/flash_window/`; the backend then best-effort flashes the `Tlamatini.exe` console/taskbar window through `window_flash.py` and prints an uppercase attention banner for the log.",
    "Installer-time registration writes a per-user HKCU Add/Remove Programs entry pointing at `Uninstaller.exe`, and frozen startup re-checks that entry through `windows_app_registration.self_heal_for_frozen()` so older installs retroactively appear in Windows' uninstall surfaces.",
    "In-app self-update uses Django views plus `self_update.py` to check GitHub releases, download and stage the selected package, then hand off the locked-file swap to the external `apply_update.ps1` helper before relaunch.",
    "Frozen-build hardening now verifies media dependencies during packaging too: `build.py` embeds numpy and OpenCV in both shipped Python runtimes and refuses to produce a release if those imports are missing.",
    "Before a heavy directory embedding run on supported NVIDIA hosts, a fail-open pre-flight guard can estimate VRAM pressure and surface a non-blocking warning in chat.",
    "Version resolution now flows through git tags, a runtime resolver module, generated build artefacts, and an open `/agent/version/` endpoint.",
    "A first-person self-knowledge file (`Tlamatini.md`) is injected into prompt construction for all chains, but loaded user context still outranks that self-reference when the request is a generic summary of the provided project.",
    "When Multi-Turn is enabled, the global planner selects context and tool stages before the executor binds only the relevant tools, including wrapped deterministic agents such as De-Compresser.",
    "Optional self-modify builds bundle `TlamatiniSourceCode/`; `copy_source_assets.py` generates that snapshot with rebuild instructions and redacted secrets, and prompt rules require her to verify that directory exists before claiming she can inspect or change her own code.",
    "The Kalier path talks directly to the MCP-Kali-Server Flask API over HTTP with Python-stdlib `urllib`, auto-seeding the default box from `kali_server_url` in Config -> URLs before any one-off per-call override is applied, and captures one atomic `INI_SECTION_KALIER` block per run.",
    "The Nmapper path resolves a user-installed `nmap` from explicit config, PATH, Program Files, or `%LOCALAPPDATA%`, constructs one safe scan action, captures XML plus normal output, parses hosts/open ports with the standard library, and emits one atomic `INI_SECTION_NMAPPER` block for Parametrizer/Forker routing.",
    "The STM32er path spawns the STM32 Template Project MCP stdio server, performs the MCP initialize handshake, runs exactly one requested tool or composite action, and emits one atomic `INI_SECTION_STM32ER` block with the result, project directory, and stage metadata.",
    "Before any flash-capable STM32er action, a critical-mission preflight validates the arm-none-eabi toolchain, STM32CubeIDE, programmer path, ST-LINK presence, and STM32F-family match; compile-only steps can run boardless, but unsafe hardware mutations are refused fail-safe.",
    "The ESP32er path resolves or bootstraps PlatformIO Core, invokes `pio` subcommands directly with Python-stdlib process control, validates project and serial-port readiness, and emits one atomic `INI_SECTION_ESP32ER` block with stage, project, port, and stdout/stderr payloads.",
    "The ESPHomer path resolves or bootstraps the `esphome` CLI, can generate a minimal valid YAML device config headlessly, validates either a serial board or OTA host before upload/log actions, and emits one atomic `INI_SECTION_ESPHOMER` block with action, config path, stage, and captured CLI output.",
    "The Windower path uses Win32 APIs plus the cross-process `AttachThreadInput` focus-transfer dance to locate windows by title and apply one lifecycle action while still returning structured geometry/state fields.",
    "The Playwrighter path loads a declarative step list, drives Playwright against Chromium/Firefox/WebKit, and emits one atomic `INI_SECTION_PLAYWRIGHTER` block with status, assertions, extracted values, and the final URL.",
    "The Unrealer path opens a TCP socket to the Unreal MCP plugin, sends one `{\"type\": command, \"params\": {...}}` payload, captures the JSON reply, and emits one `INI_SECTION_UNREALER` block for downstream logic.",
    "A boot-time command watchdog samples CPU time and I/O bytes across shell-interpreter subtrees and reaps only the ones that stay idle past the grace-and-streak window, protecting healthy long-running commands while rescuing wedged prompt waits.",
    "After spawn-capable tool calls and again after the final answer, the orphan reaper can sweep dead descendants, orphaned `conhost.exe` companions, and stale pool-linked processes without ever raising into the chat path.",
    "Tool calls execute in the backend, append observations, and may create wrapped runtime copies under `agent/agents/pools/_chat_runs_/`.",
    "On the next full start-up, `manage.py` can swap a staged database into place before Django imports, while archiving the previous live database under `DB/Older/<timestamp>/`.",
    "ACP flows deploy session-scoped pool instances, wire config values, validate NxN graph rules, and execute through Starter-driven flow semantics.",
    "Build scripts collect static assets, bundle Django/Python resources, add agent templates, and assemble `pkg.zip`, `Uninstaller.exe`, and `dist/Tlamatini_Release/`; `build.py --self-modify` additionally injects the generated source snapshot for self-rebuildable releases.",
]

HOW_TO_USE = [
    "Run from source: create a virtual environment, install requirements, migrate, create a superuser, collect static files, and start Django.",
    "Open `/agent/` for chat. Load a file or directory context before asking codebase-specific questions.",
    "Keep Multi-Turn unchecked for direct Q&A; enable Multi-Turn for tasks that need tools, wrapped agents, monitoring, or workflow seeding.",
    "Use `chat_agent_globber` to find files by pattern, `chat_agent_grepper` to locate matching content, and `chat_agent_editor` when you need an exact in-place change instead of rewriting a whole file or shelling out to grep/findstr/sed.",
    "To use the External MCP capability, open `External -> MCPs`, register or import a server into `external_mcps.json`, choose the transport/runtime fields, and let the dialog connect it before expecting its `ext__<server>__<tool>` tools to appear in Multi-Turn.",
    "When you are onboarding or debugging an external MCP, call `chat_agent_mcp_doctor` first or use the MCP Doctor workflow node; it can tell you whether the issue is transport selection, a missing runtime on PATH, placeholder secrets, or a bad endpoint before you spend time on a live connect attempt.",
    "If you want a guided external-MCP onboarding flow, use the Step-by-Step mode in the External MCP dialog so each required field is introduced progressively instead of dumping the whole connection contract at once.",
    "Use Config -> Access Keys Wizard when you need to wire or update provider credentials without editing `config.json` manually.",
    "Use About -> Check for updates on packaged installs when you want Tlamatini to fetch and stage the latest release without manually replacing the install folder.",
    "Tick `Ask Execs` when you want human approval before each state-changing Multi-Turn step; it is disabled until Multi-Turn is on, and a single Deny stops the whole chain with an explicit red interruption banner.",
    "When you are using a packaged install on Windows 10 or Windows 11, uninstall it through Settings -> Apps -> Installed apps or the legacy Programs and Features entry, not by manually deleting the folder.",
    "If an Ask Execs approval dialog or a Notifier event needs you while the browser is buried, watch for Tlamatini’s taskbar-attention flash and the matching uppercase banner in `tlamatini.log`.",
    "If you want her to inspect or modify herself, verify that `TlamatiniSourceCode/` exists in the current build first; self-modify is optional and absent builds must be treated honestly as read-only about their own code tree.",
    "For authorized Kali Linux assessments, run MCP-Kali-Server on the Kali box, set `Config -> URLs -> Kali server (Kalier)` once, and then call `chat_agent_kalier` from Multi-Turn with the desired `action` and `target` without repeating the box URL each turn.",
    "For local nmap reconnaissance, install nmap yourself or call `chat_agent_nmapper` with `action='install'` to launch the official free installer; then use `quick`, `full`, `top_ports`, `version`, `scripts`, `host_discovery`, `udp`, `custom`, or `validate` only against hosts you own or are explicitly authorized to test.",
    "For STM32 firmware work, install STM32CubeIDE, leave `Config -> URLs -> STM32 MCP server script` blank for zero-config bootstrap, and then call `chat_agent_stm32er` from Multi-Turn with one `action` at a time such as `validate`, `create_project`, `write_source`, `build`, `build_and_flash`, `serial_session`, or `live_monitor`.",
    "For ESP32 firmware work, leave `Config -> URLs -> pio_executable` blank for zero-config PlatformIO bootstrap, then call `chat_agent_esp32er` from Multi-Turn with actions like `bootstrap`, `validate`, `create_project`, `write_source`, `build`, `upload`, `build_and_upload`, `monitor`, or `monitor_session`.",
    "For ESPHome smart-home firmware work, leave `Config -> URLs -> esphome_executable` blank for zero-config bootstrap, then call `chat_agent_esphomer` from Multi-Turn with actions like `bootstrap`, `validate`, `new_config`, `config`, `compile`, `upload`, `logs`, or the one-shot `scaffold_compile_upload` flow.",
    "For desktop-window control, call `chat_agent_windower` from Multi-Turn to focus, tile, resize, list, or close a window by title, or model the same action in ACP with the Windower node.",
    "For interactive web automation, call `chat_agent_playwrighter` from Multi-Turn with a `steps_json` script, or author the same step list visually with the Playwrighter node on the canvas.",
    "For Blender work, enable the official Blender MCP add-on in Blender, make sure its TCP listener is reachable (default `localhost:9876`), then call `chat_agent_blenderer` from Multi-Turn or use the Blenderer node on the canvas.",
    "For Unreal Engine work, enable the Unreal MCP plugin inside a live UE5 project first, then call `chat_agent_unrealer` from Multi-Turn or use the visual Unrealer node on the canvas.",
    "Archive jobs can now be described directly in Multi-Turn or modeled visually in ACP: De-Compresser infers compress vs decompress from the `input` or `output` extension.",
    "If a second post-answer warning bubble ever lists surviving `name + PID` entries, treat it as an honest cleanup report and end the listed processes manually from Task Manager if needed.",
    "Use the `ACPX-Skills` navbar menu when you need to browse, enable/disable, diagnose, or reload the shipped SKILL.md catalog without asking the LLM to be your admin surface.",
    "Use the DB dropdown when you need a safe database snapshot or want to stage a different `db.sqlite3` for the next start-up without hot-swapping the live SQLite file.",
    "Open `/agentic_control_panel/` to drag agents, connect them, configure each node, validate, start, pause/resume, stop, and save `.flw` workflows.",
    "Use `python build.py`, `python build_uninstaller.py`, and `python build_installer.py` only when producing a packaged Windows release.",
]

AGENT_DESCRIPTION_GUIDE = [
    "The authoritative human-readable source for workflow-agent descriptions is `agents_descriptions.md` at the repo root, not an embedded JavaScript map or a hard-coded Django list.",
    "Current validation compares live template directories, `agents_descriptions.md` rows, and the README / Book bestiary counts so the generated dossier stays aligned with the running product.",
    "The ACP sidebar hover tooltip and the right-click Description dialog both resolve through that markdown file first; `README.md` is only a legacy fallback when the dedicated description file is absent.",
]

AGENT_RUNTIME_GUIDE = [
    "Every workflow agent follows the same operational skeleton: template directory, `config.yaml`, a session-scoped pool copy, PID/status/log files, and explicit source/target wiring.",
    "Chat-wrapped tool calls launch isolated runtime copies under `agent/agents/pools/_chat_runs_/`, while ACP uses named pool folders such as `starter_1` or `unrealer_1`.",
        "Specialized agents now stretch the platform in different directions: Globber/Grepper/Editor cover deterministic file discovery, regex search, and surgical in-place edits; Video-Analyzer closes hardware-in-the-loop video verdicts; MCP Doctor performs safe external-MCP onboarding diagnosis; ACPXer drives external coding-agent CLIs; Kalier drives a remote or tunneled Kali Linux tool server; Nmapper drives local use-only nmap scans for authorized targets; STM32er drives a zero-config STM32 firmware MCP bridge; ESP32er drives PlatformIO directly; ESPHomer drives ESPHome directly for YAML-authored smart-home devices; Blenderer drives a live Blender editor over the official MCP add-on socket; Unrealer drives a live UE5 editor; TeleTlamatini bridges full Tlamatini conversations into Telegram; and Telegrammer / Whatsapper send and receive messages over the official Telegram Bot API and Meta WhatsApp Cloud API.",
]

ACPX_SKILLS_GUIDE = [
    "The new `ACPX-Skills` navbar menu gives operators four direct actions over the skill catalog: Browse Skills, Configure Skills, Diagnostics, and Reload Registry.",
    "`Configure Skills` flips `Skill.enabled` exactly the way MCPs and Tools are toggled, so disabled skills disappear from `list_skills` and reject `invoke_skill` with `SKILL_DISABLED` instead of silently half-working.",
    "The diagnostics view cross-checks skill dependencies against disabled tools, disabled MCPs, missing ACPX agents, and orphan database rows whose SKILL.md disappeared from disk.",
]

def operator_surface_counts_guide(context: dict) -> list[str]:
    return [
        f"The live operator surface now stands at {context['workflow_agent_count']} workflow agents, {context['total_multi_turn_tools']} Multi-Turn tools, {context['acpx_tool_count']} ACPX tools, and {context['skills_count']} skills.",
        f"Source inspection confirms the total: {context['wrapped_chat_agent_count']} distinct wrapped chat-agent tools bound from `chat_agent_registry.py`, which combines with {context['core_python_tool_count']} core Python tools and {context['acpx_tool_count']} ACPX/Skill tools for {context['total_multi_turn_tools']} Multi-Turn tools overall.",
        "The count growth over older public badges now comes from several stacked waves together: the deterministic file-navigation/file-edit trio (Globber, Grepper, Editor), the ESPHomer smart-home firmware lane, MCP Doctor for External MCP onboarding, Zavuerer unified messaging, Video-Analyzer for robotic video verdicts, and Nmapper for local authorized nmap reconnaissance.",
        "The workflow-agent and wrapped-tool totals are validated from the live tree even when some handbook badges or older prose lines lag behind the newest release wave, so the dossier stays tied to source truth instead of stale summaries.",
        "This matters operationally because the planner never binds everything at once: the documented default `max_selected_tools` cap stays at 20, so breadth of capability does not mean uncontrolled tool sprawl per turn.",
    ]

CURRENT_RELEASE_GUIDE = [
    "Git resolves the product to `v1.42.0`: tag, local `main`, and `origin/main` all point to commit `c58b01ad`; the live source tree remains authoritative for generated counts and capability claims.",
    "`v1.42.0` releases STM32er's PlatformIO backend and device-aware routing for supported boards from Blue Pill/F1 through mainstream F/G/L/H7/U5/WB families.",
    "`v1.41.4` fixes External-MCP structured-output consumption: stdio and network clients now provide both text `content` blocks and `structuredContent` to the LLM instead of handing it only a short pointer.",
    "The shared formatter unwraps a sole `result` envelope, preserves errors and plain text, safely serializes structured data, caps oversized payloads, and is pinned by seven focused tests.",
    "`v1.41.3` keeps the Catalog of Prompts grouped into 13 categories with stable surviving ids, physical duplicate removal, gap-tolerant loading, and ranked fuzzy search.",
    "`v1.41.2` keeps Cancel latched to a per-user run epoch across RAG rebuilds, executor, retry, self-healing, Ask Execs, late answers, and frontend state without poisoning the next request.",
    "`v1.41.0` keeps screenshot paste/drop path-native: validated images are re-encoded in guarded Temp storage, inserted at the remembered caret, and handed to Image-Interpreter through removable chips.",
    "That Phase 1 path covers supported PlatformIO `ststm32` boards from Blue Pill/F1 through mainstream F/G/L/H7/U5/WB families, shares zero-config PlatformIO bootstrap with ESP32er, and adds ST-LINK-aware upload safeguards.",
    "The existing STM32F407 Template-MCP route remains the automatic default for blank/STM32F4 requests, while C0/H5/WBA/N6 stay explicitly unsupported until the planned ST-native CubeCLT backend exists; released migrations add camera-verified stepwise demos and a one-time contiguous catalog renumber.",
    "README.md and BookOfTlamatini.md now carry a clearer plain-Python agent disclaimer: transparency enables user control but is not a security warranty, and execution remains under the user's jurisdiction and responsibility.",
    "The README static badge and some Book release prose lag Git, so version identity comes from the live tag/commit graph rather than those historical text surfaces.",
    "The dirty worktree is counted without reproducing private values, credentials, endpoints, or machine-specific paths; this dossier does not stage, commit, or push it.",
    "The `v1.40.x` configurable-port and FlowPills work, Unreal 5.8 scaffold, Nmapper, self-healing, Create Flow, robotic loop, firmware/media agents, External MCPs, ACPX skills, and deterministic file tools remain carried product behavior.",
    "No new workflow agent type landed in this delta; live agent/tool totals are derived from the registry and repository rather than inferred from release numbers.",
    "README.md and BookOfTlamatini.md still supply the complete MIT-licensed project, easy-start installation, Ollama setup, architecture, usage, agent, and responsibility narrative rather than a narrow changelog.",
    "The regenerated PDF/PPTX treats tagged behavior and local configuration-only changes separately while preserving the full system, history, file tree, and effective-line inventory.",
]

STRUCTURED_CONTENT_1414_GUIDE = [
    "Modern MCP servers may return a short human-readable pointer in `content` while placing the actual result in `structuredContent`; the old parser discarded that machine-readable payload.",
    "Without the actual data, the model could repeat the same valid tool call until Tlamatini's repetition breaker force-stopped the run, making a successful external server appear to auto-cancel.",
    "`_format_mcp_tool_result` is now the single formatter used by both `_StdioMcpClient.call_tool` and `_NetworkMcpClientBase.call_tool`, so transport choice no longer changes result fidelity.",
    "Text blocks and non-text content are retained, while a sole `{\"result\": ...}` structured envelope is unwrapped before JSON serialization.",
    "`isError` still returns an explicit error and can now include structured-only error details; non-dict results are stringified safely.",
    "Structured payloads are capped at 24,000 characters by default with an explicit truncation marker, protecting the model context from unexpectedly huge tool responses.",
    "Seven tests cover pointer-plus-data, unchanged plain text, envelope unwrapping, text and structured-only errors, payload capping, and non-dict input.",
    "The fix is released as `v1.41.4` at `cec16594` and is present on `origin/main`.",
]

STM32ER_PLATFORMIO_WORKTREE_GUIDE = [
    "Release status: STM32er PlatformIO Phase 1 is published in tagged `v1.42.0` at `c58b01ad`, with local `main` and `origin/main` aligned on the release commit.",
    "`stm32_backend=auto` keeps the legacy Template-MCP path for blank/STM32F4 requests, but a board, a non-F4 device, or a PlatformIO-only action routes to the new direct PlatformIO backend.",
    "Friendly board aliases and device-to-board mappings include the STM32F103 Blue Pill, while raw PlatformIO board ids remain accepted for broader `ststm32` coverage.",
    "The backend mirrors ESP32er's zero-config bootstrap and shared per-user PlatformIO core, then supports environment, board, project, source, package, build, upload, monitor, QA, and artifact actions.",
    "A preflight resolves the board/family, validates `platformio.ini`, probes PlatformIO and ST-LINK readiness, and requires hardware only for upload/monitor operations; compile-only actions remain boardless-safe.",
    "`scaffold_build_flash` creates or reuses a project, writes source, builds, and flashes only when ST-LINK is confidently detected; otherwise a successful build is reported with a clear connect-and-flash next step.",
    "The broad family parser recognizes the full ST line, but PlatformIO-incompatible newest families C0/H5/U0/WBA/N6 are deliberately refused instead of risking an incorrect linker or target.",
    "The planned CubeCLT/CubeMX backend remains proposal-only for those newest devices, especially STM32N6 signing and external-flash requirements.",
    "New assets include the all-families proposal and migrations `0177`-`0179`: one-call Blue Pill build/conditional flash, two stepwise camera-verified board walkthroughs, then a deliberate category-grouped no-gap catalog renumber.",
    "Registry, contracts, wrapped-tool defaults, config, agent code, descriptions, and tests move together; focused coverage pins family/routing safety plus catalog contiguity, prompt-name identity, category ordering, and demo presence.",
]

STM32ER_STEPWISE_DEMOS_GUIDE = [
    "Migration `0178_add_stm32_stepwise_blink_camera_prompts.py` adds two `firmware_iot` walkthroughs: STM32F103 Blue Pill over an external ST-LINK V2 and STM32F407G-DISC1 over its embedded ST-LINK/V2.",
    "Both prompts require Multi-Turn, Exec Report, and Step-by-Step mode, execute exactly one stage per turn, wait for the operator's `READY`, and verify each prerequisite before moving forward.",
    "The six-stage Blue Pill path covers driver/CLI readiness, four-wire SWD wiring and probe detection, PlatformIO bootstrap, project/source/build, ST-LINK flash, and camera proof of the PC13 LED.",
    "The five-stage Discovery path uses the board's ST-LINK USB port, then bootstrap, project/source/build, embedded-probe flash, and camera proof of the green PD12 LED.",
    "Final evidence comes from `chat_agent_camcorder`; the prompt may pass the clip to Video-Analyzer or Image-Interpreter and must return a clear PASS or FAIL with the saved file path.",
    "These are seeded demonstrations and operator procedures, not a claim that hardware verification occurred during dossier generation.",
]

PROMPT_CATALOG_WORKTREE_GUIDE = [
    "Released migration `0179_regroup_resort_prompts_no_gaps.py` is a one-time deliberate override of the earlier no-renumber convention; it runs after the two new STM32 walkthroughs are appended.",
    "Every Prompt is sorted by the same category-display rank used by the UI and then by its prior id, so each category becomes one contiguous block from beginner workflows through specialized surfaces and the `other` fallback.",
    "Because `idPrompt` is the primary key, the migration first parks rows above 1,000,000 and then assigns final ids 1..N, avoiding collisions while rewriting `promptName` to `prompt-<id>`.",
    "The migration is intentionally one-way: reverse is a no-op because original ids are not stored; the source documents that no runtime foreign key references fixed prompt numbers.",
    "Future additions still append at max(idPrompt)+1, preserving contiguity until another deletion; the frontend's gap-tolerant fallback remains a defensive compatibility path.",
    "`test_prompt_catalog_contiguous.py` adds four database tests for no gaps, promptName/id agreement, nondecreasing category rank, and presence of both STM32 board demos plus Camcorder evidence.",
]

PROMPT_CATALOG_1413_GUIDE = [
    "Migration `0175_prompt_category_and_dedup.py` classifies the 106 historical prompt rows into 13 named operator categories, from Getting Started and Files/Search through ACPX, firmware, security, messaging, media, and a fail-safe More bucket.",
    "Migration `0176_delete_duplicate_acpx_prompts.py` physically deletes ids 40-52: 13 banner/Gemini variants that duplicate the seven portable ACPX demos retained at ids 33-39; surviving ids are never renumbered.",
    "Gaps are now valid. `/agent/list_prompts/` returns all visible rows regardless of id continuity, and the offline `prompt-N` fallback skips a missing id instead of terminating the catalog at the first gap.",
    "`Prompt.category` and `Prompt.hidden` remain schema-level controls, while `PROMPT_CATEGORY_ORDER` gives every known category a stable display rank and routes unknown or future values into `other` rather than dropping them.",
    "The modal renders explicit category headers and counts, then restores that grouped basic-to-advanced order whenever the search query clears.",
    "Live search accepts a prompt number, title words, mode labels, acronym, contiguous text, or fuzzy subsequence; it requires every query token, ranks best-first, highlights matched characters, and exposes Enter-to-open plus Escape/clear behavior.",
    "Prompt cards retain their One-Shot, Multi-Turn, ACPX, Exec Report, and Step-by-Step mode badges and continue to set the matching toolbar toggles when selected.",
    "The catalog stays viewport-pinned and bounded in CSS so its header and search bar remain reachable independently of the chat input height, with no JavaScript geometry coupling reintroduced.",
]

HARD_CANCEL_GUIDE = [
    "`agent/cancellation.py` mints a monotonically increasing run epoch per user and permanently latches the highest cancelled epoch; clearing the legacy setup boolean can no longer resurrect that run.",
    "Cancellation is isolated per user, so a browser tab cannot kill another user's TeleTlamatini or concurrent browser run; a missing epoch is deliberately fail-open so a dropped payload field cannot poison all future requests.",
    "`ask_rag`, both unified-chain payload rebuilds, and `CapabilityAwareToolAgentExecutor` carry and check `run_epoch`, closing the prior gap where model retries or tool loops could continue after Cancel.",
    "The executor returns a structured cancelled result instead of raising, preserving already-completed tool evidence, Exec Report data, and Create Flow inputs without triggering a fabricated transient-error fallback.",
    "Self-healing consults the same latch before retries, after watchdog waits, and before tactic announcements, so no new recovery tactic can be emitted once the operator cancels.",
    "Ask Execs polls the run latch and resolves a blocked Proceed/Deny request as deny; consumer teardown revokes the exact status emitter so a dying run cannot re-arm the UI.",
    "Frontend `userCancelledRun` is mutable per-session state: late tactic frames become strict no-ops until the next submit/reconnect, preventing the Send button from flipping back to Cancel.",
    "Coverage includes 24 cancellation contract tests, focused Ask-Execs/self-healing tests, and visible browser regression harnesses for long tool chains, model steps, repeated cancels, next-request recovery, and approval-modal cancellation.",
]

CHAT_IMAGE_1410_GUIDE = [
    "The chat accepts clipboard bitmaps through document-level Ctrl+V handling and image files dropped only onto the main chat column, avoiding conflict with the External-MCP dialog's document-level JSON drop surface.",
    "`POST /agent/paste_image/` validates the request, enforces a 25 MB ceiling, uses Pillow to flatten transparency onto white, re-encodes to JPEG, and writes a collision-safe `image_<timestamp>.jpg` under Tlamatini's guarded Temp directory.",
    "The browser inserts the saved absolute path at the remembered caret, including after focus moves to the page body during Alt+Tab, then renders one removable thumbnail chip per image.",
    "Removing a chip removes both its thumbnail and its exact path from the message, so an accidental paste can be reversed before submission.",
    "The path is the integration contract: Image-Interpreter reads local files, and `prompt.pmt` teaches Tlamatini to treat a fresh Temp image path as the user's supplied screenshot instead of asking for another attachment.",
    "`computeFormMinHeight()` measures the chips row and a ResizeObserver watches it, preventing the new row from pushing the textarea or Send button beyond the viewport.",
    "Implementation assets include `chat_image_paste.js`, the paste view/URL, template chip/drop-overlay nodes, `.chat-img-*` CSS, layout integration, self-knowledge/prompt guidance, and Temp-policy regression coverage.",
]

DJANGO_PORT_GUIDE = [
    "`django_port` moves the web UI and chat WebSocket bind away from a hardcoded 8000: edit the integer in `config.json`, restart Tlamatini, and no rebuild or source edit is required.",
    "The motivating failure is Windows `WinError 10013`: Hyper-V, WSL, or Docker can reserve port 8000 inside a dynamic exclusion range, making a frozen build unable to bind until the port is changed.",
    "Three stdlib-only helpers run before Django imports: `_resolve_config_path()` chooses `CONFIG_PATH`, the frozen executable's neighbor, or source `agent/config.json`; `_resolve_django_port()` validates `1..65535`; `_apply_configured_port()` injects the result.",
    "The completion pass calls `_apply_configured_port(sys.argv)` once from `main()`, outside the frozen branch, so frozen double-click, `.flw` association, browser auto-open, source `runserver`, and `startserver` all agree.",
    "Resolution is fail-open to 8000 for missing, unreadable, malformed, non-numeric, or out-of-range values, while an explicit CLI port such as `runserver 9100` always wins and is never double-appended.",
    "The injected source-mode value is a bare port so Django keeps its loopback host; direct Daphne/Uvicorn bypasses `manage.py`, MCP listeners `8765`/`50051` use their own settings, and TeleTlamatini keeps its own base URL.",
    "The 24-test `agent/test_django_port_config.py` suite AST-lifts the pre-Django helpers and pins path resolution, validation, fail-open behavior, CLI precedence, frozen/source wiring, and `startserver` forwarding without importing side-effectful `manage.py`.",
    "`freeingport8000.ps1` is a separate elevated Windows repair helper that resets dynamic TCP/UDP ranges, restarts WinNAT, reports excluded ranges, and performs a loopback bind test; changing `django_port` remains the safer normal remedy, and generated docs never reproduce the script's machine-specific log path.",
]

NMAPPER_GUIDE = [
    "Nmapper is the `v1.39.3` local nmap bridge for authorized pentesting and CTF recon: it is a use-only integration, not a bundled scanner.",
    "The agent never ships or redistributes nmap; it resolves a user-installed `nmap` from explicit config, PATH, Program Files, or `%LOCALAPPDATA%`, and its `install` action launches the official free installer path.",
    "The default scan path is an unprivileged TCP connect scan (`-sT`), while raw-packet features such as SYN/OS/UDP are downgraded or refused safely when Npcap or elevated packet capture is unavailable.",
    "Supported actions include `quick`, `full`, `top_ports`, `version`, `scripts`, `host_discovery`, `udp`, `custom`, `validate`, and `install`, always for targets the operator owns or is explicitly authorized to assess.",
    "Each run emits one atomic `INI_SECTION_NMAPPER` block with action, target, scan technique, ports, return code, success state, hosts-up/open-port summaries, Npcap state, XML path/content, normal output, and stage.",
    "Implementation assets include `agent/agents/nmapper/`, migrations `0170`-`0172`, `test_nmapper_agent.py`, `chat_agent_nmapper`, `update_nmapper_connection`, contracts/Parametrizer fields, ACP connector assets, FlowCreator/FlowHypervisor entries, and handbook/docs updates.",
]

STARTUP_PROMPT_POLISH_GUIDE = [
    "`v1.39.4` restored first-run/startup dialog closeability so a fresh launch can no longer be trapped behind an unclosable overlay.",
    "Commit `a45fe0e0` followed the public `v1.39.4` tag with Catalog-of-Prompts localization cleanup; that historical polish remains carried by current `v1.42.0`.",
    "The prompt catalog path stays centralized through the secure one-call `/agent/list_prompts/` endpoint ordered by category rank and stable surviving id, while the gap-tolerant probe loop remains only as an offline fallback.",
    "Frontend mutable-state tests and dialog templates continue to guard the chat/startup/overlay surfaces so future cleanup passes do not reintroduce const-poison or close-button regressions.",
]

FLOWPILLS_DISCOVERY_GUIDE = [
    "Tlamatini-FlowPills reads `HKCU\\Software\\XAIHT\\Tlamatini` first, especially the exact `AgentsRoot`, before falling back to Installed Apps, `.flw` association data, executable-relative roots, or source/preserved probes.",
    "The registry contract always rewrites six REG_SZ values — `InstallLocation`, `AgentsRoot`, `SourceAgentsRoot`, `AgentManifestPath`, `Version`, and `AgentCatalogVersion` — using an empty value when something is unknown so stale metadata cannot survive across source and frozen runs.",
    "`agent_manifest.py` discovers only complete templates, excludes `pools` and `__pycache__`, hashes script/config contents, computes a name-set catalog id, and writes atomically only when meaningful data changed.",
    "Launch publication runs first in `AgentConfig.ready()` on a daemon thread with its own idempotency gate, so optional MCP, model, ACPX, or skill import failures cannot suppress companion discovery.",
    "An agents-preserving uninstall restamps the manifest as `preserved`, writes `.tlamatini-preserved-agents.json` with the manifest SHA-256, and keeps the discovery key pointing to the preserved catalog.",
    "The contract is HKCU-only, no-admin, fail-open, and read-only with respect to agent templates; the filesystem remains the final validity authority for companion applications.",
]

UNREAL_SCAFFOLD_GUIDE = [
    "The new Unreal scaffold prompt asks for only project name and destination, locates or obtains `XaihtUnrealEngineMCP`, and invokes its deterministic `scaffold_unreal_project.py` helper.",
    "The scaffold copies and renames `MCPGameProject`, sets EngineAssociation 5.8, finds an installed UE 5.8 even when it is not registered, bundles UnrealMCP, and generates the Visual Studio 2026 solution.",
    "UE 5.8 build compatibility is explicit: V7 build settings, `Unreal5_8` include order, a `Directory.Build.targets` mitigation for the Windows environment-length limit, and a pre-fixed Visual Studio Tools plugin.",
    "The prompt instructs operators to build the project target rather than the entire solution, then open the editor; UnrealMCP starts its TCP listener at `127.0.0.1:55557` for Unrealer.",
    "Unrealer normalizes disk-style `/Content` references to `/Game`, maps the friendly material `slot` input to `slot_index`, and forwards one command per TCP connection into the live editor.",
    "Migrations `0173` and `0174`, Unrealer code/config comments, README/Book entries, and the current release metadata carry the scaffold workflow across source, catalog, and operator documentation.",
]

RESPONSIVENESS_HARDENING_GUIDE = [
    "The `v1.39.5` smoothness wave focuses on bounded waits, concurrency isolation, accurate partial-result reporting, and lossless final answers rather than a new agent count.",
    "Image-Interpreter now bounds Ollama connect/read gaps; the System-Metrics WebSocket bounds receive time; External MCP warm-connect/supervisor gates recover if thread start fails; and command/ACPX output decoding is UTF-8-safe.",
    "Nmapper assigns per-host budgets to slow actions, keeps the outer kill budget above them, preserves partial output, uses collision-proof filenames for parallel scans, and labels hard timeouts honestly.",
    "Multi-Turn folds deferred completion deliverables into every exit path, while orphan-survivor state is keyed by conversation user so simultaneous requests cannot consume or clear each other's evidence.",
    "`.flw` export redaction covers ordinary secret paths and URI userinfo, and the command parser distinguishes Windows trailing backslashes from escaped inner quotes at assignment boundaries.",
    "These changes are now committed release behavior; the user's pre-refresh uncommitted configuration values remain private and are not quoted in this dossier.",
]

ROBOTIC_LOOP_GUIDE = [
    "`v1.38.0` is the milestone where the Robotic-Loop-Training story became concrete: Tlamatini demonstrated a closed hardware loop by programming a robotic arm from a blank page and two cameras.",
    "The loop is intentionally composable rather than one monolith: STM32er writes/builds/flashes firmware, Camcorder records the physical attempt, Video-Analyzer judges the captured motion, and Forker routes the next branch.",
    "Video-Analyzer keeps the loop conservative: a deterministic OpenCV gate rejects no-motion clips before model spending, then two independent Ollama cloud vision interpreters must agree before `PASS_OK` is emitted.",
    "Forker branches on substring-safe `TLM_VERDICT::<TOKEN>` markers, so `FAIL_NO_MOTION`, `FAIL_WRONG_MOTION`, `UNCLEAR`, or `ANALYSIS_ERROR` can never accidentally match the success route.",
    "The PDF and PPTX now describe that loop as a system capability, not merely a release-note flourish, because it explains how Tlamatini can iteratively improve real hardware with observed evidence.",
]

FRONTEND_HOTFIX_GUIDE = [
    "`v1.38.1` was the same-week frontend-state-recovery hotfix: `package.json` was aligned at the tagged commit `08efa1d2`, while the functional fix landed in `af356c31` after the `85ee4e6c` const-poison incident.",
    "The core contract is explicit: cross-file runtime globals in `agent_page_state.js` and `acp-globals.js` that other modules reassign must remain `let`, because per-file ESLint cannot see those cross-file writes.",
    "`agent/test_frontend_mutable_state.py` now guards both source files and collected staticfiles so an automated cleanup cannot silently turn chat state, ACP state, tools, agents, skills, history, or busy flags back into `const`.",
    "The Catalog of Prompts now loads through one secure `GET /agent/list_prompts/` endpoint ordered by `idPrompt`, eliminating expected-404 console spam and preventing an `idPrompt` gap from hiding later prompts.",
    "The legacy prompt probe loop remains as an offline fallback, and the same fix also hardened dialogs: Configure-Mcps probe loops exit cleanly, About-video `play()` promises are guarded, and Esc closes About/Update overlays.",
]

V136_RELEASE_GUIDE = [
    "Release identity: current public tag `v1.42.0` at `c58b01ad`, the STM32er PlatformIO expansion; v1.41.4 External-MCP structured output, v1.41.3 prompt-catalog organization, v1.41.2 Hard Cancel, v1.41.0 image ingestion, and earlier waves remain part of this release lineage.",
    "New agent: Video-Analyzer becomes the current media-verdict workflow agent and wrapped `chat_agent_video_analyzer`, complementing Image-Interpreter with video-specific motion analysis.",
    "Implementation assets: `agent/agents/video_analyzer/`, migrations `0166_add_video_analyzer.py`, `0167_add_chat_agent_video_analyzer_tool.py`, `0168_add_video_analyzer_demo_prompt.py`, `test_video_analyzer_agent.py`, `chat_agent_registry.py`, `mcp_agent.py`, and `services/agent_contracts.py` all move together.",
    "Model strategy: `interpreter_model_1` defaults to `qwen3-vl:235b-cloud`, `interpreter_model_2` defaults to `qwen3.5:cloud`, and `merging_model` defaults to `glm-5.2:cloud`, with independent calls merged only after both interpreters report.",
    "Routing contract: every run emits `INI_SECTION_VIDEO_ANALYZER` plus `TLM_VERDICT::<TOKEN>` markers such as `PASS_OK`, `FAIL_NO_MOTION`, `FAIL_WRONG_MOTION`, `UNCLEAR`, and `ANALYSIS_ERROR` for Forker and Parametrizer.",
    "Adjacent UI work: prompt search moved from exact-title hunting to substring, word-start, and fuzzy matching, and generated `.flw` files now use a serpentine layout to reduce visual congestion.",
]

VIDEO_ANALYZER_GUIDE = [
    "Video-Analyzer is Tlamatini's video-verdict agent: she receives `video_pathfilenames` as a direct file, wildcard, folder-newest rule, or Camcorder pool name and resolves the video before model calls begin.",
    "The first gate is deterministic and cheap: OpenCV/numpy frame sampling computes a motion score and returns `FAIL_NO_MOTION` without spending LLM calls when the recording clearly shows no movement.",
    "When motion exists, two Ollama cloud vision models run in parallel. One specializes in temporal/action evidence, the other provides an independent holistic judgement, and a merger model produces the final report.",
    "The final answer is intentionally conservative: PASS requires both interpreters to agree; disagreement or weak evidence becomes `UNCLEAR`, and wrong movement becomes `FAIL_WRONG_MOTION` rather than a false pass.",
    "Structured output is `INI_SECTION_VIDEO_ANALYZER` with `video_path`, `verdict`, `verdict_token`, `confidence`, `motion_score`, `frames_analyzed`, model names, `status`, and the body report.",
    "The robotics loop is now explicit: STM32er can flash firmware, Camcorder can record the board, Video-Analyzer can judge the recorded motion, and Forker can branch on the `TLM_VERDICT::` token to retry or finish.",
]

PROMPT_SEARCH_AND_FLOW_GUIDE = [
    "Prompt search now behaves like an operator tool instead of a static dropdown: category grouping, numeric/acronym/substring/word-start/subsequence matching, best-first ranking, highlighted hits, and mode badges make the deduplicated catalog navigable.",
    "The prompt-card rendering keeps One-Shot, Multi-Turn, ACPX, Exec Report, and Step-by-Step families distinguishable while migrations `0175`-`0176` remove redundant ACPX variants without renumbering surviving prompt ids.",
    "Generated `.flw` workflows now use a serpentine canvas layout with row capacity and alternating direction, preserving the logical Starter -> Agent -> ... -> Ender order without pushing large chains into an unreadable line.",
    "This matters for documentation because the PDF and deck must describe not only new agents, but also the operator usability improvements that make the larger system actually usable.",
]

SELF_HEALING_GUIDE = [
    "Every Multi-Turn model step now goes through `agent/self_healing.py::SelfHealingInvoker`, so model calls are bounded by `unified_agent_llm_step_timeout_seconds` and retried with distinct tactics instead of hanging silently.",
    "Recovery tactics include plain retry, short back-off, trimming oldest messages with `trim_messages`, and a tool-less summary fallback; the ladder can run up to `unified_agent_llm_step_max_tactics` unless the user presses Cancel.",
    "A status broadcaster registered by `consumers.py` sends live first-person recovery messages to the user's chat while the executor works in a worker thread.",
    "If recovery is exhausted after agents already ran, `mcp_agent.py` builds a degraded but truthful answer from real `ToolMessage` results, preserves Exec Report/Create Flow evidence, and prepends `recovery_preamble(...)` instead of claiming no tools ran.",
    "Coverage includes `agent/test_self_healing.py`, `agent/tests.py`, `Tlamatini/tests_e2e/test_self_healing_visual.py`, and `Tlamatini/tests_e2e/test_create_flow_visual.py` for visible browser validation.",
    "Frontend follow-up: `agent_page_ui.js::isSelfHealingStatusMessage()` anchors on leading `Tactic #` or `Tactic '` status frames so `appendChatMessage()` keeps controls disabled and the Send button on Cancel while the run is still executing.",
    "The matcher deliberately avoids substring matching because the final `recovery_preamble(...)` quotes tactic lines; a loose `includes` match would misclassify the final answer and trap the UI on Cancel forever.",
]

CREATE_FLOW_GUIDE = [
    "The whole-answer SUCCESS/FAILURE classifier `agent/services/answer_analizer.py` was removed on 2026-07-06, so no extra LLM round trip computes `answer_success`.",
    "The browser now shows Create Flow when Multi-Turn ran, at least one successful tool call resolves to a registered canvas agent, and the user is not anonymous.",
    "Successful tool calls are resolved through a punctuation/space/case-insensitive registry key, so display names such as `File Creator`, `File-Creator`, and `filecreator` converge to the live Agents sidebar entry.",
    "Generated `.flw` downloads keep only successful, resolvable tool-call entries, post a successful-only `tool_calls_log` to `/agent/flow_from_tool_calls/`, and drop failed or unregistered executions rather than turning them into broken workflow nodes.",
    "If the registry cannot load, the frontend fails open: the button stays available and the backend flow normalizer remains the final validation layer.",
    "Exec Report remains per-tool evidence, not a whole-answer verdict; the checkbox is disabled/greyed unless Multi-Turn is checked.",
]

FRONTEND_RECOVERY_GUIDE = [
    "`agent_page_chat.js::appendChatMessage()` now has a self-healing status branch before the final-answer branch, so status frames render without calling `enableControlsAfterOperation()`.",
    "`disableControlsDuringOperation()` is re-applied idempotently for each live tactic status line, preserving the user's Cancel button while the worker thread continues.",
    "`agent_page_ui.js::isSelfHealingStatusMessage()` strips leading icons/symbols and then anchors on `Tactic #` / `Tactic '`, matching standalone status frames but not the final recovery summary.",
    "Create Flow validation now uses `_resolveSuccessfulAgents()`, `_agentNameKey()`, and `_buildRegistryKeyMap()` to normalize display names against the live registry and skip only unresolved successful entries.",
    "`eslint.config.mjs` declares `isSelfHealingStatusMessage` as a global, while `docs/claude/frontend.md`, `docs/claude/multi-turn.md`, and `docs/claude/recent-fixes.md` document the gotcha and the no-substring rule.",
]

DISCOVERER_PDCP_GUIDE = [
    "Discoverer now has an optional ProjectDiscovery Cloud Platform key (`pdcp_api_key` / `PDCP_API_KEY`) for cvemap/vulnx rate limits and nuclei `-ai` or cloud upload features.",
    "Config -> Access Keys Wizard exposes `Security Recon (ProjectDiscovery)`, writes the key to `config.json`, `data.keys`, and `agent/agents/discoverer/config.yaml`, and blank fields preserve existing values.",
    "`tools._seed_global_agent_defaults` auto-injects the configured key into every `chat_agent_discoverer` run so prompts never paste the credential.",
    "`agent_contracts.py` redacts `pdcp_api_key` from `.flw` exports, and `regen_secrets.py` scrubs it back to `PDCP_API_KEY` before a push-able tree.",
    "The new `0169_add_discoverer_cvemap_latest_demo_prompt.py` migration seeds a passive latest-CVE briefing prompt that uses cvemap/vulnx and reports whether `pdcp_used` was active.",
]

DISCOVERER_VULNX_GO_GUARD_GUIDE = [
    "ProjectDiscovery retired cvemap's CVE API in August 2025, so Tlamatini keeps the operator-facing `cvemap` tool key but installs and runs ProjectDiscovery's successor binary, `vulnx`.",
    "`discoverer.py` now maps `cvemap` to `github.com/projectdiscovery/cvemap/cmd/vulnx@latest`, resolves the installed binary as `vulnx`, and builds `vulnx id <CVE>` or `vulnx search --severity ... --product ... --limit ...` arguments.",
    "The findings counter now understands vulnx object JSON such as `{\"count\": N, \"results\": [...]}`, so Exec Report and Forker routing see the real result count instead of a misleading one-line object.",
    "The private Go compiler and ProjectDiscovery tool binaries still install under `<install_dir>/Go` / `<install_dir>/Go/bin-tools`, keeping the runtime self-contained and avoiding any system Go or PATH mutation.",
    "The `.gitignore` Go-deny block, tracked `git_deny_go.py`, and its managed pre-commit hook prevent `Go/`, `bin-tools/`, `go-build/`, `pkg/mod/`, and downloaded Go archives from entering source control; ignored `Go/` content is not counted as project code.",
    "The committed asset wave also includes `0169_add_discoverer_cvemap_latest_demo_prompt.py`, `.claude/skills/tlamatini-daily-chat-test/harness/discoverer_1000.py`, and `Tlamatini/agent/test_discoverer_thousand.py` for latest-CVE and high-volume Discoverer validation.",
]

V1332_RELEASE_GUIDE = [
    "Release identity: latest reachable public tag `v1.33.2`; the generated artifacts now describe the Zavuerer release family plus post-release cleanup while preserving the v1.32.0 quality-and-identity background.",
    "New agent: Zavuerer becomes the 83rd workflow-agent type and the 60th wrapped chat-agent, adding `chat_agent_zavuerer` for Zavu unified messaging across SMS, WhatsApp, Telegram, Email, and Voice.",
    "Configuration: Config -> Access Keys Wizard now includes `Unified Messaging (Zavu)` and persists `zavu_api_key`, which the wrapped runtime seeds into Zavuerer without exposing the secret in prompts.",
    "Canvas/runtime support: `agent_contracts.py`, `views.py`, `capability_registry.py`, `chat_agent_registry.py`, `tools.py`, frontend ACP JS/CSS, and migrations `0159`-`0164` all move together to make Zavuerer usable from both surfaces.",
    "Follow-up defaults: `config.json` now favors `glm-5.2:cloud` across the primary chat/model slots, and the latest commits also adjust runtime parsing/middleware/settings surfaces around the release.",
    "Cost and safety wording: sign-up for Zavu is free, but sends are pay-as-you-go per message; the docs keep the authorized, opted-in recipient boundary explicit for A2P, WhatsApp-window, consent, and GDPR-style rules.",
]

IMAGE_INTERPRETER_GUIDE = [
    "Image-Interpreter is now a triple-model vision analyst, not a single generic image describer: each image goes through two parallel interpreter calls and one merger pass.",
    "`interpreter_model_1` defaults to `qwen3.5:cloud` and is tuned for forensic OCR, mockup/GUI element inventories, percent-based positions/sizes, colors, fonts, and verbatim text.",
    "`interpreter_model_2` defaults to `gemma4:cloud` and reads the image holistically: design intent, visual hierarchy, scene meaning, people, and reasoned identity hypotheses.",
    "`merging_model` defaults to `glm-5.2:cloud`; it waits behind a barrier until both interpretations arrive, then emits one definitive report with union-of-facts, conflict notes, and discrepancy handling.",
    "All four prompt surfaces (`prompt_user`, `prompt_interpreter_model_1`, `prompt_interpreter_model_2`, and `prompt_merging_model`) receive the image file name as an identity clue, because a file named after a person often depicts that person.",
    "Fail-safe behavior is explicit: one failed interpreter still lets the merger work from the survivor; a failed merger returns both raw interpretations concatenated instead of losing the analysis.",
    "The structured output is now `INI_SECTION_IMAGE_INTERPRETER` with `file_path`, `interpreter_model_1`, `interpreter_model_2`, `merging_model`, `status`, and the merged report body for Parametrizer/Forker routing.",
    "Config -> Models now exposes all three Image-Interpreter model slots (`image_interpreter_model`, `image_interpreter_model_2`, `image_merging_model`), and older preserved configs receive defaults so Save is not blocked by empty new fields.",
]

ZAVUERER_GUIDE = [
    "Zavuerer is Tlamatini's new Zavu unified-messaging agent: one workflow node and one wrapped `chat_agent_zavuerer` tool can send SMS, WhatsApp, Telegram, Email, or Voice messages through Zavu's `/v1/messages` API.",
    "The operator configures one `zavu_api_key` through Config -> Access Keys Wizard -> Unified Messaging (Zavu), and the runtime seeds that key into Zavuerer automatically so prompts never need to repeat or expose the secret.",
    "The agent speaks direct HTTP through the Python standard library, has a safe `health` action, refuses sends when the key/recipient/body/channel are invalid, and always emits `INI_SECTION_ZAVUERER` for Parametrizer/Forker branching.",
    "Canvas integration is complete: `agent_contracts.py` redacts `zavu_api_key`, `views.py` handles Zavuerer connections, frontend CSS/JS gives the node a visible identity, and migrations `0159`-`0164` seed the Agent, Tool, demo/catalog prompts, and setup-wizard dedupe.",
    "Cost and safety boundary: Zavu sign-up is free, but sending is pay-as-you-go per message; Zavuerer must only message authorized, opted-in recipients, with A2P, WhatsApp 24-hour-window, consent, and GDPR-style responsibilities called out explicitly.",
]

EXTERNAL_MCPS_GUIDE = [
    "The `v1.26.0` headline feature is the External MCP universal client: Tlamatini can now consume any MCP server declared in `external_mcps.json` instead of depending only on bundled MCP integrations.",
    "Four transports are supported in the current implementation — stdio, streamable HTTP, legacy SSE, and WebSocket — and the active server set is intentionally bounded so operators can expose high-value remote tools without flooding the planner.",
    "Once connected, remote MCP tools are wrapped into the same Multi-Turn execution surface under `ext__<server>__<tool>`, so they participate in planning, execution reporting, Ask Execs, and the rest of the operator guardrails instead of becoming a parallel hidden subsystem.",
]

MCP_DOCTOR_GUIDE = [
    "MCP Doctor is the new onboarding and triage specialist added with the External MCP release: she inspects a configured server entry and tells the operator what is wrong before a live connection attempt wastes time.",
    "The diagnostic checks cover transport shape, runtime command availability on PATH, placeholder or missing secrets, endpoint/source/docs links, and the concrete next step the operator should take to make the server connectable.",
    "Operators can reach the same behavior from both surfaces that matter: the MCP Doctor workflow agent on the canvas and the wrapped `chat_agent_mcp_doctor` tool inside Multi-Turn.",
]

EXTERNAL_MCP_ASSETS_GUIDE = [
    "The release is backed by real tracked assets rather than markdown-only claims: `agent/external_mcp_manager.py`, `agent/external_mcps.json`, the `agent/agents/mcp_doctor/` tree, the `0141`-`0143` migrations, `static/agent/js/external_mcps_dialog.js`, and `static/agent/css/external_mcps_dialog.css` are the core implementation wave.",
    "The current test surface proves this is not a shallow UI feature: `test_external_mcp_universal.py`, `test_external_mcp_transports.py`, `test_external_mcp_e2e.py`, `test_external_mcp_add_flow.py`, and `test_parametrizer_mcp_doctor.py` exercise the universal-client path and the onboarding diagnosis path.",
    "Handbook evidence now exists alongside the code too: the README tutorial, Book sections, capability/planner wiring, and the dialog assets all moved together in the same release wave, with later cleanup removing obsolete draft documentation from the tracked tree.",
]

BLENDERER_GUIDE = [
    "Blenderer was introduced in `v1.20.0` as the 77th workflow agent and remains the direct bridge to the official Blender MCP add-on, letting Tlamatini operate a live Blender session from either Multi-Turn chat or the visual workflow canvas.",
    "The bridge talks over a TCP socket (default `localhost:9876`) and supports both raw `execute_code` requests and higher-level scene/object/material/render actions, so operators can mix deterministic verbs with precise Python-driven 3D automation.",
    "This extends Tlamatini beyond code and firmware orchestration into DCC / 3D-production work: asset inspection, scene mutation, material changes, camera setup, and render-triggering now live inside the same operator surface as the rest of the system.",
]

SELF_UPDATE_GUIDE = [
    "The packaged application now includes an in-app self-update flow exposed from About -> Check for updates, so operators can refresh a frozen install without manually unpacking and replacing the entire release folder.",
    "The backend checks the latest GitHub release, downloads and stages the package, then hands off the locked-file swap to `apply_update.ps1`, which performs the replacement after the running executable exits and relaunches the updated app.",
    "The update path is state-preserving by design: it keeps `config.json`, user content, and one `agents_backup` generation, and it stages the user's database through `DB/ToLoad/` plus `post_update_migrate.flag` so the next launch restores and migrates that DB instead of wiping it.",
]

SOURCE_SNAPSHOT_GUIDE = [
    "Self-modify builds now rely on `copy_source_assets.py` to generate `TlamatiniSourceCode/` as a fresh rebuild-oriented snapshot of the repository.",
    "That snapshot includes source, build scripts, docs, skills, and small required binaries while deliberately omitting or redacting heavy media and secrets; `_REBUILD_INSTRUCTIONS.md` plus `_SOURCE_SNAPSHOT_MANIFEST.json` explain how to restore the missing payloads.",
    "Because the snapshot is optional, the runtime contract remains strict: Tlamatini must verify that `TlamatiniSourceCode/` exists before claiming she can inspect, modify, or rebuild herself.",
]

API_KEYS_WIZARD_GUIDE = [
    "Config now includes an Access Keys Wizard / API-Keys Wizard path so operators can manage provider credentials from the browser instead of editing `config.json` directly.",
    "The backend exposes masked status plus explicit save endpoints, which keeps the operator flow convenient without echoing raw secrets back into the page.",
    "That feature arrives as a real asset wave, not just a hidden helper: `agent/access_key_wizard.py`, `static/agent/js/access_keys_wizard.js`, `static/agent/css/access_keys_wizard.css`, and the updated `templates/agent/agent_page.html` now form a dedicated setup surface.",
    "This matters operationally because ACPX agents, cloud LLM providers, and adjacent integrations often fail for simple missing-key reasons; the wizard turns that setup into a first-class UI workflow.",
]

FILE_CREATOR_HARDENING_GUIDE = [
    "The `v1.19.5` File-Creator hardening pass fixes the wrong-symbol / wrong-escape corruption path that hit backslash-dense Java, JSON, CSS, and regex files during wrapped chat-agent execution.",
    "Two byte-exact channels now exist: plain `content` is re-extracted verbatim from the raw request, and `content_b64` can carry source or binary bytes through a parser-immune base64 path when exact escaping matters most.",
    "For operators, the consequence is straightforward: code generation and document generation can now be described as byte-complete file writes rather than best-effort convenience output.",
    "This belongs in the dossier because File-Creator is one of Tlamatini’s central deterministic execution surfaces and one of the safest alternatives to GUI typing or editor-driving automation.",
]

MEDIA_VOICE_GUIDE = [
    "The current media family spans capture, playback, and speech: Shoter (screen), Camcorder (camera-in), Recorder (mic-in), AudioPlayer (speakers-out), VideoPlayer (screen-out with audio), Talker (text-to-speech), and Whisperer (speech-to-text).",
    "Talker remains female-only by design and uses an Ollama neural TTS model with SNAC decoding, while Whisperer records the microphone itself or transcribes a file through faster-whisper locally or cloud Whisper APIs.",
    "Recent stability work matters here too: Talker now chunks long input by sentence for long-form speech, and media-output defaults were moved into the application `Temp` directory rather than user content folders.",
]

COMMAND_WATCHDOG_GUIDE = [
    "A boot-time autonomous command watchdog now protects the chat against shell wrappers that stay alive while making zero CPU or I/O progress, the classic signature of a mangled prompt waiting forever on stdin.",
    "It never kills on elapsed time alone: the subtree must outlive the grace window and remain idle for a configured streak, so long builds, downloads, and compiles are preserved while only genuinely wedged interpreters are reaped.",
    "This watchdog complements rather than replaces the orphan reaper: the watchdog rescues blocked synchronous tool calls before they return, while the orphan reaper still handles post-return survivors and stale console companions.",
]

NEW_ASSETS_GUIDE = [
    "The tagged `v1.41.4` wave adds the shared `_format_mcp_tool_result` path in `external_mcp_manager.py` plus seven focused cases in `test_external_mcp_universal.py` for machine-readable MCP results.",
    "The tagged `v1.42.0` STM32er/prompt wave adds over one thousand lines to `stm32er.py`, expands config/registry/contracts/tools/tests/docs, and tracks the all-families proposal, migrations `0177`-`0179`, and `test_prompt_catalog_contiguous.py`.",
    "The clearer disclaimer commit for README.md and BookOfTlamatini.md is included in the `v1.42.0` ancestry; neither handbook is modified by this generation pass.",
    "The `v1.41.0` image-ingestion wave adds `chat_image_paste.js`, paste view/URL wiring, chat chip/drop-overlay template nodes, image CSS, layout observation, Temp-policy coverage, and matching README/Book/self-knowledge guidance.",
    "The `v1.41.2` cancellation wave adds `agent/cancellation.py`, `test_cancellation.py`, `test_ask_execs_allowlist.py`, expanded self-healing/frontend tests, and coordinated changes across consumers, executor, RAG, retry, permission, settings, and chat-state surfaces.",
    "Two browser regression harnesses under `.claude/skills/tlamatini-daily-chat-test/harness/` exercise repeated cancels, next-request recovery, model/tool-chain cancellation, Ask-Execs cancellation, and the destructive/human-contacting/network allowlist contract.",
    "The `v1.41.3` prompt-catalog wave adds migrations `0175`-`0176`, Prompt category/hidden fields, ordered category metadata in `views.py`, and grouped, gap-tolerant, fuzzy-searchable catalog behavior in `tools_dialog.js`/`.css`.",
    "The same prompt wave physically removes 13 duplicate ACPX rows while retaining seven portable originals and stable surviving ids; the primary endpoint and offline fallback both tolerate the resulting gap.",
    "`freeingport8000.ps1` is the new Windows repair asset for resetting dynamic ranges, restarting WinNAT, reporting exclusions, and bind-testing 8000; its machine-specific scratch path is deliberately excluded from generated prose.",
    "The committed Discoverer PDCP assets include `access_key_wizard.py` Security Recon fields, `tools.py` default seeding, `agent_contracts.py` `.flw` redaction, `regen_secrets.py` scrubbing, and migration `0169_add_discoverer_cvemap_latest_demo_prompt.py` for the latest-CVE prompt.",
    "Discoverer hardening assets now include `discoverer.py` changes for `cvemap` -> `vulnx`, `.gitignore` Go-deny patterns, tracked `git_deny_go.py` guard/pre-commit installer, `discoverer_1000.py`, and `test_discoverer_thousand.py` validation coverage.",
    "The project inventory is rebuilt from `git ls-files` plus `git ls-files --others --exclude-standard` on every run, so any git-unignored local assets are counted while ignored runtime caches such as `Go/` remain outside the dossier.",
    "The newest frontend assets include `agent_page_chat.js`, `agent_page_ui.js`, `eslint.config.mjs`, and the Claude frontend/Multi-Turn/recent-fixes notes that document self-healing status frames and Create Flow name resolution.",
    "The `v1.40.1` port assets span committed `manage.py`/`config.json`, source/startserver integration, the committed 24-test `agent/test_django_port_config.py` suite, version surfaces in `package.json`/`VERSIONING.md`, and operator contracts across README, Book, self-knowledge, prompt guidance, and `docs/claude`.",
    "The newest agent assets are concrete: `agent/agents/zavuerer/`, migrations `0159_add_zavuerer.py` through `0164_dedup_zavuerer_setup_wizards.py`, `test_zavuerer_agent.py`, Access Keys Wizard Zavu fields, planner capability hints, wrapped-tool registration, and ACP canvas styling/connection support.",
    "The latest cleanup/assets span also includes `GEMINI.md`, `FirstFinalPlanToSpeedUp.md`, `docs/claude/recent-fixes.md`, response-parser/runtime/settings/middleware touch-ups, and removal of the stale `Tlamatini/db.sqlite3.bak-prereseat` backup from the tracked surface.",
    "The same wave preserves evidence-oriented tests and build guards such as `test_private_data_guard.py`, performance/visual checks, About-window authorship tests, and public-release verification rules that distinguish sensitive PII from valid creator names.",
    "The same span refreshes shipped visual/media assets: `TlamatiniAbout.png` replaces the old `TlamatiniAbout.jpg`, and `agent/images/TlamatiniAndKyber.mp4` is part of the repository asset set described by the dossier.",
    "The same recent window also retains the earlier self-modify/browser-setup asset wave — `copy_source_assets.py`, `agent/access_key_wizard.py`, `static/agent/js/access_keys_wizard.js`, `static/agent/css/access_keys_wizard.css`, and the Blender control surface in `agent/agents/blenderer/`.",
    "Key operator/runtime files such as `prompt.pmt`, `chat_agent_registry.py`, `tools.py`, `views.py`, `urls.py`, `manage.py`, `file_extractor.py`, and the File-Creator/File-Extractor templates also changed, so the visible features are backed by concrete implementation assets rather than documentation-only promises.",
    "Because the dossier already includes the full repository inventory (git-tracked files plus git-unignored working-tree additions) and the full line-count inventory, these named assets serve as the human-readable shortlist of what changed most materially in the latest release wave.",
]

PROMPT_CATALOG_GUIDE = [
    "Version `1.3.2` tightened the HTML answer contract with a Prime Directive on visual readability: explicit background and text color, no grey-on-dark body text, and safer table-body defaults.",
    "The seeded `Prompts` dropdown was also re-sorted into a learner path: context-only Q&A first, then metrics, files search, shell, code generation, vision, specialized single-tool actions, agent control, Unrealer, and heavier Multi-Turn/ACPX demos last.",
    "The `v1.35.0` prompt-search pass then makes that larger catalog easier to operate: prompt cards support substring, word-start, and fuzzy matching, with mode badges that keep one-shot, Multi-Turn, ACPX, Exec Report, and Step-by-Step demos visually distinct.",
    "Those readability rules remain in force in the current documentation set; `v1.42.0` carries faithful External-MCP structured-output delivery on top of v1.41.3 category grouping, physical duplicate removal, gap-tolerant loading, and ranked fuzzy search while preserving the broader operator context.",
]

SELF_KNOWLEDGE_GUIDE = [
    "Version `1.8.0` gives Tlamatini a first-person self-knowledge map in `Tlamatini/agent/Tlamatini.md`, so she can answer more accurately about her own architecture, runtime modes, open ports, pages, and internal capability surface.",
    "That self-reference is injected into all prompt chains through `prompt.pmt` and `agent/rag/config.py`, but it fails open if the file is missing or unreadable and it never overrides a user-loaded project when the request is a generic summary of the provided context.",
    "The language contract matters too: when a pronoun is used for Tlamatini in the documentation or prompt guidance, she is referred to as `she` / `her`, matching the updated handbook identity rules.",
]

SELF_MODIFY_GUIDE = [
    "Self-modification is a separate capability axis from frozen-vs-source runtime: only builds created with `python build.py --self-modify` bundle `TlamatiniSourceCode/` next to the application.",
    "When that directory is present, she can inspect and modify her own shipped source tree; when it is absent, she must say so plainly and fall back to her injected self-knowledge plus the surrounding docs.",
    "The build pipeline now announces that choice explicitly, so operators can tell whether a release is self-modify-capable before asking her to work on herself.",
]

MULTITURN_4096_GUIDE = [
    "The unified-agent loop now defaults to 4096 iterations instead of 256, giving long autonomous operator runs much more room before they exhaust the turn budget.",
    "That expansion is about conversational/tool-loop depth, not about blindly firing more tools at once: the planner’s selected-tool cap still keeps the tool surface bounded per request.",
    "Operationally, the bigger ceiling helps long workflows, while duplicate-call guards and the dedicated `chat_agent_sleeper` tool remain the antidote to accidental busy-polling loops.",
]

ASK_EXECS_GUIDE = [
    "Introduced in `v1.10.0` and still part of the current `v1.26.0` surface, `Ask Execs` is the Multi-Turn-only safety modifier that makes Tlamatini ask before each state-changing Tool, MCP, wrapped agent, or skill-backed execution instead of running it immediately.",
    "The permission dialog is explicit and auditable: it names the Tool or Agent family, the underlying raw tool name, the full parameters, the program or command to be executed, and the shell or execution surface involved.",
    "Proceed runs that one step and then prompts again at the next state-changing step; Deny halts the entire chain immediately and appends a red `Execution interrupted` banner even when Exec Report itself is off.",
]

ASK_EXECS_PIPELINE_GUIDE = [
    "Under the hood, `agent/exec_permission.py` provides `ExecPermissionBroker`, which lets the synchronous worker-thread executor emit a permission request onto the WebSocket event loop and then block on a `threading.Event` until the browser replies.",
    "The gate sits after deduplication and quota checks, so skipped calls never prompt and denied calls never appear as executed rows; only work that truly ran lands in Exec Report.",
    "The round-trip is fail-safe: browser disconnect, emit failure, cancel, or broker shutdown all resolve to Deny, so an unconfirmed state-changing action never slips through just because the UI vanished at the wrong time.",
]

WINDOWS_ATTENTION_GUIDE = [
    "The current Windows attention path is explicit and local: the browser calls `POST /agent/flash_window/`, and the backend routes that request through `agent/window_flash.py`.",
    "That helper can flash the `Tlamatini.exe` console/taskbar window with `FlashWindowEx` and always prints an uppercase attention banner, so the signal survives in `tlamatini.log` even when the browser is minimized.",
    "Two concrete reasons are wired today: Ask Execs execution approval prompts and Notifier notifications. The path is best-effort and fail-safe, degrading cleanly on non-Windows or windowless launches.",
]

WINDOWS_APP_REGISTRATION_GUIDE = [
    "Introduced in `v1.11.0` and still carried by the current `v1.26.0` release, the frozen install now behaves like a real Windows application: `install.py` writes a per-user HKCU Add/Remove Programs entry so Tlamatini appears in Settings -> Apps -> Installed apps and in the legacy Programs and Features list.",
    "The entry carries `DisplayName`, `DisplayVersion`, `InstallLocation`, `DisplayIcon`, `UninstallString`, `QuietUninstallString`, `NoModify`, `NoRepair`, and best-effort `EstimatedSize`, all pointing at the bundled `Uninstaller.exe` without requiring administrator rights.",
    "The matching runtime self-heal in `agent/apps.py` calls `windows_app_registration.self_heal_for_frozen()` on every frozen launch, so installs created before this feature existed can appear in Windows' uninstall UI after the next normal app start.",
]

UNREAL_EXTENDED_GUIDE = [
    "Unrealer’s documentation now points at the public `XAIHT/XaihtUnrealEngineMCP` fork, the Unreal Engine MCP variant developed specifically for Tlamatini.",
    "That fork exposes the extended 53-command, nine-category surface documented in README and Book: not only the base editor/Blueprint/UMG verbs, but also system, level, asset, material, screenshot, viewport, and in-editor Python paths.",
    "The practical message for operators is simple: Unrealer forwards the connected plugin’s verb surface directly, so Tlamatini’s client does not need a new release every time the plugin adds another supported command.",
]

REVIEWER_ANALYZER_GUIDE = [
    "The workflow catalog now includes two new high-value specialists: Reviewer and Analyzer, lifting the canvas inventory to 64 agents and the seed-skill catalog to 23 SKILL.md packages.",
    "Reviewer is LLM-powered: it resolves a git diff for `repo_path`, reviews it with a senior-engineer prompt, and emits an `INI_SECTION_REVIEWER` block whose first routable field is `verdict = APPROVE | REQUEST_CHANGES | COMMENT`.",
    "Analyzer is deterministic and scanner-driven: it runs whichever of `bandit`, `semgrep`, `ruff`, `eslint`, `gitleaks`, and `pip-audit` are installed on PATH over `target_path`, then emits an `INI_SECTION_ANALYZER` block with `status` and `total_findings` for downstream routing.",
]

REVIEWER_ANALYZER_SURFACES = [
    "Both agents always trigger `target_agents`, so a downstream Forker can branch on `{verdict}`, `{status}`, or `{total_findings}` instead of scraping prose.",
    "They are intentionally canvas-only workflow agents: there is no wrapped `chat_agent_reviewer` or `chat_agent_analyzer`, and therefore no duplicate Exec Report row family for those names.",
    "The chat-side counterparts live in the ACPX skill catalog instead: `code-review` exposes senior-engineer git-diff review, while `security-audit` exposes the deterministic multi-scanner sweep.",
]

REVIEWER_PRECISION_GUIDE = [
    "The `v1.4.1` Reviewer refinement tightened accuracy rather than adding a new surface: when `diff_ref` is empty, the review prompt labels the diff as the uncommitted working tree plus staged area, so the model must not describe those findings as already committed or pushed.",
    "The same patch teaches the model Tlamatini’s managed-secret convention: `agent/config.json` and selected `agent/agents/*/config.yaml` files can hold live local credentials in a keyed working copy, while `regen_secrets.py --mode push-able` scrubs them back to placeholders before commit.",
    "That guidance is mirrored into the `code-review` SKILL.md package too, keeping the chat-surface review behavior and the canvas Reviewer agent aligned on commit-state wording and secret-severity expectations.",
]

NATIVE_DIALOGS_GUIDE = [
    "The current tagged release `v1.4.2` removes Tkinter from the unstable runtime-facing dialog path and replaces it with `Tlamatini/agent/native_dialogs.py`, a native Windows dialog bridge used by browser-triggered pickers.",
    "This change pairs with the existing DB and operator dialogs: file and folder selection still feels local and GUI-first, but the fragile Tkinter dependency is no longer part of the interactive runtime path that users trigger from chat or ACP surfaces.",
    "The patch arrived with dedicated tests (`test_native_dialogs.py`) and with follow-up orphan-reaper/runtime adjustments, so the release reads as a stability pass rather than a cosmetic refactor.",
]

PLAYWRIGHTER_GUIDE = [
    "Playwrighter is the new 65th workflow agent in `v1.5.0`: a deterministic Playwright-powered browser automator for Chromium, Firefox, or WebKit that executes ordered interactive steps instead of static fetches.",
    "It covers the gap between Crawler and Googler: logins, multi-step forms, wizard clicks, JS-rendered SPA scraping, screenshots after interaction, assertions, downloads, and authenticated end-to-end UI checks.",
    "The agent emits `INI_SECTION_PLAYWRIGHTER` with `start_url`, `final_url`, `status`, `steps_run`, `assert_result`, and extracted values so downstream Forker or Parametrizer logic can branch on pass/fail or reuse scraped data.",
]

PLAYWRIGHTER_SURFACES_GUIDE = [
    "Two operator surfaces ship in lock-step: the wrapped Multi-Turn tool `chat_agent_playwrighter` takes the whole script as `steps_json`, while the visual Playwrighter canvas node stores the same declarative step list in YAML.",
    "Session continuity is built in: `headless: false` lets operators watch it drive, and `storage_state_in` / `storage_state_out` carry login state across runs without forcing a manual browser setup each time.",
    "Because Playwrighter is state-changing, its executions appear in Exec Report and it always triggers `target_agents` whether the run succeeds or fails.",
]

WINDOWER_GUIDE = [
    "Windower is the new 66th workflow agent: a deterministic Win32 window manager that finds an application window by title and performs one lifecycle operation on the window itself instead of clicking inside it.",
    "It supports `focus`, `minimize`, `maximize`, `restore`, `move`, `resize`, `move_resize`, `close`, `topmost`, `untopmost`, `arrange`, and `list`, with `substring`, `exact`, or `regex` title matching plus `match_index` for duplicate titles.",
    "The agent emits `INI_SECTION_WINDOWER` with `action`, `window_title`, `matched`, `match_count`, `state`, `left`, `top`, `width`, and `height`, so downstream Forker or Parametrizer logic can branch on presence, state, or geometry.",
]

WINDOWER_SURFACES_GUIDE = [
    "Two operator surfaces ship in lock-step: the wrapped Multi-Turn tool `chat_agent_windower` accepts free-form key=value requests, while the visual Windower canvas node stores the same operation fields in YAML.",
    "Windower is the desktop-window sibling of Mouser and Keyboarder: use Windower when the goal is the window as a whole, Mouser for controls inside it, and Keyboarder for text entry into it.",
    "Because Windower changes real window state, its executions appear in Exec Report and it always triggers `target_agents` whether the action succeeds or fails.",
]

KALIER_GUIDE = [
    "Kalier is Tlamatini’s current Kali Linux bridge: a direct client for MCP-Kali-Server that lets her drive authorized recon, enumeration, web scanning, and offensive-security workflows from chat or canvas.",
    "It can issue one capability per run, including `nmap`, `gobuster`, `dirb`, `nikto`, `sqlmap`, `metasploit`, `hydra`, `john`, `wpscan`, `enum4linux`, arbitrary `command`, or a safe `health` probe of the remote server.",
    "The agent emits `INI_SECTION_KALIER` with `action`, `endpoint`, `subject`, `return_code`, `success`, `timed_out`, and `server_url`, so downstream Forker or Parametrizer logic can branch on results without scraping prose.",
    "The practical operator result is simpler prompts: after the Kali box URL is configured once, normal Multi-Turn requests no longer repeat `server_url` unless you intentionally override it for a one-off target.",
]

KALIER_SURFACES_GUIDE = [
    "Two operator surfaces ship in lock-step: the wrapped Multi-Turn tool `chat_agent_kalier` now auto-injects the configured `kali_server_url`, while the visual Kalier canvas node still stores `server_url` explicitly in YAML per node.",
    "Kalier talks straight to the Kali-side Flask API over HTTP using Python-stdlib `urllib`, so it stays self-contained in the pool subprocess and works the same in source or frozen builds; the embedded-client injection fails open if the config value is blank or unreadable.",
    "Authorized use only: the intended operator flow is an in-scope lab, CTF, or permitted engagement, often with `ssh -L 5000:localhost:5000 user@KALI_IP` tunneling a remote Kali box back to `http://127.0.0.1:5000`, and the repo now includes `Tlamatini-Kali-Setup.md` for the zero-client walkthrough.",
]

STM32ER_GUIDE = [
    "STM32er is Tlamatini's critical-mission STM32 bridge. The established route talks to the `STM32 Template Project MCP` for STM32F407, while released `v1.42.0` adds a direct PlatformIO route for supported mainstream STM32 boards without removing that flow.",
    "Both backends keep zero-config intent: the template route bootstraps its MCP dependencies, and the new route resolves or installs PlatformIO Core in the same shared per-user location used by ESP32er.",
    "Backend-specific preflight preserves target safety: compile-only work may proceed without attached hardware, but upload, reset, serial, and live operations are refused when the required toolchain, project, board mapping, programmer, or ST-LINK evidence is missing.",
]

STM32ER_SURFACES_GUIDE = [
    "Two operator surfaces ship in lock-step: the wrapped Multi-Turn tool `chat_agent_stm32er` takes one `action` per call, while the visual STM32er canvas node stores the same fields in YAML and triggers downstream agents on both success and failure.",
    "The tool surface retains the 23 MCP verbs and locally adds PlatformIO environment, board, project, source, build, flash, monitor, package, QA, artifact, and safe scaffold composites; every run still emits one `INI_SECTION_STM32ER` block for routing.",
    "Config seeding supplies both MCP and PlatformIO defaults, while `stm32_backend=auto`, `board`, and `device` keep the normal prompt focused on firmware intent. Newest-silicon CubeCLT/N6 handling remains proposal-only and is not advertised as implemented.",
]

ESP32ER_GUIDE = [
    "ESP32er is Tlamatini’s current PlatformIO Core bridge: she can scaffold, author, build, upload, and monitor ESP32-class firmware without relying on an external MCP server or IDE.",
    "The operator promise is zero-config bootstrap: leave `pio_executable` blank and ESP32er downloads or pip-falls-back to PlatformIO Core on first use, validates it, and caches it under a per-user directory so the user installs only the board USB driver plus Tlamatini.",
    "Before any build-or-upload action, ESP32er runs a serial-aware preflight over `pio` resolution, project shape, and connected ports; upload and monitor require a real serial device, while non-espressif32 targets are warned about rather than hard-refused because PlatformIO is intentionally multi-target.",
]

ESP32ER_SURFACES_GUIDE = [
    "Two operator surfaces ship in lock-step: the wrapped Multi-Turn tool `chat_agent_esp32er` takes one `action` per call, while the visual ESP32er canvas node stores the same fields in YAML and triggers downstream agents on both success and failure.",
    "The tool surface covers environment/meta (`bootstrap`, `validate`, `system_info`, `boards`), project lifecycle (`create_project`, `write_source`, `read_source`, `list_sources`, `clean`), build and flash (`build`, `upload`, `build_and_upload`, `list_artifacts`), serial HIL (`device_list`, `monitor`, `monitor_session`), and package / QA paths (`pkg_install`, `pkg_list`, `pkg_update`, `check`, `test`), with every run emitting an `INI_SECTION_ESP32ER` block for Forker or Parametrizer routing.",
    "Config -> URLs now seeds the chat path with `pio_executable` and `pio_core_dir`, so firmware prompts usually describe only the board, project, and task while the wrapped tool auto-injects the PlatformIO runtime plumbing.",
]

ESP32_TEMPLATE_GUIDE = [
    "The new `ESP32TemplateProject` repository is the known-good baseline documented in BookOfTlamatini’s bonus chapter: a plain PlatformIO project, not a server, meant to prove an ESP32 board and toolchain are healthy before larger firmware work.",
    "It mirrors ESP32er’s grain: `platformio.ini`, `src/`, `include/`, `lib/`, `test/`, a blinking `main.cpp`, serial output at 115200, and GitHub-ready docs/CI so the reference project does not silently rot.",
    "Operationally, ESP32er can either point at a checkout of that template (`project_dir`) or scaffold an equivalent from scratch with `action='create_project'`, then carry the directory through build, upload, and monitor steps.",
]

ESPHOMER_GUIDE = [
    "ESPHomer is Tlamatini’s ESPHome bridge: she can author YAML-based smart-home device configs, validate them, compile firmware, upload over USB or OTA, and observe logs without introducing an extra MCP server or IDE.",
    "The operator promise is zero-config bootstrap: leave `esphome_executable` blank and ESPHomer installs or resolves ESPHome on first use, so the user installs only the board USB driver plus Tlamatini.",
    "Before any compile-or-upload action, ESPHomer runs a fail-safe preflight over `esphome` resolution, YAML existence, and serial-or-OTA readiness; upload, run, and logs require a real serial board or an OTA host because the first flash is always USB.",
]

ESPHOMER_SURFACES_GUIDE = [
    "Two operator surfaces ship in lock-step: the wrapped Multi-Turn tool `chat_agent_esphomer` takes one `action` per call, while the visual ESPHomer canvas node stores the same fields in YAML and triggers downstream agents on both success and failure.",
    "The tool surface covers environment/meta (`bootstrap`, `validate`, `version`), device-YAML lifecycle (`new_config`, `write_config`, `read_config`, `config`, `clean`), build and flash (`compile`, `upload`, `run`, `list_artifacts`), bounded observation (`logs`), and the one-shot `scaffold_compile_upload` lifecycle.",
    "Config -> URLs now seeds the chat path with `esphome_executable`, so smart-home firmware prompts usually describe only the device, board, and task while the wrapped tool auto-injects the ESPHome runtime plumbing.",
]

ESPHOME_TEMPLATE_GUIDE = [
    "The bundled `ESPHomeTemplateProject` is the known-good ESPHome baseline documented in README and Book: a phone-controllable light with native `api`, `ota`, `wifi`, and a board LED output, shipped as `agent/agents/esphomer/ESPHomeTemplateProject/tlamatini-light.yaml`.",
    "It matters because ESPHomer’s source-of-truth is a YAML device file, not a C++ project tree: that sample proves the first real workflow of `new_config` -> `config` -> `compile` -> `upload` without forcing users to invent a valid starter config from memory.",
    "Operationally, the sample closes the gap between the new agent and a practical first build a user can validate, compile, flash, and then control from a smart-home hub such as Home Assistant.",
]

DESIGN_PRINCIPLES = [
    "Evidence-first answers: Tlamatini grounds responses in selected project context and hybrid retrieval rather than freeform model memory.",
    "Explicit orchestration: checked Multi-Turn uses a visible tool loop, capability scoring, and staged planning instead of a single opaque call.",
    "Operational reversibility: risky changes such as database replacement are staged, archived, and delayed to the only safe window instead of hot-swapped mid-session.",
    "Fail-open diagnostics: GPU pressure probes and session-restore safeguards warn early without breaking CPU-only or degraded environments.",
    "Runtime isolation: wrapped chat-agent copies run in session-scoped folders so template agents remain pristine while live runs stay inspectable.",
    "Self-knowledge with scope discipline: she can talk accurately about herself without letting self-reference override a user-loaded project context.",
    "Operator truth over vibes: Exec Report tables, tlamatini.log, skill audits, and ACPX transcripts make the system auditable after execution.",
]

INSTALLATION_GUIDE = [
    "The easiest path for most users is the packaged installer from GitHub Releases: no manual Python install is required because the release already carries Python 3.12.10 and the project dependencies.",
    "Source mode remains the developer path: clone the repo, create a virtual environment, install `requirements.txt`, run migrations, create a superuser, collect static files, and then launch Django.",
    "Packaged Windows installs open the browser at `http://127.0.0.1:8000/` and create the default `user / changeme` login; manual source installs use your own `createsuperuser` account instead.",
    "Port 8000 is only the default: `config.json`'s `django_port` moves the web UI to any free port on every launch path, which is the fix when Windows or Hyper-V has RESERVED port 8000 and startup fails with `WinError 10013`.",
    "When a newer packaged release exists, the intended upgrade path is in-app: `About -> Check for updates`, not manual folder replacement.",
    "You can run either the checked-in cloud/back-end defaults from `Tlamatini/agent/config.json` or a local/remote Ollama-backed configuration with matching model names.",
]

CONFIGURATION_GUIDE = [
    "Source mode resolves `Tlamatini/agent/config.json`; frozen builds resolve `config.json` next to the executable; `CONFIG_PATH` overrides both.",
    "Core keys include `embeding-model`, `chained-model`, `ollama_base_url`, `ollama_token`, `enable_unified_agent`, `unified_agent_model`, and `unified_agent_max_iterations`.",
    "The checked-in default model baseline moved again in the recent Git window: the shared config now favors `glm-5.2:cloud`, so source or frozen installs that keep the shipped config should be documented as cloud-first unless the operator intentionally swaps models.",
    "URL configuration now also includes `kali_server_url`, the STM32er bootstrap fields `stm32_mcp_server_script`, `stm32_mcp_python`, `stm32_template_dir`, `stm32_ide_root`, `stm32_mcp_repo_url`, and `stm32_mcp_install_dir`, plus ESP32er’s `pio_executable` and `pio_core_dir`, all edited from `Config -> URLs` and inherited automatically by the chat-side wrapped tools.",
    "Credential configuration is no longer hand-edit-only: Config -> Access Keys Wizard provides a browser-side path for ACPX, provider secrets, unified messaging, and Security Recon (ProjectDiscovery) keys such as `pdcp_api_key` while preserving masked status in the UI.",
    "The chat-side Config -> Models and Config -> URLs dialogs are now first-class configuration surfaces, and they can explicitly ask the operator to reconnect when saved values change live-session assumptions.",
    "The separate DB dropdown is not a config editor: it is a maintenance surface for copying the live SQLite database out or staging a replacement for the next full start-up.",
    "Multi-Turn is toggled from the chat toolbar, but it depends on the unified-agent configuration and the selected model/base-url pairing being valid; the current default iteration ceiling is 4096, and Ask Execs only becomes available when Multi-Turn itself is on.",
    "Image-Interpreter now uses three model slots from Config -> Models: `image_interpreter_model` for the first forensic interpreter, `image_interpreter_model_2` for the second holistic interpreter, and `image_merging_model` for the final report merger.",
]

START_HERE_GUIDE = [
    "BookOfTlamatini now leads with a five-step onboarding path because the easiest way to succeed with Tlamatini is to treat setup as one guided sequence instead of reading the whole handbook first.",
    "Recommended path for most operators: install the packaged release from GitHub Releases, launch the Start-menu shortcut, and let the bundled Python 3.12.10 plus dependencies carry the runtime without asking the user to install Python manually.",
    "Developer path stays available: clone the repo, create a virtual environment, install `requirements.txt`, run migrations, and start Django with `python Tlamatini/manage.py runserver` (the `--noreload` flag is optional since 2026-07-11 — plain `runserver` now boots clean and auto-reloads).",
    "After the app opens, the first in-app surfaces that matter are `Config -> Models`, `Config -> URLs`, and `Config -> Access Keys Wizard`; those now form the real beginner path, not manual JSON editing.",
    "For ordinary question-answering keep Multi-Turn off; turn it on only when you want Tlamatini to execute tools, wrapped agents, or remote MCP capabilities instead of answering directly.",
]

FIRST_RUN_CONFIG_GUIDE = [
    "Three new onboarding screenshots now anchor the first-run path: `Tlamatini/agent/images/MenuConfig.jpg`, `ConfigureModels.jpg`, and `ACPXKeysConfigureWizard.jpg`.",
    "`Config -> Models` is the place where operators map embedding, chat, vision, and auxiliary model names to what Ollama actually exposes on the host machine.",
    "`Config -> Access Keys Wizard` now carries an especially important rule: a localhost Ollama usually needs no Ollama token, while a remote Ollama endpoint may require one.",
    "Saved model, URL, or credential changes can invalidate the assumptions of the current session, so reconnecting after major Config edits is now part of the honest operator guidance.",
]

RUNNING_GUIDE = [
    "Development server: `python Tlamatini/manage.py runserver` — the `--noreload` flag is now OPTIONAL (plain `runserver` boots clean and auto-reloads since the 2026-07-11 `agent/apps.py` reloader-aware fix; it used to double-start the MCP helper ports :8765 / :50051 and crash with WinError 10048).",
    "Preferred async/dev bootstrap: `python Tlamatini/manage.py startserver`, which starts MCP services before the Django server.",
    "Production-style ASGI entrypoint: `daphne -b 127.0.0.1 -p 8000 tlamatini.asgi:application` — this path bypasses `manage.py`, so it does not read `django_port` and the port must be passed on the command line.",
    "Current startup also re-applies GPU-performance / Ollama-pinning hooks in the background on supported NVIDIA Windows hosts, so restart-time behavior stays closer to the tuned development baseline.",
    "That same early startup window is also where a staged `DB/ToLoad/db.sqlite3` file is promoted into the live database path, before Django opens SQLite.",
    "In frozen mode, startup now also self-heals the per-user Windows Installed-apps entry when `Uninstaller.exe` is present beside `Tlamatini.exe`, so old installs can gain a standard uninstall surface without a reinstall.",
    "Startup now also prints a `--- [VERSION] Tlamatini ...` banner, making the running build visible in both the console and `tlamatini.log` without an HTTP call.",
    "Startup cleans pool state, repopulates the agent registry, launches MCP metrics/file-search servers, and then serves HTTP plus WebSocket traffic.",
]

DB_MENU_GUIDE = [
    "The new DB dropdown gives operators two GUI-first maintenance paths: `Backup database` copies the live SQLite file out, and `Set DB` stages a chosen `db.sqlite3` for the next full start-up.",
    "Both dialogs are live-validated in the browser and now expose Browse buttons that open native host-side pickers for folders or `db.sqlite3` files.",
    "Backup is read-only and uses the currently live database path; Set DB never hot-swaps the file mid-session because Django already holds SQLite open.",
]

DB_SWAP_GUIDE = [
    "Set DB writes the selected file to `DB/ToLoad/db.sqlite3`; the real swap happens only at the top of `manage.py` before Django imports anything.",
    "When that swap runs, the previous live database is moved into `DB/Older/<timestamp>/db.sqlite3`, creating a built-in rollback trail instead of overwriting history.",
    "Reconnect is not enough for this path: the operator must fully restart Tlamatini so the pre-Django swap window opens again.",
    "If the staged file is bad or locked, startup fails open: Tlamatini logs the error and continues with the previous live database.",
]

VERSIONING_GUIDE = [
    "Tlamatini now follows Semantic Versioning 2.0.0 with git tags as the single source of truth: you tag, then you build, instead of hand-editing version strings across files.",
    "The build path resolves a version once and propagates it into generated runtime metadata, Win32 VERSIONINFO resources, and the release-folder naming convention.",
    "On untagged commits the resolver falls back honestly to a git-derived development version instead of failing the build or lying about the release state.",
]

VERSION_SURFACES_GUIDE = [
    "Operators can now see the running version in the About dialog, the startup banner, the open `GET /agent/version/` health-check endpoint, and Windows file properties on the built executables.",
    "Packaged installs also project the same release identity into Windows' Installed-apps metadata through the ARP `DisplayVersion` field, so the uninstall surface and the binaries agree on the version being removed.",
    "The runtime resolver lives in `Tlamatini/agent/version.py`, while the build-oriented coordination logic lives in the repo-root `versioning.py` and the longer policy notes live in `VERSIONING.md`.",
    "Build overrides are supported in a predictable order: CLI `--version`, then `TLAMATINI_VERSION`, then git describe, then the sentinel `0.0.0+unknown`.",
]

DE_COMPRESSER_GUIDE = [
    "De-Compresser is the deterministic archive worker for compression and decompression tasks, deciding direction from whichever side exposes a recognized archive extension.",
    "Supported archive families are `.gz`, `.zip`, `.7z`, `.tar.gz`, and `.gz.tar`; file-to-folder extraction and file-or-directory packing are both documented in the README and Book.",
    "Password handling is explicit: `passwordless=true` skips it, while `passwordless=false` requires the `DE_COMPRESSER_PWD` environment variable and fails fast if the secret is missing.",
]

DE_COMPRESSER_INTEGRATION_GUIDE = [
    "The agent is reachable from both operator surfaces: ACP canvas nodes wire through dedicated connection-update views, and checked Multi-Turn can invoke it through `chat_agent_de_compresser`.",
    "Format engines are practical rather than magical: stdlib `gzip` and `zipfile` cover core cases, `7z` is preferred for encrypted `7z` or `zip`, and `py7zr` was added to `requirements.txt` as the Python fallback.",
    "Every run emits an `INI_SECTION_DE_COMPRESSER<<< ... >>>END_SECTION_DE_COMPRESSER` block and still triggers `target_agents`, so downstream Parametrizer or Raiser logic can branch on `success=true|false` instead of guessing from prose.",
]

UNREAL_MCP_GUIDE = [
    "Unreal MCP is a UE5 plugin that runs inside the editor and listens on `127.0.0.1:55557` for one JSON command per TCP connection; Tlamatini is the client side of that link, not the plugin host.",
    "The Unrealer workflow agent forwards any command the connected plugin build exposes, up to a 53-command surface across nine categories: actor manipulation (incl. viewport screenshots), Blueprint authoring and node wiring, input mappings, UMG widget building, in-editor Python/console execution, level I/O, asset import, and material authoring. Headless build/cook/test is out of scope (it needs UnrealEditor-Cmd, not the editor socket).",
    "The same integration is available in both operator surfaces: checked Multi-Turn via `chat_agent_unrealer`, and ACP canvas flows via the visual Unrealer node plus Parametrizer chaining.",
]

UNREAL_INSTALL_GUIDE = [
    "Install the upstream plugin into `<YourProject>/Plugins/UnrealMCP/`, enable it from `Edit -> Plugins`, restart UE5, and confirm the Output Log says `UnrealMCP listening on 127.0.0.1:55557`.",
    "Tlamatini does not compile or embed the plugin; the Unreal editor must already be running with the listener bound before any Unrealer call can succeed.",
    "A seeded smoke test ships in the Prompts table: `Unreal MCP End-to-End Editor Drive` walks through sanity-check, actor spawn, Blueprint creation/compile, and UMG widget assembly.",
]

UNREAL_RUNTIME_GUIDE = [
    "Each Unrealer run is one command: load `config.yaml`, open a fresh TCP socket, send `{\"type\": command, \"params\": params}`, read until valid JSON arrives, and log one atomic `INI_SECTION_UNREALER<<< ... >>>END_SECTION_UNREALER` block.",
    "The wrapped-tool path stores chat runs under `agent/agents/pools/_chat_runs_/unrealer_<seq>_<id>/`, while visual flows use normal `unrealer_<n>` pool folders and can chain response fields through Parametrizer mappings.",
    "Exec Report treats Unrealer as its own row family, and troubleshooting stays concrete: connection-refused means the plugin is not listening, while read-timeout usually means UE5's game thread is busy.",
]

ORPHAN_REAPER_GUIDE = [
    "Tlamatini now ships a three-tier orphan reaper in `Tlamatini/agent/orphan_reaper.py` focused on Windows console-host leftovers such as `conhost.exe` and `openconsole.exe`.",
    "Tier 1 runs after spawn-capable Multi-Turn tool calls, Tier 2 runs once after the final answer in a background thread, and Tier 3 runs again during application shutdown through the same cleanup path that already tears down pools.",
    "The candidate set stays narrow on purpose: dead descendants of the current process tree, orphaned console hosts tied to that tree, and pool-linked processes whose `cmdline` still points into `agents/pools/...` even though tracking records are gone.",
]

ORPHAN_PREVENTION_GUIDE = [
    "Prevention landed alongside cleanup: Windows spawn sites now use `CREATE_NO_WINDOW | DETACHED_PROCESS | CREATE_NEW_PROCESS_GROUP` with stdio piped to `DEVNULL`, so the console host is usually never allocated in the first place.",
    "The ACPX runtime now kills full process trees instead of only top-level wrappers, and pool-agent scripts gained a conservative `subprocess.Popen` guard so forgotten descendants still inherit `CREATE_NO_WINDOW`.",
    "If something survives both cleanup tiers, Tlamatini sends a second chat bubble listing each surviving `name + PID`; the reaper never raises into the user path because a cleanup crash would be worse than the leftovers it tried to remove.",
]

EMBEDDING_GUARD_GUIDE = [
    "README and BookOfTlamatini now document an embedding-memory pre-flight guard for GPU hosts before a directory-context load starts its FAISS embedding burst.",
    "On supported NVIDIA hosts it estimates embedding-model VRAM pressure, and when the projected load is too high it emits a non-blocking warning chat bubble instead of silently letting RAM<->VRAM thrash surprise the operator.",
    "CPU-only, AMD, and Apple-Silicon hosts stay fail-open: the guard becomes a no-op when the NVIDIA probe does not apply.",
    "The practical operator response is straightforward: switch to a smaller embedding model, reconnect if needed, or proceed knowingly with the heavier model.",
]

RECENT_RUNTIME_SAFEGUARDS = [
    "Config -> Models and Config -> URLs dialogs now track their pre-edit baseline and can show a reconnect-required dialog when the saved values change what the live chat session should trust.",
    "The restored-session autoload path now buffers early WebSocket frames so context-loading spinners and disabled-input state are not lost during automatic reconnect/restore flows.",
    "Startup and restart behavior now also re-apply GPU performance and Ollama keep-alive hooks in the background on supported NVIDIA Windows hosts, improving warm-model readiness without blocking Django boot.",
    "Browser-driven attention routing is now part of that operator-safety layer too: Ask Execs prompts and Notifier events can raise a taskbar flash plus a log banner without relying on the browser window already being visible.",
    "Windows process hygiene is now part of that safety story too: detached no-window spawns and the three-tier orphan reaper reduce the chance that Task Manager shows stale Tlamatini-icon console helpers after long runs.",
    "Ask Execs extends that safety story into execution approval itself: the operator can now stop a destructive chain before the next mutation instead of only auditing it after the fact in Exec Report.",
]

RELEASE_GUIDE = [
    "Release production is a three-step pipeline: `build.py` -> `build_uninstaller.py` -> `build_installer.py`.",
    "The final distributable is the full versioned release folder `dist/Tlamatini_Release_v<version>/`, not a stray executable copied outside its payload.",
    "Use `build.py --self-modify` when you intentionally want a release that ships `TlamatiniSourceCode/` so she can inspect or modify herself at runtime.",
    "Current `build.py` treats `README.md` and `jd-cli/` as required post-build assets and fails hard if those payloads are missing.",
    "Bundled support scripts cover shortcut creation/removal, `.flw` association, the PowerShell launcher, Windows-specific installer ergonomics, and the per-user Installed-apps registration path.",
]

EXEC_REPORT_GUIDE = [
    "Exec Report is a Multi-Turn-only transparency layer that appends one operation table per state-changing agent family to the final answer.",
    "Rows are recorded from the live tool-call stream rather than guessed from the LLM prose, so the report is the operational ground truth.",
    "Each row receives a SUCCESS/FAILURE verdict from raw tool returns, making long installs, deployments, and remediations inspectable after the fact.",
    "When Ask Execs is enabled, Exec Report and the red denial banner complement each other: already-executed steps still render as tables, while the denied step stays out of the tables and is surfaced only through the interruption banner.",
]

ACPX_GUIDE = [
    "ACPX lets Tlamatini spawn external coding-agent CLIs such as Codex, Claude Code, Cursor, Gemini, Qwen, and others as managed child processes.",
    "It pairs those agents with markdown-driven `SKILL.md` packages, validated I/O contracts, permission gating, and append-only audit logs.",
    "Transport-aware drain rules and bounded event bodies reduce latency while protecting the LLM context budget during external-agent relays.",
    "The operator-facing references are `README.md` and `ACPX.md`, while the implementation lives under `agent/acpx/`, `agent/skills/`, and `agent/skills_pkg/`.",
]

APP_LOG_GUIDE = [
    "The built-in `tlamatini.log` file captures both stdout and stderr through a tee stream initialized in `manage.py` before Django starts.",
    "In source mode the log sits next to `manage.py`; in frozen mode it lives next to the executable.",
    "Immediate flush behavior makes the log the primary forensic artifact for startup problems, warnings, tracebacks, and runtime diagnostics.",
]

OLLAMA_COMMANDS = "\n".join(
    [
        '$env:OLLAMA_INSTALL_DIR = "$env:LOCALAPPDATA\\Programs\\Ollama"',
        "irm https://ollama.com/install.ps1 | iex",
        "ollama --version",
        "ollama signin",
        "ollama serve",
        "Invoke-WebRequest http://127.0.0.1:11434/api/tags -UseBasicParsing",
        "ollama pull Nomic-Embed-Text:latest",
        "ollama pull glm-5.2:cloud",
        "ollama pull qwen3.5:cloud",
        "ollama pull gpt-oss:120b-cloud",
        "ollama pull qwen3.5:397b-cloud",
        "ollama pull glm-5.1:cloud",
    ]
)

OLLAMA_GUIDE = [
    "Open a normal PowerShell window, not an elevated one, for the safest no-admin Windows installation path.",
    "Install into `%LOCALAPPDATA%\\Programs\\Ollama` with the official PowerShell installer script and then reopen PowerShell so PATH updates are visible.",
    "If you plan to use the default `:cloud` models from the shipped config, sign in with `ollama signin` so the host is linked to your Ollama account before you test the app.",
    "Verify the CLI with `ollama --version`, start `ollama serve` if the background service is not already active, and confirm `http://127.0.0.1:11434/api/tags` responds.",
    "Pull the default repository model tags exactly as written if you want the shipped config and agent templates to work unchanged.",
    "The Book now clarifies the token rule: a localhost Ollama usually needs no Ollama bearer token in Tlamatini, while a remote Ollama endpoint may require one in `Config -> Access Keys Wizard` or the matching config key.",
]

ARCHITECTURE_LAYERS = [
    ("Browser interfaces", "Chat page plus Agentic Control Panel templates and JavaScript modules."),
    ("Django/Channels", "Authentication, views, WebSockets, session state, message persistence, and ASGI startup."),
    ("RAG and context", "Metadata extraction, text splitting, FAISS/BM25 retrieval, context budgeting, and fallback behavior."),
    ("Multi-Turn engine", "Capability registry, global execution planner, explicit tool loop, answer parsing, and answer-success classification."),
    ("Tools and agents", "Core tools, MCP context providers, wrapped chat-agent launchers, and the current visual workflow agent templates."),
    ("Packaging", "PyInstaller build scripts, shortcut registration, `.flw` association, installer, uninstaller, and release folder assembly."),
]

AGENT_CATEGORIES = [
    ("Control", "starter, ender, stopper, cleaner, barrier, flowbacker"),
    ("Execution and files", "executer, pythonxer, pser, file_creator, file_extractor, file_interpreter, de_compresser, playwrighter, windower, unrealer, kalier, stm32er, esp32er, esphomer, arduiner, mover, deleter"),
    ("DevOps and infra", "gitter, dockerer, kuberneter, jenkinser, ssher, scper"),
    ("Data and APIs", "sqler, mongoxer, apirer, crawler, googler"),
    ("Monitoring and routing", "monitor_log, monitor_netstat, flowhypervisor, forker, asker, counter, and, or"),
    ("Communication", "notifier, emailer, recmailer, telegrammer, teletlamatini, whatsapper"),
    ("Security and media", "kyber_keygen, kyber_cipher, kyber_decipher, image_interpreter, video_analyzer, shoter, camcorder, recorder, audioplayer, videoplayer, j_decompiler"),
    ("Workflow intelligence", "flowcreator, gatewayer, gateway_relayer, node_manager, parametrizer, prompter, summarizer, acpxer"),
]


def pdf_styles() -> dict[str, ParagraphStyle]:
    base = getSampleStyleSheet()
    return {
        "title": ParagraphStyle(
            "TlamatiniTitle",
            parent=base["Title"],
            alignment=TA_CENTER,
            fontName="Helvetica-Bold",
            fontSize=28,
            leading=34,
            textColor=colors.HexColor("#17342d"),
            spaceAfter=14,
        ),
        "subtitle": ParagraphStyle(
            "TlamatiniSubtitle",
            parent=base["Normal"],
            alignment=TA_CENTER,
            fontName="Helvetica",
            fontSize=12,
            leading=16,
            textColor=colors.HexColor("#6b4a34"),
            spaceAfter=20,
        ),
        "h1": ParagraphStyle(
            "TlamatiniH1",
            parent=base["Heading1"],
            fontName="Helvetica-Bold",
            fontSize=18,
            leading=22,
            textColor=colors.HexColor("#0f3b31"),
            spaceBefore=14,
            spaceAfter=8,
        ),
        "h2": ParagraphStyle(
            "TlamatiniH2",
            parent=base["Heading2"],
            fontName="Helvetica-Bold",
            fontSize=13,
            leading=16,
            textColor=colors.HexColor("#8f5c35"),
            spaceBefore=10,
            spaceAfter=6,
        ),
        "body": ParagraphStyle(
            "TlamatiniBody",
            parent=base["BodyText"],
            fontName="Helvetica",
            fontSize=9.5,
            leading=12.5,
            textColor=colors.HexColor("#1f2933"),
            spaceAfter=5,
        ),
        "bullet": ParagraphStyle(
            "TlamatiniBullet",
            parent=base["BodyText"],
            fontName="Helvetica",
            fontSize=9.2,
            leading=12.2,
            leftIndent=13,
            firstLineIndent=-8,
            textColor=colors.HexColor("#1f2933"),
            spaceAfter=4,
        ),
        "mono": ParagraphStyle(
            "TlamatiniMono",
            parent=base["Code"],
            fontName="Courier",
            fontSize=6.7,
            leading=7.7,
            textColor=colors.HexColor("#17231f"),
        ),
    }


def p(text: str, style: ParagraphStyle) -> Paragraph:
    return Paragraph(text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), style)


def bullet(text: str, style: ParagraphStyle) -> Paragraph:
    return p(f"- {text}", style)


def _table_cell(text: str, font_size: float, *, header: bool) -> Paragraph:
    """Wrap a string cell in a Paragraph so ReportLab word-wraps it inside the
    column width instead of letting a long single-line string overflow the page.
    Long unbreakable tokens (e.g. deep slash-paths) are force-split by
    Paragraph's default splitLongWords behavior."""
    style = ParagraphStyle(
        "TableHeaderCell" if header else "TableBodyCell",
        fontName="Helvetica-Bold" if header else "Helvetica",
        fontSize=font_size,
        leading=font_size + 2,
        textColor=colors.white if header else colors.HexColor("#1f2933"),
    )
    safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return Paragraph(safe, style)


def table(data: list[list], widths: list[float] | None = None, font_size: int = 8) -> Table:
    # Cells are wrapped in Paragraphs (the only flowable that word-wraps within a
    # column). Plain string cells passed straight to Table render on a single
    # line and overflow narrow columns; Paragraphs wrap to the column width.
    wrapped = [
        [
            _table_cell(cell, font_size, header=(row_index == 0)) if isinstance(cell, str) else cell
            for cell in row
        ]
        for row_index, row in enumerate(data)
    ]
    tbl = Table(wrapped, colWidths=widths, repeatRows=1)
    tbl.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#17342d")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), font_size),
                ("LEADING", (0, 0), (-1, -1), font_size + 2),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#a8b1aa")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#f8faf7"), colors.HexColor("#eef4ef")]),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    return tbl


def split_lines(text: str, per_page: int) -> list[str]:
    lines = text.splitlines()
    return ["\n".join(lines[index : index + per_page]) for index in range(0, len(lines), per_page)]


def split_items(items: list, size: int) -> list[list]:
    return [items[index : index + size] for index in range(0, len(items), size)]


def pdf_page_footer(canvas, doc) -> None:
    canvas.saveState()
    canvas.setStrokeColor(colors.HexColor("#8f5c35"))
    canvas.setLineWidth(0.4)
    canvas.line(doc.leftMargin, 0.48 * inch, A4[0] - doc.rightMargin, 0.48 * inch)
    canvas.setFillColor(colors.HexColor("#17342d"))
    canvas.setFont("Helvetica", 7)
    canvas.drawString(doc.leftMargin, 0.32 * inch, "Tlamatini complete project dossier")
    canvas.drawRightString(A4[0] - doc.rightMargin, 0.32 * inch, f"Page {doc.page}")
    canvas.restoreState()


def cover_image_flowable(image_path: Path, max_width: float, max_height: float) -> Image:
    """Return a reportlab Image scaled to fit WITHIN ``max_width`` x ``max_height``
    (points) while PRESERVING the source image's natural aspect ratio — so the
    cover is letter-boxed, never stretched. Falls back to the box size if the
    image's dimensions cannot be read."""
    try:
        src_w, src_h = ImageReader(str(image_path)).getSize()
    except Exception:
        src_w = src_h = 0
    if src_w <= 0 or src_h <= 0:
        img = Image(str(image_path), width=max_width, height=max_height)
    else:
        scale = min(max_width / src_w, max_height / src_h)
        img = Image(str(image_path), width=src_w * scale, height=src_h * scale)
    img.hAlign = "CENTER"
    return img


def build_pdf(context: dict) -> None:
    styles = pdf_styles()
    doc = SimpleDocTemplate(
        str(PDF_OUTPUT),
        pagesize=A4,
        rightMargin=0.55 * inch,
        leftMargin=0.55 * inch,
        topMargin=0.58 * inch,
        bottomMargin=0.68 * inch,
        title="Tlamatini App Summary",
        author="Tlamatini",
    )
    story: list = []

    cover_image = context["reference_media"][0] if context["reference_media"] else REPO_ROOT / "Tlamatini.jpg"
    story.append(p("TLAMATINI", styles["title"]))
    story.append(
        p(
            "Complete Project Dossier: what the system does, how it works, how to use Tlamatini, complete repository file tree, and effective line inventory",
            styles["subtitle"],
        )
    )
    if cover_image.exists():
        try:
            story.append(cover_image_flowable(cover_image, 6.8 * inch, 3.8 * inch))
            story.append(Spacer(1, 12))
        except Exception:
            pass
    story.append(
        table(
            [
                ["Measure", "Value"],
                ["Generated", context["generated_at"]],
                ["Current HEAD", f"{context['head_short']} - {context['head_subject']}"],
                ["Resolved version", f"{context['version_info']['version']} ({context['version_info']['source']})"],
                ["Repository inventory files", str(context["inventory_files"])],
                ["Tracked files", str(context["tracked_files"])],
                ["Git-unignored working-tree additions", str(context["untracked_files"])],
                ["Workflow agents", str(context["workflow_agent_count"])],
                ["Multi-Turn tools", str(context["total_multi_turn_tools"])],
                ["Skills", str(context["skills_count"])],
                ["Total effective lines", f"{context['total_effective_lines']:,}"],
                ["Total physical text lines", f"{context['total_lines']:,}"],
            ],
            widths=[1.9 * inch, 4.8 * inch],
            font_size=8,
        )
    )
    story.append(PageBreak())

    story.append(p("1. What Tlamatini Is", styles["h1"]))
    for item in SYSTEM_OVERVIEW:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Agent-directory disclaimer: user jurisdiction and responsibility", styles["h2"]))
    for item in AGENT_DIRECTORY_DISCLAIMER:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("What the system does", styles["h2"]))
    for item in WHAT_IT_DOES:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("How it works", styles["h2"]))
    for item in HOW_IT_WORKS:
        story.append(bullet(item, styles["bullet"]))
    story.append(PageBreak())

    story.append(p("2. Architecture Layers", styles["h1"]))
    arch_rows = [["Layer", "Role"]] + [[layer, desc] for layer, desc in ARCHITECTURE_LAYERS]
    story.append(table(arch_rows, widths=[1.75 * inch, 5.0 * inch], font_size=8))
    story.append(p("Design principles", styles["h2"]))
    for item in DESIGN_PRINCIPLES:
        story.append(bullet(item, styles["bullet"]))
    story.append(PageBreak())

    story.append(p("3. Installation, Configuration, and Everyday Use", styles["h1"]))
    story.append(p("Start here - the easiest path", styles["h2"]))
    for item in START_HERE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Installation essentials", styles["h2"]))
    for item in INSTALLATION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("First configuration screens", styles["h2"]))
    for item in FIRST_RUN_CONFIG_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Configuration essentials", styles["h2"]))
    for item in CONFIGURATION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("DB menu and database swap-in", styles["h2"]))
    for item in DB_MENU_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    for item in DB_SWAP_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Versioning system", styles["h2"]))
    for item in VERSIONING_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Version surfaces", styles["h2"]))
    for item in VERSION_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p(f"Current release focus in {context['version_info']['version']}", styles["h2"]))
    for item in CURRENT_RELEASE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.41.4 External-MCP structured output", styles["h2"]))
    for item in STRUCTURED_CONTENT_1414_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.42.0 STM32er PlatformIO expansion", styles["h2"]))
    for item in STM32ER_PLATFORMIO_WORKTREE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.42.0 stepwise STM32 camera-verification demos", styles["h2"]))
    for item in STM32ER_STEPWISE_DEMOS_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.42.0 category-grouped prompt catalog with no gaps", styles["h2"]))
    for item in PROMPT_CATALOG_WORKTREE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.41.3 categorized and deduplicated prompt catalog", styles["h2"]))
    for item in PROMPT_CATALOG_1413_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.41.2 per-user Hard Cancel", styles["h2"]))
    for item in HARD_CANCEL_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.41.0 screenshot paste and image drop", styles["h2"]))
    for item in CHAT_IMAGE_1410_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.40.1 configurable web port", styles["h2"]))
    for item in DJANGO_PORT_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Tlamatini-FlowPills companion discovery", styles["h2"]))
    for item in FLOWPILLS_DISCOVERY_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Unreal Engine 5.8 one-prompt scaffold", styles["h2"]))
    for item in UNREAL_SCAFFOLD_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.39.5 responsiveness and safety hardening", styles["h2"]))
    for item in RESPONSIVENESS_HARDENING_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Nmapper local nmap bridge", styles["h2"]))
    for item in NMAPPER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Startup dialog and prompt-catalog polish", styles["h2"]))
    for item in STARTUP_PROMPT_POLISH_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.38.0 robotic loop closure", styles["h2"]))
    for item in ROBOTIC_LOOP_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.38.1 frontend-state recovery hotfix", styles["h2"]))
    for item in FRONTEND_HOTFIX_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Post-v1.36.0 self-healing Multi-Turn reliability", styles["h2"]))
    for item in SELF_HEALING_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Create Flow and Exec Report gating", styles["h2"]))
    for item in CREATE_FLOW_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Frontend recovery controls and Create Flow name resolution", styles["h2"]))
    for item in FRONTEND_RECOVERY_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Discoverer PDCP key integration", styles["h2"]))
    for item in DISCOVERER_PDCP_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Discoverer vulnx and Go-toolchain Git-deny guard", styles["h2"]))
    for item in DISCOVERER_VULNX_GO_GUARD_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.36.0 Video-Analyzer release delta", styles["h2"]))
    for item in V136_RELEASE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Video-Analyzer motion-verdict agent", styles["h2"]))
    for item in VIDEO_ANALYZER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Prompt search and generated .flw layout", styles["h2"]))
    for item in PROMPT_SEARCH_AND_FLOW_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("v1.33.2 Zavuerer release delta", styles["h2"]))
    for item in V1332_RELEASE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Zavuerer unified-messaging agent", styles["h2"]))
    for item in ZAVUERER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Image-Interpreter triple-model vision pipeline", styles["h2"]))
    for item in IMAGE_INTERPRETER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("External MCP universal client", styles["h2"]))
    for item in EXTERNAL_MCPS_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("MCP Doctor agent and wrapped tool", styles["h2"]))
    for item in MCP_DOCTOR_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("External MCP implementation assets", styles["h2"]))
    for item in EXTERNAL_MCP_ASSETS_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Blenderer", styles["h2"]))
    for item in BLENDERER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("In-app self-update", styles["h2"]))
    for item in SELF_UPDATE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Self-modify source snapshot", styles["h2"]))
    for item in SOURCE_SNAPSHOT_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("API-Keys Wizard", styles["h2"]))
    for item in API_KEYS_WIZARD_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("New assets in the current release wave", styles["h2"]))
    for item in NEW_ASSETS_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("File-Creator hardening", styles["h2"]))
    for item in FILE_CREATOR_HARDENING_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Self-knowledge and identity contract", styles["h2"]))
    for item in SELF_KNOWLEDGE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Self-modify builds", styles["h2"]))
    for item in SELF_MODIFY_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Multi-Turn 4096-turn autonomy", styles["h2"]))
    for item in MULTITURN_4096_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Ask Execs in Multi-Turn", styles["h2"]))
    for item in ASK_EXECS_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Ask Execs runtime path", styles["h2"]))
    for item in ASK_EXECS_PIPELINE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Windows attention issuing", styles["h2"]))
    for item in WINDOWS_ATTENTION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Windows Installed-apps registration", styles["h2"]))
    for item in WINDOWS_APP_REGISTRATION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("De-Compresser agent", styles["h2"]))
    for item in DE_COMPRESSER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("De-Compresser integration and fallback behavior", styles["h2"]))
    for item in DE_COMPRESSER_INTEGRATION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Unreal MCP and the Unrealer agent", styles["h2"]))
    for item in UNREAL_MCP_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Extended Unreal MCP surface", styles["h2"]))
    for item in UNREAL_EXTENDED_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Installing the UE5 plugin and smoke-testing it", styles["h2"]))
    for item in UNREAL_INSTALL_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Orphan-process cleanup", styles["h2"]))
    for item in ORPHAN_REAPER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Orphan-process prevention and survivor reporting", styles["h2"]))
    for item in ORPHAN_PREVENTION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("How to use it", styles["h2"]))
    for item in HOW_TO_USE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Agent descriptions and catalog source of truth", styles["h2"]))
    for item in AGENT_DESCRIPTION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Reviewer and Analyzer", styles["h2"]))
    for item in REVIEWER_ANALYZER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    for item in REVIEWER_ANALYZER_SURFACES:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Operator surface counts", styles["h2"]))
    for item in operator_surface_counts_guide(context):
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Kalier current role", styles["h2"]))
    for item in KALIER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    for item in KALIER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("STM32er current role", styles["h2"]))
    for item in STM32ER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    for item in STM32ER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("ESP32er current role", styles["h2"]))
    for item in ESP32ER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    for item in ESP32ER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("ESP32 Template Project reference baseline", styles["h2"]))
    for item in ESP32_TEMPLATE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("ESPHomer current role", styles["h2"]))
    for item in ESPHOMER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    for item in ESPHOMER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("ESPHome template baseline", styles["h2"]))
    for item in ESPHOME_TEMPLATE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Windower on Multi-Turn and canvas", styles["h2"]))
    for item in WINDOWER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    for item in WINDOWER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Playwrighter current role", styles["h2"]))
    for item in PLAYWRIGHTER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    for item in PLAYWRIGHTER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Reviewer precision patch in v1.4.1", styles["h2"]))
    for item in REVIEWER_PRECISION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Native dialogs and Tkinter removal in v1.4.2", styles["h2"]))
    for item in NATIVE_DIALOGS_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("ACPX-Skills menu", styles["h2"]))
    for item in ACPX_SKILLS_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Source-mode bootstrap commands", styles["h2"]))
    story.append(
        Preformatted(
            "\n".join(
                [
                    "python -m venv venv",
                    "venv\\Scripts\\activate",
                    "pip install -r requirements.txt",
                    "python Tlamatini/manage.py migrate",
                    "python Tlamatini/manage.py createsuperuser",
                    "python Tlamatini/manage.py collectstatic --noinput",
                    "python Tlamatini/manage.py runserver --noreload",
                ]
            ),
            styles["mono"],
        )
    )
    story.append(PageBreak())

    story.append(p("4. Ollama Setup Without Administrative Rights", styles["h1"]))
    for item in OLLAMA_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("No-admin Ollama commands and default model pulls", styles["h2"]))
    story.append(Preformatted(OLLAMA_COMMANDS, styles["mono"]))
    story.append(PageBreak())

    story.append(p("5. Runtime, Release, and Operator Diagnostics", styles["h1"]))
    story.append(p("Running the application", styles["h2"]))
    for item in RUNNING_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Embedding-memory pre-flight guard", styles["h2"]))
    for item in EMBEDDING_GUARD_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Reconnect and restart safeguards", styles["h2"]))
    for item in RECENT_RUNTIME_SAFEGUARDS:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Media and voice family", styles["h2"]))
    for item in MEDIA_VOICE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Autonomous command watchdog", styles["h2"]))
    for item in COMMAND_WATCHDOG_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Prompt catalog and answer readability discipline", styles["h2"]))
    for item in PROMPT_CATALOG_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Ask Execs and execution approval", styles["h2"]))
    for item in ASK_EXECS_GUIDE + ASK_EXECS_PIPELINE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Windows attention issuing", styles["h2"]))
    for item in WINDOWS_ATTENTION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Unreal MCP runtime behavior", styles["h2"]))
    for item in UNREAL_RUNTIME_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Orphan reaper runtime behavior", styles["h2"]))
    for item in ORPHAN_REAPER_GUIDE + ORPHAN_PREVENTION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Release pipeline", styles["h2"]))
    for item in RELEASE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Windows uninstall registration", styles["h2"]))
    for item in WINDOWS_APP_REGISTRATION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Exec Report", styles["h2"]))
    for item in EXEC_REPORT_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("ACPX and skills", styles["h2"]))
    for item in ACPX_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Application log", styles["h2"]))
    for item in APP_LOG_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(PageBreak())

    story.append(p("6. Agent Catalog and Runtime Model", styles["h1"]))
    story.append(p(f"Tlamatini currently exposes {context['workflow_agent_count']} workflow-agent templates.", styles["body"]))
    story.append(table([["Category", "Representative agents"]] + AGENT_CATEGORIES, widths=[1.85 * inch, 4.9 * inch], font_size=7.8))
    story.append(p("All workflow agents follow a common deployment pattern: template directory, YAML configuration, session-scoped pool copy, PID/status/log files, target/source wiring, and optional reanimation state.", styles["body"]))
    story.append(p("Agent catalog validation", styles["h2"]))
    for item in AGENT_DESCRIPTION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("User jurisdiction over plain-Python agents", styles["h2"]))
    for item in AGENT_DIRECTORY_DISCLAIMER:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("How agent runtimes are shaped", styles["h2"]))
    for item in AGENT_RUNTIME_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("MCP Doctor spotlight", styles["h2"]))
    for item in MCP_DOCTOR_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Blenderer spotlight", styles["h2"]))
    for item in BLENDERER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Self-update spotlight", styles["h2"]))
    for item in SELF_UPDATE_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Reviewer and Analyzer spotlight", styles["h2"]))
    for item in REVIEWER_ANALYZER_GUIDE + REVIEWER_ANALYZER_SURFACES:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Operator surface counts", styles["h2"]))
    for item in operator_surface_counts_guide(context):
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Kalier spotlight", styles["h2"]))
    for item in KALIER_GUIDE + KALIER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("STM32er spotlight", styles["h2"]))
    for item in STM32ER_GUIDE + STM32ER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Windower spotlight", styles["h2"]))
    for item in WINDOWER_GUIDE + WINDOWER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Playwrighter spotlight", styles["h2"]))
    for item in PLAYWRIGHTER_GUIDE + PLAYWRIGHTER_SURFACES_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Reviewer precision spotlight", styles["h2"]))
    for item in REVIEWER_PRECISION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Native-dialog spotlight", styles["h2"]))
    for item in NATIVE_DIALOGS_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Unrealer spotlight", styles["h2"]))
    for item in UNREAL_MCP_GUIDE + UNREAL_RUNTIME_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("De-Compresser spotlight", styles["h2"]))
    for item in DE_COMPRESSER_GUIDE + DE_COMPRESSER_INTEGRATION_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(p("Orphan-reaper spotlight", styles["h2"]))
    for item in ORPHAN_REAPER_GUIDE:
        story.append(bullet(item, styles["bullet"]))
    story.append(PageBreak())

    story.append(p("7. Repository Facts and Git Changes", styles["h1"]))
    repo_rows = [
        ["Metric", "Value"],
        ["Repository inventory files", f"{context['inventory_files']}"],
        ["Tracked files in git", f"{context['tracked_files']}"],
        ["Git-unignored working-tree additions", f"{context['untracked_files']}"],
        ["Workflow agents", f"{context['workflow_agent_count']}"],
        ["Multi-Turn tools", f"{context['total_multi_turn_tools']}"],
        ["Wrapped chat-agent tools", f"{context['wrapped_chat_agent_count']}"],
        ["Skills", f"{context['skills_count']}"],
        ["agents_descriptions.md rows", f"{context['agent_description_rows']}"],
        ["Django migrations", f"{context['migrations']}"],
        ["Frontend JavaScript modules", f"{context['js_modules']}"],
        ["Frontend CSS files", f"{context['css_files']}"],
        ["HTML templates", f"{context['html_templates']}"],
        ["Python requirements", f"{context['requirements_count']}"],
        ["Binary/asset inventory files skipped from line count", f"{context['binary_count']}"],
        ["Resolved version", f"{context['version_info']['version']}"],
        ["Version source", f"{context['version_info']['source']}"],
    ]
    story.append(table(repo_rows, widths=[3.0 * inch, 3.7 * inch], font_size=8))
    story.append(p("Latest commits", styles["h2"]))
    commit_rows = [["Date", "Commit", "Subject"]]
    for commit in context["recent_commits"]:
        commit_rows.append([iso_date(commit.committed_at), commit.short_hash, commit.subject])
    story.append(table(commit_rows, widths=[1.0 * inch, 0.8 * inch, 4.9 * inch], font_size=7))
    baseline = context["visual_doc_baseline"]
    if baseline is not None:
        story.append(p("Changes since the last committed PDF/PPTX refresh", styles["h2"]))
        story.append(
            p(
                f"Last committed visual-dossier refresh: {baseline.short_hash} on {iso_date(baseline.committed_at)} — {baseline.subject}",
                styles["body"],
            )
        )
        for item in context["visual_doc_highlights"]:
            story.append(bullet(item, styles["bullet"]))
        visual_chunks = split_items(context["visual_doc_commits"], 12)
        for index, chunk in enumerate(visual_chunks, 1):
            story.append(p(f"Visual-dossier change appendix {index} of {len(visual_chunks)}", styles["h2"]))
            visual_rows = [["Date", "Commit", "Subject"]]
            for commit in chunk:
                visual_rows.append([iso_date(commit.committed_at), commit.short_hash, commit.subject])
            story.append(table(visual_rows, widths=[1.0 * inch, 0.8 * inch, 4.9 * inch], font_size=7))
            if index != len(visual_chunks):
                story.append(PageBreak())
    git_window_heading = "Git changes from today" if RECENT_GIT_WINDOW_LABEL == "today" else f"Git changes from the {RECENT_GIT_WINDOW_LABEL}"
    story.append(p(git_window_heading, styles["h2"]))
    for item in context["weekly_highlights"]:
        story.append(bullet(item, styles["bullet"]))
    weekly_chunks = split_items(context["weekly_commits"], 12)
    for index, chunk in enumerate(weekly_chunks, 1):
        appendix_heading = (
            f"Today's commit appendix {index} of {len(weekly_chunks)}"
            if RECENT_GIT_WINDOW_LABEL == "today"
            else f"{RECENT_GIT_WINDOW_DAYS}-day commit appendix {index} of {len(weekly_chunks)}"
        )
        story.append(p(appendix_heading, styles["h2"]))
        weekly_rows = [["Date", "Commit", "Subject"]]
        for commit in chunk:
            weekly_rows.append([iso_date(commit.committed_at), commit.short_hash, commit.subject])
        story.append(table(weekly_rows, widths=[1.0 * inch, 0.8 * inch, 4.9 * inch], font_size=7))
        if index != len(weekly_chunks):
            story.append(PageBreak())
    story.append(PageBreak())

    story.append(p("8. Effective Line Inventory by Language", styles["h1"]))
    story.append(
        p(
            "Methodology: tracked text files only. Blank lines and comment-only lines are excluded. Python counts also remove module, class, and function docstrings detected through AST parsing.",
            styles["body"],
        )
    )
    line_rows = [["Language", "Files", "Effective lines", "Total lines", "Share"]]
    for row in context["language_rows"]:
        share = row.effective_lines / max(context["total_effective_lines"], 1)
        line_rows.append(
            [
                row.language,
                f"{row.files}",
                f"{row.effective_lines:,}",
                f"{row.total_lines:,}",
                f"{share:.1%}",
            ]
        )
    story.append(table(line_rows, widths=[1.65 * inch, 0.75 * inch, 1.35 * inch, 1.25 * inch, 0.85 * inch], font_size=7.5))
    story.append(p(f"Total effective lines: {context['total_effective_lines']:,}", styles["h2"]))
    story.append(PageBreak())

    story.append(p("9. Largest Effective Source Files", styles["h1"]))
    largest_rows = [["Path", "Language", "Effective", "Total"]]
    for file_stat in context["file_rows"][:25]:
        largest_rows.append([file_stat.path, file_stat.language, f"{file_stat.effective_lines:,}", f"{file_stat.total_lines:,}"])
    story.append(table(largest_rows, widths=[4.0 * inch, 1.2 * inch, 0.75 * inch, 0.75 * inch], font_size=6.7))
    story.append(PageBreak())

    story.append(p("10. Complete Repository File Tree (Repository Appendix)", styles["h1"]))
    TREE_OUTPUT.write_text(context["tree_text"], encoding="utf-8")
    tree_chunks = split_lines(context["tree_text"], 76)
    for index, chunk in enumerate(tree_chunks, 1):
        story.append(p(f"Tree appendix {index} of {len(tree_chunks)}", styles["h2"]))
        story.append(Preformatted(chunk, styles["mono"]))
        if index != len(tree_chunks):
            story.append(PageBreak())

    doc.build(story, onFirstPage=pdf_page_footer, onLaterPages=pdf_page_footer)


def rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.color.rgb = color


def add_dark_background(slide, accent: RGBColor, image_path: Path | None = None) -> None:
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))
        veil = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))
        fill(veil, THEME["obsidian"], 22)
        veil.line.fill.background()
    else:
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))
        fill(bg, THEME["obsidian"])
        bg.line.fill.background()

    for x in [0.35, SLIDE_W - 0.52]:
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(0.42), Inches(0.04), Inches(6.65))
        fill(line, accent, 25)
        line.line.fill.background()
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.42), Inches(12.2), Inches(0.03))
    fill(top, accent, 15)
    top.line.fill.background()
    bottom = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.04), Inches(12.2), Inches(0.03))
    fill(bottom, accent, 40)
    bottom.line.fill.background()


def add_run(paragraph, text: str, size: int, color: RGBColor, bold: bool = False, font: str = "Aptos") -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    audit: list[tuple[float, float, float, float, str]],
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 18,
    color: RGBColor | None = None,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    name: str = "text",
    font: str = "Aptos",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p0 = frame.paragraphs[0]
    p0.alignment = align
    add_run(p0, text, size, color or THEME["white"], bold, font)
    audit.append((x, y, w, h, name))


def fit_bullet_size(
    bullets: list[str],
    box_w_in: float,
    box_h_in: float,
    max_size: float,
    min_size: float = 9.0,
    space_after_pt: float = 4.0,
) -> int:
    """Deterministically pick the LARGEST font size in [min_size, max_size] at
    which the bullet block is estimated to fit inside ``box_w_in`` x ``box_h_in``.

    This does NOT rely on ``TextFrame.fit_text`` — that path silently fails when
    PowerPoint cannot resolve the display font (Aptos) on the build host, which is
    exactly why dense cards used to overflow. The estimate is intentionally
    CONSERVATIVE (slightly over-counts wrapped lines and reserves a vertical
    gutter) so rendered text never spills its card. ``space_after_pt`` must match
    the paragraph spacing used by :func:`add_bullets`."""
    usable_w_pt = max(box_w_in * 72.0 - 30.0, 36.0)   # minus bullet glyph + indent + margins
    budget_h_pt = box_h_in * 72.0 * 0.93              # keep a safety gutter
    lo, hi = int(round(min_size)), int(round(max_size))
    for size in range(hi, lo - 1, -1):
        chars_per_line = max(int(usable_w_pt / (0.52 * size)), 6)
        line_h = 1.26 * size
        total = 0.0
        for text in bullets:
            n_lines = max(1, -(-len(str(text)) // chars_per_line))  # ceil
            total += n_lines * line_h + space_after_pt
        if total <= budget_h_pt:
            return size
    return lo


def add_bullets(
    slide,
    audit: list[tuple[float, float, float, float, str]],
    x: float,
    y: float,
    w: float,
    h: float,
    bullets: list[str],
    size: int = 16,
    color: RGBColor | None = None,
    name: str = "bullets",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    # Backstop only: PowerPoint may re-shrink on edit. The deterministic size
    # computed below is the primary guarantee against overflow.
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Pt(3)
    frame.margin_right = Pt(3)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    space_after_pt = 4.0
    fitted = fit_bullet_size(bullets, w, h, size, space_after_pt=space_after_pt)
    for idx, item in enumerate(bullets):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.text = item
        para.level = 0
        para.font.name = "Aptos"
        para.font.size = Pt(fitted)
        para.font.color.rgb = color or THEME["muted"]
        para.space_after = Pt(space_after_pt)
        para.bullet = True
    audit.append((x, y, w, h, name))


def add_title(slide, audit: list[tuple[float, float, float, float, str]], title: str, kicker: str, accent: RGBColor) -> None:
    add_text(slide, audit, 0.78, 0.6, 11.7, 0.33, kicker.upper(), 10, accent, True, name="kicker", font="Aptos")
    add_text(slide, audit, 0.76, 0.93, 11.9, 0.62, title, 27, THEME["white"], False, name="title", font="Aptos Display")


def add_panel(
    slide,
    audit: list[tuple[float, float, float, float, str]],
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    bullets: list[str],
    accent: RGBColor,
    name: str,
    size: int = 15,
) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, THEME["panel"], 8)
    shape.line.color.rgb = THEME["line"]
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.12))
    fill(bar, accent, 0)
    bar.line.fill.background()
    add_text(slide, audit, x + 0.22, y + 0.22, w - 0.44, 0.34, title, 14, accent, True, name=f"{name}-title")
    add_bullets(slide, audit, x + 0.22, y + 0.72, w - 0.44, h - 0.88, bullets, size=size, name=f"{name}-body")
    audit.append((x, y, w, h, name))


def add_metric_card(
    slide,
    audit: list[tuple[float, float, float, float, str]],
    x: float,
    y: float,
    w: float,
    label: str,
    value: str,
    accent: RGBColor,
    name: str,
) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.02))
    fill(shape, THEME["panel2"], 4)
    shape.line.color.rgb = accent
    add_text(slide, audit, x + 0.16, y + 0.15, w - 0.32, 0.24, label.upper(), 8, accent, True, name=f"{name}-label")
    add_text(slide, audit, x + 0.16, y + 0.42, w - 0.32, 0.38, value, 20, THEME["white"], True, name=f"{name}-value")
    audit.append((x, y, w, 1.02, name))


def add_flow_boxes(
    slide,
    audit: list[tuple[float, float, float, float, str]],
    x: float,
    y: float,
    labels: list[str],
    accent: RGBColor,
) -> None:
    box_w = 1.74
    gap = 0.27
    for idx, label in enumerate(labels):
        bx = x + idx * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(bx), Inches(y), Inches(box_w), Inches(0.75))
        fill(shape, THEME["stone"], 3)
        shape.line.color.rgb = accent
        add_text(slide, audit, bx + 0.08, y + 0.18, box_w - 0.16, 0.28, label, 10, THEME["white"], True, PP_ALIGN.CENTER, name=f"flow-{idx}")
        audit.append((bx, y, box_w, 0.75, f"flowbox-{idx}"))
        if idx < len(labels) - 1:
            add_text(slide, audit, bx + box_w, y + 0.19, gap, 0.2, ">", 12, accent, True, PP_ALIGN.CENTER, name=f"arrow-{idx}")


def audit_layout(audit: list[tuple[float, float, float, float, str]], slide_no: int) -> None:
    for x, y, w, h, name in audit:
        if x < -0.01 or y < -0.01 or x + w > SLIDE_W + 0.01 or y + h > SLIDE_H + 0.01:
            raise RuntimeError(f"Slide {slide_no}: {name} is outside slide bounds")
    major = [
        item
        for item in audit
        if not item[4].endswith("-label")
        and not item[4].endswith("-value")
        and not item[4].startswith("flow-")
        and not item[4].startswith("arrow-")
    ]
    for i, a in enumerate(major):
        for b in major[i + 1 :]:
            if a[4].split("-")[0] == b[4].split("-")[0]:
                continue
            if rects_overlap(a, b):
                raise RuntimeError(f"Slide {slide_no}: {a[4]} overlaps {b[4]}")


def rects_overlap(a: tuple[float, float, float, float, str], b: tuple[float, float, float, float, str]) -> bool:
    ax, ay, aw, ah, _ = a
    bx, by, bw, bh, _ = b
    return ax < bx + bw and ax + aw > bx and ay < by + bh and ay + ah > by


def add_slide(prs: Presentation, title: str, kicker: str, accent: RGBColor, image: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    audit: list[tuple[float, float, float, float, str]] = []
    add_dark_background(slide, accent, image)
    add_title(slide, audit, title, kicker, accent)
    return slide, audit


def add_themed_column_slides(
    prs: Presentation,
    title: str,
    kicker: str,
    accent: RGBColor,
    columns: list[tuple[str, RGBColor, list[str]]],
    per_column: int = 5,
    top: float = 1.58,
    height: float = 5.0,
    size: int = 14,
) -> None:
    """Render labelled columns of bullets, PAGINATING onto as many slides as
    needed so no column ever holds more than ``per_column`` bullets — the
    "split into more slides" rule that keeps dense content from overflowing its
    cards. Each ``columns`` entry is ``(header, accent, items)``; a column's
    items beyond ``per_column`` continue, in the same lane, on the next page."""
    margin, gap = 0.72, 0.3
    n = len(columns)
    col_w = (SLIDE_W - 2 * margin - (n - 1) * gap) / n
    max_items = max((len(items) for _, _, items in columns), default=0)
    pages = max(1, -(-max_items // per_column))  # ceil
    for page in range(pages):
        page_title = title if pages == 1 else f"{title} ({page + 1}/{pages})"
        slide, audit = add_slide(prs, page_title, kicker, accent)
        for ci, (header, col_accent, items) in enumerate(columns):
            seg = items[page * per_column:(page + 1) * per_column]
            x = margin + ci * (col_w + gap)
            add_panel(slide, audit, x, top, col_w, height, header, seg, col_accent, f"col-{page}-{ci}", size)
        audit_layout(audit, len(prs.slides))


def language_table_text(rows: list[LineStats], total_effective: int) -> str:
    lines = ["LANGUAGE                 FILES   EFFECTIVE      TOTAL   SHARE"]
    for row in rows:
        share = row.effective_lines / max(total_effective, 1)
        lines.append(
            f"{row.language[:22]:22} {row.files:5d} {row.effective_lines:11,d} {row.total_lines:10,d} {share:6.1%}"
        )
    return "\n".join(lines)


def file_table_text(rows: list[FileStats], limit: int = 14) -> str:
    lines = ["PATH                                                     LANGUAGE     EFFECTIVE   TOTAL"]
    for row in rows[:limit]:
        path = row.path
        if len(path) > 55:
            path = "..." + path[-52:]
        lines.append(f"{path:55} {row.language[:10]:10} {row.effective_lines:9,d} {row.total_lines:7,d}")
    return "\n".join(lines)


def build_ppt(context: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    cover = context["reference_media"][0] if context["reference_media"] else None

    slide, audit = add_slide(
        prs,
        "TLAMATINI",
        "Complete project dossier: what she does, how she works, how to use Tlamatini",
        THEME["copper"],
        cover,
    )
    add_text(slide, audit, 0.9, 2.0, 5.5, 0.6, "El Saber Cosmico del Desarrollo", 24, THEME["white"], False, name="cover-tag", font="Aptos Display")
    add_text(
        slide,
        audit,
        0.9,
        2.76,
        6.3,
        1.2,
        f"Self-hosted AI developer assistant (cloud LLMs by default) with RAG, Multi-Turn orchestration, {context['workflow_agent_count']} agents, ACPX delegation, visual workflows, self-knowledge, and Windows packaging.",
        17,
        THEME["muted"],
        False,
        name="cover-body",
    )
    add_metric_card(slide, audit, 0.9, 4.25, 1.75, "Files", str(context["inventory_files"]), THEME["jade"], "cover-m1")
    add_metric_card(slide, audit, 2.85, 4.25, 1.75, "Agents", str(context["workflow_agent_count"]), THEME["copper"], "cover-m2")
    add_metric_card(slide, audit, 4.8, 4.25, 1.95, "Effective", f"{context['total_effective_lines']:,}", THEME["amber"], "cover-m3")
    add_text(slide, audit, 0.9, 6.35, 6.2, 0.32, f"Generated {context['generated_at']} at HEAD {context['head_short']}", 9, THEME["muted"], name="cover-foot")
    audit_layout(audit, 1)

    slide, audit = add_slide(prs, "What Tlamatini Is", "system identity", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.65, 5.85, 4.95, "Definition", SYSTEM_OVERVIEW, THEME["jade"], "identity-a", 16)
    add_panel(slide, audit, 6.92, 1.65, 5.55, 4.95, "Core interfaces", [
        "Chat page at `/agent/` for context-grounded Q&A, code generation, tool calls, and Multi-Turn execution.",
        "Agentic Control Panel at `/agentic_control_panel/` for visual workflow design.",
        "Django admin and config dialogs for MCPs, tools, agents, users, and persistent settings.",
    ], THEME["copper"], "identity-b", 16)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Agent Directory Disclaimer", "plain-Python agents are user-jurisdiction code", THEME["amber"])
    add_panel(
        slide,
        audit,
        0.78,
        1.6,
        5.9,
        4.95,
        "User jurisdiction",
        AGENT_DIRECTORY_DISCLAIMER[:2],
        THEME["amber"],
        "agent-disclaimer-a",
        12,
    )
    add_panel(
        slide,
        audit,
        6.95,
        1.6,
        5.55,
        4.95,
        "Responsibility boundary",
        AGENT_DIRECTORY_DISCLAIMER[2:],
        THEME["copper"],
        "agent-disclaimer-b",
        12,
    )
    audit_layout(audit, len(prs.slides))

    add_themed_column_slides(prs, "What The System Does", "capability map", THEME["copper"], [
        ("Knowledge", THEME["jade"], WHAT_IT_DOES[:6]),
        ("Action", THEME["copper"], WHAT_IT_DOES[6:12]),
        ("Delivery", THEME["amber"], WHAT_IT_DOES[12:]),
    ], per_column=3)

    slide, audit = add_slide(prs, "How It Works", "execution pipeline", THEME["jade"])
    add_flow_boxes(slide, audit, 0.95, 2.6, ["Browser", "Channels", "RAG", "Planner", "Tools", "Answer"], THEME["jade"])
    add_panel(slide, audit, 0.82, 4.0, 11.55, 2.45, "Request-to-answer flow", [
        "A browser request flows through Django Channels into the RAG/context layer, the Multi-Turn planner, the tool executor, and back as a synthesized answer.",
        "The next two pages detail the request path (intake, retrieval, permission gating) and the runtime path (planning, tool execution, agent bridges, packaging).",
    ], THEME["jade"], "works-intro", 15)
    audit_layout(audit, len(prs.slides))

    add_themed_column_slides(prs, "How It Works — Detail", "execution pipeline", THEME["jade"], [
        ("Request path", THEME["jade"], HOW_IT_WORKS[:11]),
        ("Runtime path", THEME["copper"], HOW_IT_WORKS[11:]),
    ], per_column=6)

    slide, audit = add_slide(prs, "Design Principles", "how the software is shaped", THEME["amber"])
    add_panel(slide, audit, 0.82, 1.72, 11.55, 4.75, "Core design choices", DESIGN_PRINCIPLES, THEME["amber"], "design", 16)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "RAG And Context Engine", "retrieval core", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.8, 4.95, "RAG responsibilities", [
        "Load selected file or directory context from the chat interface.",
        "Extract metadata, split text, rank chunks with hybrid retrieval, and respect context budgets.",
        "Fallback to explicit loaded files when retrieval cannot provide enough memory.",
        "Preserve the difference between direct one-shot chat and checked Multi-Turn orchestration.",
    ], THEME["amber"], "rag-a", 15)
    add_panel(slide, audit, 6.9, 1.6, 5.6, 4.95, "Why it matters", [
        "The assistant answers from project evidence instead of generic memory.",
        "Large codebases remain navigable without injecting the whole repository into every prompt.",
        "Multi-Turn can prefetch only the contexts required by the current execution plan.",
    ], THEME["jade"], "rag-b", 16)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Multi-Turn Oracle", "agentic execution", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Checked mode", [
        "Capability registry scores context providers, tools, and wrapped agents for the current request.",
        "Global planner builds prefetch, execute, monitor, and answer stages.",
        "Explicit tool loop runs tool calls, appends observations, and asks again until final answer or the 4096-turn limit.",
        "Create Flow no longer uses the removed answer classifier; it appears when Multi-Turn has at least one successful agent call.",
    ], THEME["copper"], "mt-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Unchecked mode", [
        "Keeps the original prompt validation and legacy prefetch behavior.",
        "Maintains compatibility for fast Q&A and simple context-grounded answers.",
        "Avoids forcing every chat request into agentic execution.",
    ], THEME["jade"], "mt-b", 16)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Ask Execs", "v1.10.0 safety modifier still active in v1.42.0", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Operator contract", ASK_EXECS_GUIDE, THEME["amber"], "ask-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Runtime mechanics", ASK_EXECS_PIPELINE_GUIDE, THEME["jade"], "ask-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Windows Attention Issuing", "taskbar flash and uppercase log banner when she needs you", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Current mechanism", WINDOWS_ATTENTION_GUIDE, THEME["jade"], "attention-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Why operators care", [
        "If the browser is buried, Tlamatini still has a local way to pull the operator back at the exact moment an approval prompt or notification matters.",
        "The path is concrete and auditable: `POST /agent/flash_window/`, `window_flash.py`, `FlashWindowEx`, and the matching uppercase banner in `tlamatini.log`.",
        "Because the helper is best-effort and fail-safe, a missed flash never breaks the request path or blocks the surrounding workflow.",
    ], THEME["amber"], "attention-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Windows Installed-App Registration", "v1.11.0 uninstall integration carried into v1.42.0", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What changed", WINDOWS_APP_REGISTRATION_GUIDE, THEME["copper"], "arp-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Why operators care", [
        "Packaged installs now show up in normal Windows uninstall surfaces instead of only leaving behind shortcuts and a loose `Uninstaller.exe` in the install folder.",
        "The registration is HKCU-only and non-elevated, matching the installer’s per-user design on Windows 10 and Windows 11.",
        "Because frozen startup self-heals the entry, even older installs can gain the uninstall surface after a later app launch without a reinstall.",
    ], THEME["jade"], "arp-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Current Release Focus", "v1.42.0 - STM32er expansion plus carried reliability", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Release line", CURRENT_RELEASE_GUIDE[:3], THEME["amber"], "rel-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Reliability and flow", CURRENT_RELEASE_GUIDE[3:6], THEME["jade"], "rel-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Current Release Context", "image paths, port recovery, handbook truth, and carried foundations", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Flow and recovery", CURRENT_RELEASE_GUIDE[6:10], THEME["jade"], "rel-c", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Security and source control", CURRENT_RELEASE_GUIDE[10:12], THEME["amber"], "rel-d", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Release Continuity", "older waves still carried by the current dossier", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Carried product story", CURRENT_RELEASE_GUIDE[12:14], THEME["copper"], "rel-e", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Documentation contract", CURRENT_RELEASE_GUIDE[14:], THEME["jade"], "rel-f", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "External MCP Structured Results", "v1.41.4 - successful server data reaches the model", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Failure mechanism", STRUCTURED_CONTENT_1414_GUIDE[:3], THEME["jade"], "mcp-1414-a", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Unified transport path", STRUCTURED_CONTENT_1414_GUIDE[3:5], THEME["amber"], "mcp-1414-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Structured Output Boundaries", "context protection, errors, and regression proof", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Payload safety", STRUCTURED_CONTENT_1414_GUIDE[5:6], THEME["copper"], "mcp-1414-c", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Coverage and release", STRUCTURED_CONTENT_1414_GUIDE[6:], THEME["jade"], "mcp-1414-d", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "STM32er PlatformIO Release", "published in tagged v1.42.0", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Routing and coverage", STM32ER_PLATFORMIO_WORKTREE_GUIDE[:4], THEME["amber"], "stm32-pio-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Preflight and safe composite", STM32ER_PLATFORMIO_WORKTREE_GUIDE[4:6], THEME["jade"], "stm32-pio-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "PlatformIO Coverage Boundary", "what STM32er Phase 1 implements and what remains planned", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Deliberate refusal", STM32ER_PLATFORMIO_WORKTREE_GUIDE[6:8], THEME["copper"], "stm32-pio-c", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "New assets and tests", STM32ER_PLATFORMIO_WORKTREE_GUIDE[8:], THEME["jade"], "stm32-pio-d", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "STM32 Stepwise Proof Demos", "Blue Pill and F407 Discovery from driver to camera evidence", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Operator contract", STM32ER_STEPWISE_DEMOS_GUIDE[:3], THEME["jade"], "stm32-step-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Hardware and evidence", STM32ER_STEPWISE_DEMOS_GUIDE[3:], THEME["amber"], "stm32-step-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Prompt Catalog No-Gap Renumber", "v1.42.0 - deliberate one-time primary-key migration", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Migration mechanics", PROMPT_CATALOG_WORKTREE_GUIDE[:3], THEME["copper"], "prompt-live-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Boundaries and proof", PROMPT_CATALOG_WORKTREE_GUIDE[3:], THEME["jade"], "prompt-live-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Prompt Catalog Reorganized", "v1.41.3 - 13 categories and duplicate removal", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Migration contract", PROMPT_CATALOG_1413_GUIDE[:4], THEME["jade"], "prompt-1413-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Stable identity", PROMPT_CATALOG_1413_GUIDE[4:5], THEME["amber"], "prompt-1413-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Prompt Catalog Search", "grouped at rest, ranked and flattened while searching", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Search behavior", PROMPT_CATALOG_1413_GUIDE[5:7], THEME["amber"], "prompt-1413-c", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Viewport contract", PROMPT_CATALOG_1413_GUIDE[7:], THEME["jade"], "prompt-1413-d", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Hard Cancel Run Latch", "v1.41.2 - cancelled runs stay cancelled", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Per-user epoch", HARD_CANCEL_GUIDE[:4], THEME["copper"], "hard-cancel-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Evidence preservation", HARD_CANCEL_GUIDE[4:5], THEME["jade"], "hard-cancel-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Hard Cancel Boundaries", "approval, recovery, frontend, and regression coverage", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "No resurrection", HARD_CANCEL_GUIDE[5:7], THEME["jade"], "hard-cancel-c", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Coverage", HARD_CANCEL_GUIDE[7:], THEME["amber"], "hard-cancel-d", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Screenshot To Chat", "v1.41.0 - paste or drop an image path into the prompt", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Ingestion path", CHAT_IMAGE_1410_GUIDE[:4], THEME["amber"], "image-chat-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Reversible input", CHAT_IMAGE_1410_GUIDE[4:5], THEME["jade"], "image-chat-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Image Ingestion Boundaries", "vision handoff, layout safety, and implementation assets", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Vision and layout", CHAT_IMAGE_1410_GUIDE[5:6], THEME["copper"], "image-chat-c", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Concrete assets", CHAT_IMAGE_1410_GUIDE[6:], THEME["jade"], "image-chat-d", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Configurable Web Port", "v1.40.1 — escape reserved port 8000 without rebuilding", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Why and where", DJANGO_PORT_GUIDE[:3], THEME["amber"], "django-port-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Every manage.py launch path", DJANGO_PORT_GUIDE[3:5], THEME["jade"], "django-port-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Web Port Boundaries", "fail-open resolution, CLI precedence, and focused coverage", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Deliberate boundaries", DJANGO_PORT_GUIDE[5:6], THEME["copper"], "django-port-c", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Regression contract", DJANGO_PORT_GUIDE[6:], THEME["jade"], "django-port-d", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "FlowPills Companion Discovery", "v1.40.0 — find agent templates without Python or drive scans", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Lookup contract", FLOWPILLS_DISCOVERY_GUIDE[:3], THEME["jade"], "flowpills-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Publication and preservation", FLOWPILLS_DISCOVERY_GUIDE[3:], THEME["amber"], "flowpills-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Unreal 5.8 Project Scaffold", "v1.39.5 — two prompt fields to a ready-to-build C++ project", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Deterministic scaffold", UNREAL_SCAFFOLD_GUIDE[:3], THEME["copper"], "unreal-scaffold-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Build and editor handoff", UNREAL_SCAFFOLD_GUIDE[3:], THEME["jade"], "unreal-scaffold-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "v1.39.5 Runtime Hardening", "bounded waits, accurate partial results, and request isolation", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Responsiveness", RESPONSIVENESS_HARDENING_GUIDE[:3], THEME["amber"], "responsive-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Safety and result integrity", RESPONSIVENESS_HARDENING_GUIDE[3:], THEME["jade"], "responsive-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Nmapper Local Recon", "use-only nmap bridge for authorized targets", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Operator contract", NMAPPER_GUIDE[:3], THEME["jade"], "nmapper-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Actions and assets", NMAPPER_GUIDE[3:], THEME["amber"], "nmapper-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Startup And Prompt Polish", "v1.39.4 closeability plus current prompt localization", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Startup dialog", STARTUP_PROMPT_POLISH_GUIDE[:2], THEME["amber"], "startup-prompt-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Prompt catalog", STARTUP_PROMPT_POLISH_GUIDE[2:], THEME["jade"], "startup-prompt-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Robotic Loop Closed", "v1.38.0 — from blank page to observed hardware behavior", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Closed-loop chain", ROBOTIC_LOOP_GUIDE[:2], THEME["copper"], "robot-loop-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Verdict safety", ROBOTIC_LOOP_GUIDE[2:], THEME["jade"], "robot-loop-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "v1.38.1 Frontend Hotfix", "mutable-state recovery and one-call prompt catalog", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Const-poison recovery", FRONTEND_HOTFIX_GUIDE[:3], THEME["amber"], "frontend-hotfix-a", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Prompt catalog path", FRONTEND_HOTFIX_GUIDE[3:], THEME["jade"], "frontend-hotfix-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Self-Healing Multi-Turn", "watchdog-bounded model steps and truthful recovery", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Recovery loop", SELF_HEALING_GUIDE[:3], THEME["copper"], "heal-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Evidence preservation", SELF_HEALING_GUIDE[3:5], THEME["jade"], "heal-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Create Flow Gate", "successful-only flows without answer_success", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What changed", CREATE_FLOW_GUIDE[:2], THEME["jade"], "flow-gate-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Operator result", CREATE_FLOW_GUIDE[2:], THEME["amber"], "flow-gate-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Frontend Recovery Controls", "status frames stay busy until the real final answer", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Self-healing status frames", FRONTEND_RECOVERY_GUIDE[:3], THEME["copper"], "frontend-recovery-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Create Flow name resolution", FRONTEND_RECOVERY_GUIDE[3:], THEME["jade"], "frontend-recovery-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Discoverer PDCP Wiring", "ProjectDiscovery key without prompt-pasted secrets", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Configuration path", DISCOVERER_PDCP_GUIDE[:3], THEME["amber"], "pdcp-a", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Redaction and prompt seed", DISCOVERER_PDCP_GUIDE[3:], THEME["copper"], "pdcp-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Discoverer Vulnx And Go Guard", "current CVE search plus source-control protection", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "vulnx CVE lane", DISCOVERER_VULNX_GO_GUARD_GUIDE[:3], THEME["copper"], "vulnx-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Go-deny guard", DISCOVERER_VULNX_GO_GUARD_GUIDE[3:], THEME["jade"], "vulnx-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "v1.36.0 Release Delta", "Video-Analyzer, prompts, and generated workflow layout", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Exact release checklist", V136_RELEASE_GUIDE[:3], THEME["copper"], "v136-a", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Runtime contract", V136_RELEASE_GUIDE[3:], THEME["jade"], "v136-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Video-Analyzer", "motion-verdict agent for robotic-loop training", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What she does", VIDEO_ANALYZER_GUIDE[:3], THEME["jade"], "video-analyzer-a", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Verdicts and routing", VIDEO_ANALYZER_GUIDE[3:], THEME["amber"], "video-analyzer-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Prompt Search And .flw Layout", "operator usability after the catalog grew", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Prompt search", PROMPT_SEARCH_AND_FLOW_GUIDE[:2], THEME["amber"], "prompt-flow-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Generated flows", PROMPT_SEARCH_AND_FLOW_GUIDE[2:], THEME["copper"], "prompt-flow-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "v1.33.2 Release Delta", "Zavuerer, model defaults, cleanup, and safety", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Exact release checklist", V1332_RELEASE_GUIDE, THEME["copper"], "v1332-a", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Recent implementation assets", [
        "`agent/agents/zavuerer/` carries the stdlib-only Zavu REST client and config template.",
        "`0159`-`0164` seed the Agent, Tool, prompts, and setup-wizard dedupe.",
        "`access_key_wizard.py`, `tools.py`, and `config.json` wire the Zavu key path.",
        "`config.json` now favors `glm-5.2:cloud` for the shipped cloud model baseline.",
        "`views.py`, ACP JS/CSS, and `agent_contracts.py` make canvas wiring and redaction work.",
    ], THEME["jade"], "v1332-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Zavuerer", "new unified messaging agent for Zavu", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What she does", ZAVUERER_GUIDE[:3], THEME["amber"], "zavu-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Implementation and safety", ZAVUERER_GUIDE[3:], THEME["jade"], "zavu-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Image-Interpreter", "triple-model vision analysis with fail-safe merging", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Parallel vision pipeline", IMAGE_INTERPRETER_GUIDE[:4], THEME["copper"], "image-interpreter-a", 11)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Prompts, output, and config", IMAGE_INTERPRETER_GUIDE[4:], THEME["jade"], "image-interpreter-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "External MCPs", "how Tlamatini now reaches tools outside her bundled runtime", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Operator model", EXTERNAL_MCPS_GUIDE, THEME["jade"], "xmcp-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Concrete implementation", EXTERNAL_MCP_ASSETS_GUIDE, THEME["amber"], "xmcp-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "MCP Doctor", "safe onboarding and diagnostics before live external MCP use", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What she does", MCP_DOCTOR_GUIDE, THEME["copper"], "doctor-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "How operators use it", [
        "Open `External -> MCPs` to add or import a server, then run MCP Doctor before the first real connect if the runtime, endpoint, or secret story is still uncertain.",
        "In Multi-Turn, call `chat_agent_mcp_doctor` when you want the LLM to triage a server declaratively instead of guessing from prose about PATH, env vars, or transport settings.",
        "Because the diagnosis path is static and fail-safe, it surfaces onboarding mistakes early without leaking into a half-connected remote-tool session.",
    ], THEME["jade"], "doctor-b", 11)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Start Here", "the new easy-follow onboarding path", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Five-step path", START_HERE_GUIDE, THEME["amber"], "start-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "First configuration surfaces", FIRST_RUN_CONFIG_GUIDE, THEME["jade"], "start-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Blenderer And Self-Update", "how to use the newest release surfaces", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Blenderer", BLENDERER_GUIDE, THEME["jade"], "blend-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Self-update", SELF_UPDATE_GUIDE, THEME["amber"], "blend-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Self-Modify Source Snapshot", "generated `TlamatiniSourceCode/` and rebuild contract", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Snapshot contract", SOURCE_SNAPSHOT_GUIDE, THEME["jade"], "snap-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Access Keys, Updates, And File Writes", API_KEYS_WIZARD_GUIDE[:2] + SELF_UPDATE_GUIDE[:1] + FILE_CREATOR_HARDENING_GUIDE[:1], THEME["amber"], "snap-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Media, Voice, And Runtime Resilience", "Talker / Whisperer family plus watchdog hardening", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Media and voice", MEDIA_VOICE_GUIDE, THEME["copper"], "media-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Autonomous watchdog", COMMAND_WATCHDOG_GUIDE, THEME["jade"], "media-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Newest Assets And Surfaces", "backend, frontend, and build files added or upgraded recently", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Named assets", NEW_ASSETS_GUIDE[:3], THEME["jade"], "assets-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Hardening assets", NEW_ASSETS_GUIDE[3:6], THEME["amber"], "assets-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Newest Assets Appendix", "remaining current asset deltas", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Discoverer and frontend", NEW_ASSETS_GUIDE[6:9], THEME["amber"], "assets-c", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Frontend, port, and agent assets", NEW_ASSETS_GUIDE[9:12], THEME["jade"], "assets-d", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Newest Assets Source Trail", "why the inventory and line counts changed", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Media and browser setup", NEW_ASSETS_GUIDE[12:15], THEME["copper"], "assets-e", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Inventory meaning", NEW_ASSETS_GUIDE[15:], THEME["jade"], "assets-f", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Agentic Control Panel", "visual workflow temple", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "ACP workflow", [
        "Drag prebuilt agents onto the canvas and connect them visually.",
        "Configure deployed pool instances, not only static template folders.",
        "Validate structural rules before execution.",
        "Start through Starter, pause/resume through reanimation state, and stop through Ender semantics.",
    ], THEME["jade"], "acp-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Flow files", [
        "Workflows save and load as `.flw` files.",
        "Generated flows from chat are starter drafts, not a replacement for ACP validation.",
        "LED indicators and logs show runtime health across the session pool.",
    ], THEME["copper"], "acp-b", 16)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, f"The {context['workflow_agent_count']} Guardians", "workflow agent catalog", THEME["copper"])
    left = [f"{name}: {desc}" for name, desc in AGENT_CATEGORIES[:4]]
    right = [f"{name}: {desc}" for name, desc in AGENT_CATEGORIES[4:]]
    add_panel(slide, audit, 0.72, 1.56, 5.95, 5.1, "Agent families", left, THEME["copper"], "agents-a", 13)
    add_panel(slide, audit, 6.92, 1.56, 5.65, 5.1, "More guardians", right, THEME["jade"], "agents-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Agent Catalog Integrity", "count and description source of truth", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Count alignment", [
        f"Templates on disk with config.yaml: {context['workflow_agent_count']}",
        f"Description rows in `agents_descriptions.md`: {context['agent_description_rows']}",
        f"This dossier reconciles README and Book wording back to the same {context['workflow_agent_count']}-agent inventory when older badges or legacy prose lines lag behind the live tree.",
    ], THEME["jade"], "agent-proof-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Description source", AGENT_DESCRIPTION_GUIDE, THEME["copper"], "agent-proof-b", 14)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Reviewer And Analyzer", "new review and security surfaces", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What they do", REVIEWER_ANALYZER_GUIDE, THEME["amber"], "reviewer-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "How operators reach them", REVIEWER_ANALYZER_SURFACES, THEME["jade"], "reviewer-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Operator Surface Counts", "README header and planner-facing inventory", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Current counts", operator_surface_counts_guide(context), THEME["copper"], "surface-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Why the counts matter", [
        "The README now surfaces the same operator picture the app exposes in practice: broad capability, selective planner binding, and a capped tool budget per request.",
        f"Those counts complement the {context['workflow_agent_count']}-agent bestiary instead of replacing it: skills, wrapped tools, and ACPX tools are different layers of the same operating surface.",
        "For dossier readers, this closes a gap between the capability narrative and the quick-glance repo badges at the top of the handbook.",
    ], THEME["jade"], "surface-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Kalier", "embedded-client Kali Linux control for chat and canvas", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What it adds", KALIER_GUIDE, THEME["jade"], "kalier-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "How operators reach it", KALIER_SURFACES_GUIDE, THEME["amber"], "kalier-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "STM32er", "critical-mission STM32 firmware control with dual backends", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What it adds", STM32ER_GUIDE, THEME["copper"], "stm32-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "How operators reach it", STM32ER_SURFACES_GUIDE, THEME["jade"], "stm32-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "ESP32er", "PlatformIO-driven ESP32 firmware control", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What it adds", ESP32ER_GUIDE, THEME["jade"], "esp32-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "How operators reach it", ESP32ER_SURFACES_GUIDE, THEME["amber"], "esp32-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "ESP32 Template Project", "known-good PlatformIO baseline for ESP32er", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Why it matters", ESP32_TEMPLATE_GUIDE, THEME["amber"], "esp32-c", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Operator result", [
        "The reference project gives Tlamatini a stable build/upload/monitor proving ground before a user asks her to generate larger ESP32 firmware.",
        "Because it is a plain PlatformIO repo, it matches the exact grain of `chat_agent_esp32er` and the visual ESP32er node instead of introducing another sidecar protocol.",
        "This closes the gap between the new agent and a practical first project a user can build, flash, and watch over serial on real silicon.",
    ], THEME["jade"], "esp32-d", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "ESPHomer", "ESPHome-driven smart-home firmware control", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What it adds", ESPHOMER_GUIDE, THEME["copper"], "esphome-a", 12)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "How operators reach it", ESPHOMER_SURFACES_GUIDE, THEME["jade"], "esphome-b", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "ESPHome Template Project", "known-good YAML baseline for ESPHomer", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Why it matters", ESPHOME_TEMPLATE_GUIDE, THEME["jade"], "esphome-c", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Operator result", [
        "The bundled sample gives Tlamatini a stable validate/compile/upload proving ground before a user asks her to generate a custom ESPHome device.",
        "Because the source-of-truth is a single YAML file, it matches the exact grain of `chat_agent_esphomer` and the visual ESPHomer node instead of introducing a separate project-server protocol.",
        "This closes the gap between the new agent and a practical first device a user can build, flash, and then control from a smart-home hub.",
    ], THEME["amber"], "esphome-d", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Self-Knowledge And Self-Modify", "who she is and how she can improve herself", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What changed", SELF_KNOWLEDGE_GUIDE, THEME["amber"], "self-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Self-modify and autonomy", SELF_MODIFY_GUIDE + MULTITURN_4096_GUIDE[:2], THEME["jade"], "self-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Windower In Multi-Turn", "desktop window management for chat and canvas", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What it adds", WINDOWER_GUIDE, THEME["amber"], "window-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "How operators reach it", WINDOWER_SURFACES_GUIDE, THEME["jade"], "window-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Playwrighter", "real-browser automation for chat and canvas", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What it adds", PLAYWRIGHTER_GUIDE, THEME["jade"], "play-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "How operators reach it", PLAYWRIGHTER_SURFACES_GUIDE, THEME["amber"], "play-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Reviewer Precision In v1.4.1", "commit-state and secret-handling refinement", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Behavioral accuracy patch", REVIEWER_PRECISION_GUIDE, THEME["jade"], "reviewer-c", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Why it matters", [
        "Local working-copy credentials in managed config files are no longer described as already committed when the diff is still uncommitted or only staged.",
        "Review findings stay stricter on true secrets in source code or outside the managed scrub-path set, so the patch reduces noise without weakening real security findings.",
        "The same rules apply in both the canvas Reviewer agent and the `code-review` skill, keeping the two review surfaces behaviorally aligned.",
    ], THEME["amber"], "reviewer-d", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Native Dialogs In v1.4.2", "Tkinter removed from the unstable runtime path", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What changed", NATIVE_DIALOGS_GUIDE, THEME["amber"], "native-a", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Release meaning", [
        "This stability-focused patch sits immediately before the current `v1.5.0` Playwrighter release and remains part of the current operator/runtime story.",
        "It preserves the operator experience of Browse-driven file and folder picking while removing a UI technology that was destabilizing the application.",
        "Because the fix landed with tests and runtime cleanup updates, it belongs in the technical dossier even though the markdown handbook has not yet been fully rewritten around it.",
    ], THEME["jade"], "native-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Gatewayer And External Signals", "inbound automation boundary", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Ingress role", [
        "Receives HTTP webhook events or optional folder-drop files.",
        "Normalizes, persists, queues, and dispatches events into downstream workflow agents.",
        "Supports bearer/HMAC-style auth patterns, IP allowlists, dedup files, and crash recovery state.",
    ], THEME["amber"], "gate-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Operating note", [
        "Gatewayer is a workflow trigger, not a chat-router clone.",
        "Declared config fields are now documented with clear caveats where enforcement is not complete.",
        "Stable log markers let Monitor Log and Summarizer watch gateway health.",
    ], THEME["jade"], "gate-b", 16)
    audit_layout(audit, len(prs.slides))

    add_themed_column_slides(prs, "How To Use Tlamatini", "operator path", THEME["jade"], [
        ("Daily use", THEME["jade"], HOW_TO_USE[:10]),
        ("Workflows and releases", THEME["copper"], HOW_TO_USE[10:]),
    ], per_column=5)

    slide, audit = add_slide(prs, "Source Mode Bootstrap", "commands that matter", THEME["copper"])
    commands = "\n".join([
        "python -m venv venv",
        "venv\\Scripts\\activate",
        "pip install -r requirements.txt",
        "python Tlamatini/manage.py migrate",
        "python Tlamatini/manage.py createsuperuser",
        "python Tlamatini/manage.py collectstatic --noinput",
        "python Tlamatini/manage.py runserver --noreload",
    ])
    add_text(slide, audit, 0.92, 1.75, 11.35, 3.0, commands, 18, THEME["white"], False, name="commands", font="Cascadia Mono")
    add_panel(slide, audit, 0.92, 5.1, 11.35, 1.18, "First run checklist", [
        "Open the browser, log in with your source-mode superuser, load a context path, and decide whether the task needs Multi-Turn.",
    ], THEME["jade"], "checklist", 15)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Installation And Configuration", "README-backed operator path", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.85, 4.95, "Install essentials", INSTALLATION_GUIDE, THEME["jade"], "install-a", 15)
    add_panel(slide, audit, 6.92, 1.6, 5.55, 4.95, "Config essentials", CONFIGURATION_GUIDE[:4], THEME["copper"], "install-b", 14)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Configuration Essentials (continued)", "README-backed operator path", THEME["jade"])
    add_panel(slide, audit, 0.82, 1.6, 11.55, 4.95, "Config essentials", CONFIGURATION_GUIDE[4:], THEME["copper"], "install-c", 15)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "ACPX-Skills And Prompts", "recent operator-surface documentation updates", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "ACPX-Skills menu", ACPX_SKILLS_GUIDE, THEME["amber"], "skills-a", 14)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Prompt catalog and readability", PROMPT_CATALOG_GUIDE, THEME["jade"], "skills-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "DB Menu And Startup Swap", "operator-facing database maintenance surface", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Backup and Set DB", DB_MENU_GUIDE, THEME["copper"], "db-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "What happens on next start-up", DB_SWAP_GUIDE, THEME["jade"], "db-b", 14)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Versioning System", "release-identity work across runtime and builds", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.58, "SemVer and resolver", VERSIONING_GUIDE, THEME["amber"], "ver-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.58, "Where version appears", VERSION_SURFACES_GUIDE, THEME["jade"], "ver-b", 14)
    add_text(
        slide,
        audit,
        0.92,
        6.34,
        11.1,
        0.22,
        f"Resolved current version: {context['version_info']['version']} | build: {context['version_info']['build']}",
        9,
        THEME["muted"],
        False,
        name="version-foot",
        font="Cascadia Mono",
    )
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "De-Compresser Agent", "archive compression and decompression worker", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Operator contract", DE_COMPRESSER_GUIDE, THEME["copper"], "decomp-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Integration and fallbacks", DE_COMPRESSER_INTEGRATION_GUIDE, THEME["jade"], "decomp-b", 14)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Unreal MCP And Unrealer", "UE5 editor bridge for chat and canvas", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What it adds", UNREAL_MCP_GUIDE, THEME["jade"], "unreal-a", 14)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Install and runtime path", UNREAL_INSTALL_GUIDE + UNREAL_RUNTIME_GUIDE[:1], THEME["copper"], "unreal-b", 13)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Extended Unrealer Surface", "the 53-command fork and why it matters", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Extended fork", UNREAL_EXTENDED_GUIDE, THEME["amber"], "unreal-c", 13)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Runtime consequences", UNREAL_RUNTIME_GUIDE[1:] + ["The seeded Unreal demo prompts now cover screenshots, scene-building, and in-editor Python/introspection paths on top of the original blueprint flow."], THEME["jade"], "unreal-d", 12)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Orphan-Process Cleanup", "Windows process-hygiene and survivor reporting", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Three-tier reaper", ORPHAN_REAPER_GUIDE, THEME["amber"], "reaper-a", 14)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Prevention and survivor reporting", ORPHAN_PREVENTION_GUIDE, THEME["jade"], "reaper-b", 14)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Ollama Without Admin Rights", "local model setup on Windows", THEME["amber"])
    add_text(slide, audit, 0.85, 1.72, 11.55, 3.2, OLLAMA_COMMANDS, 9, THEME["white"], False, name="ollama-commands", font="Cascadia Mono")
    add_panel(slide, audit, 0.85, 5.02, 11.55, 1.9, "Checklist", OLLAMA_GUIDE[:2], THEME["amber"], "ollama-check", 14)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Ollama Readiness", "service, API, and model pulls", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Service and API", OLLAMA_GUIDE[2:], THEME["jade"], "ollama-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Default pull set", [
        "Nomic-Embed-Text:latest",
        "glm-5.2:cloud",
        "qwen3.5:cloud",
        "gpt-oss:120b-cloud",
        "qwen3.5:397b-cloud",
        "glm-5.1:cloud",
    ], THEME["copper"], "ollama-b", 15)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "GPU Guard And Reconnect UX", "README and Book driven operator safeguards", THEME["amber"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Embedding-memory pre-flight guard", EMBEDDING_GUARD_GUIDE, THEME["amber"], "guard-a", 14)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Reconnect and restart safeguards", RECENT_RUNTIME_SAFEGUARDS, THEME["jade"], "guard-b", 14)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Running Modes", "development, MCP bootstrap, and ASGI", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "How to run", RUNNING_GUIDE, THEME["copper"], "run-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "What startup does", [
        "Initializes Django and runtime guards in manage.py.",
        "Cleans pool state and repopulates the Agent table from current disk templates.",
        "Launches MCP metrics and gRPC file-search servers before steady-state traffic.",
        "Self-heals the frozen Windows Installed-apps entry when `Uninstaller.exe` is present beside the executable.",
        "Handles shutdown by killing tracked/untracked agent processes and clearing pool artifacts.",
    ], THEME["jade"], "run-b", 15)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Packaging Path", "source usage is separate from release building", THEME["amber"])
    add_flow_boxes(slide, audit, 1.15, 2.0, ["build.py", "pkg.zip", "build_uninstaller", "Uninstaller", "build_installer", "Release"], THEME["amber"])
    add_panel(slide, audit, 0.92, 3.35, 11.35, 2.55, "Release rule", [
        "Run packaging only when preparing a Windows distribution.",
        "The final distributable is the full `dist/Tlamatini_Release_v<version>/` folder, not one executable copied out of context.",
        "Installer scripts register shortcuts, `.flw` file associations, the bundled uninstaller, and the per-user Installed-apps entry.",
    ], THEME["amber"], "packaging", 16)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Exec Report", "show-your-work visibility for Multi-Turn", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Why it exists", EXEC_REPORT_GUIDE, THEME["jade"], "exec-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Operator effect", [
        "Tables appear only for state-changing agent families that actually fired.",
        "Verdicts are derived from real tool returns, not inferred from prose summaries.",
        "This is the audit surface that makes long jobs debuggable from the chat output itself.",
    ], THEME["amber"], "exec-b", 15)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "ACPX And Skills", "external coding-agent runtime", THEME["copper"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "What ACPX adds", ACPX_GUIDE, THEME["copper"], "acpx-a", 15)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Files and operator model", [
        "`agent/acpx/` hosts the runtime, permission gate, registry, and tools.",
        "`agent/skills/` and `agent/skills_pkg/` host the in-process skill harness and markdown catalog.",
        "Transcripts and audit logs are persisted so external delegation stays replayable.",
    ], THEME["jade"], "acpx-b", 15)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Application Log", "tlamatini.log as forensic truth", THEME["amber"])
    add_panel(slide, audit, 0.85, 1.75, 11.45, 4.6, "How the log works", APP_LOG_GUIDE, THEME["amber"], "log", 17)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Repository Facts", "current head inventory", THEME["jade"])
    metrics = [
        ("Repo files", context["inventory_files"]),
        ("Agents", context["workflow_agent_count"]),
        ("Migrations", context["migrations"]),
        ("JS", context["js_modules"]),
        ("CSS", context["css_files"]),
        ("HTML", context["html_templates"]),
    ]
    for idx, (label, value) in enumerate(metrics):
        add_metric_card(slide, audit, 0.82 + idx * 2.05, 1.75, 1.75, label, str(value), THEME["jade"] if idx % 2 == 0 else THEME["copper"], f"repo-{idx}")
    add_panel(slide, audit, 1.05, 3.1, 10.85, 3.35, "Current HEAD", [
        f"{context['head_short']} - {context['head_subject']}",
        f"Resolved version: {context['version_info']['version']} ({context['version_info']['source']})",
        f"Generated on {context['generated_at']}",
        f"Inventory scope: {context['inventory_files']} files = {context['tracked_files']} tracked + {context['untracked_files']} git-unignored working-tree additions",
        f"Multi-Turn tools: {context['total_multi_turn_tools']}; wrapped chat-agent tools: {context['wrapped_chat_agent_count']}; skills: {context['skills_count']}",
        f"Python requirements: {context['requirements_count']}; authoritative agent-description rows: {context['agent_description_rows']}",
        f"Binary or asset inventory files skipped from line count: {context['binary_count']}",
    ], THEME["amber"], "repo-head", 15)
    audit_layout(audit, len(prs.slides))

    baseline = context["visual_doc_baseline"]
    if baseline is not None:
        visual_highlights = context["visual_doc_highlights"]
        slide, audit = add_slide(prs, "Since Last Dossier Refresh", "all important changes since the last committed PDF/PPTX update", THEME["copper"])
        add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Baseline", [
            f"{baseline.short_hash} on {iso_date(baseline.committed_at)}",
            baseline.subject,
            f"Commits since then: {len(context['visual_doc_commits'])}",
        ], THEME["copper"], "since-a", 13)
        add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Key changes", visual_highlights[:3], THEME["jade"], "since-b", 10)
        audit_layout(audit, len(prs.slides))

        remaining_highlights = visual_highlights[3:]
        for offset in range(0, len(remaining_highlights), 6):
            group = remaining_highlights[offset:offset + 6]
            slide, audit = add_slide(
                prs,
                "Dossier Delta Continued",
                f"verified implementation changes {offset + 4}-{offset + 3 + len(group)}",
                THEME["jade"],
            )
            if len(group) == 1:
                add_panel(slide, audit, 0.78, 1.6, 11.72, 4.95, "Additional verified change", group, THEME["jade"], f"since-more-{offset}", 12)
            else:
                split_at = (len(group) + 1) // 2
                add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Runtime and UI", group[:split_at], THEME["jade"], f"since-more-a-{offset}", 10)
                add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Data and operator contract", group[split_at:], THEME["amber"], f"since-more-b-{offset}", 10)
            audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Recent Platform Additions", "tagged v1.42.0 release lineage", THEME["jade"])
    add_panel(slide, audit, 0.78, 1.6, 5.9, 4.95, "Recent agents and execution surfaces", [
        "External MCP results (v1.41.4): stdio and network clients surface text plus structuredContent, unwrap common envelopes, cap payloads, and preserve errors.",
        "STM32er (v1.42.0): device-aware PlatformIO routing adds Blue Pill/mainstream-family build and safe flash; new stepwise Blue Pill/F407 demos finish with camera evidence.",
        "Prompt catalog (v1.42.0): migration 0179 deliberately regroups and renumbers every row to contiguous 1..N category blocks, backed by four database invariants tests.",
        "Prompt catalog (v1.41.3): 13 categories, 13 duplicate ACPX rows removed, stable surviving ids, gap-tolerant loading, and ranked fuzzy search.",
        "Hard Cancel (v1.41.2): per-user run epochs stop executor/retry/self-healing resurrection while preserving completed tool evidence and the next request.",
        "Screenshot chat, configurable port, FlowPills discovery, and Unreal 5.8 scaffolding remain carried foundations; no new workflow-agent type landed.",
    ], THEME["copper"], "monday-a", 10)
    add_panel(slide, audit, 6.95, 1.6, 5.55, 4.95, "Lifecycle, policy, and monitoring", [
        "Resolved identity: `origin/main`, local `main`, and tag v1.42.0 all point to `c58b01ad`; only local configuration values remain modified and are not reproduced here.",
        "The stronger disclaimer says plain-Python transparency enables user control but is not a security warranty; the operator owns authorization, permissions, review, and consequences.",
        "The dossier distinguishes tagged behavior from local configuration-only changes, preserves private-data discipline, and does not stage, commit, or push.",
        "Operator setup: easy-start install, Ollama guidance, Config dialogs, DB menu, and Windows Installed-apps registration stay in the dossier.",
        f"Catalog now stands at {context['workflow_agent_count']} workflow agents and {context['total_multi_turn_tools']} Multi-Turn tools ({context['wrapped_chat_agent_count']} wrapped chat-agent + {context['acpx_tool_count']} ACPX/Skill + {context['core_python_tool_count']} core), with {context['skills_count']} skills.",
    ], THEME["jade"], "monday-b", 10)
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Effective Lines By Language", "no comments, no blanks", THEME["copper"])
    add_metric_card(slide, audit, 0.9, 1.64, 2.35, "Total effective", f"{context['total_effective_lines']:,}", THEME["copper"], "lines-m1")
    add_metric_card(slide, audit, 3.55, 1.64, 2.35, "Total physical", f"{context['total_lines']:,}", THEME["jade"], "lines-m2")
    add_metric_card(slide, audit, 6.2, 1.64, 2.35, "Text languages", str(len(context["language_rows"])), THEME["amber"], "lines-m3")
    add_text(slide, audit, 0.92, 3.05, 11.35, 3.2, language_table_text(context["language_rows"], context["total_effective_lines"]), 8, THEME["white"], False, name="language-table", font="Cascadia Mono")
    audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "Largest Files By Effective Lines", "where most authored text lives", THEME["jade"])
    add_text(slide, audit, 0.85, 1.72, 11.7, 4.85, file_table_text(context["file_rows"]), 8, THEME["white"], False, name="largest-table", font="Cascadia Mono")
    audit_layout(audit, len(prs.slides))

    recent_highlight_chunks = split_items(context["weekly_highlights"], 5)
    for idx, chunk in enumerate(recent_highlight_chunks, 1):
        slide_title = (
            RECENT_GIT_WINDOW_TITLE
            if len(recent_highlight_chunks) == 1
            else f"{RECENT_GIT_WINDOW_TITLE} ({idx}/{len(recent_highlight_chunks)})"
        )
        slide, audit = add_slide(prs, slide_title, "recent changes according to git history", THEME["amber"])
        add_panel(slide, audit, 0.82, 1.65, 11.55, 4.9, RECENT_GIT_HIGHLIGHT_TITLE, chunk, THEME["amber"], f"latest-{idx}", 15)
        audit_layout(audit, len(prs.slides))

    visual_chunks = split_items(context["visual_doc_commits"], 6)
    for idx, chunk in enumerate(visual_chunks, 1):
        slide, audit = add_slide(
            prs,
            f"Visual Dossier Change Appendix {idx}/{len(visual_chunks)}",
            "commits since the last committed PDF/PPTX refresh",
            THEME["jade"] if idx % 2 else THEME["copper"],
        )
        visual_lines = [f"{iso_date(c.committed_at)} | {c.short_hash} | {c.subject}" for c in chunk]
        add_panel(slide, audit, 0.82, 1.68, 11.55, 4.86, "Commit timeline", visual_lines, THEME["jade"] if idx % 2 else THEME["copper"], f"visual-{idx}", 12)
        audit_layout(audit, len(prs.slides))

    weekly_chunks = split_items(context["weekly_commits"], 6)
    for idx, chunk in enumerate(weekly_chunks, 1):
        slide, audit = add_slide(
            prs,
            (
                f"Today's Commit Appendix {idx}/{len(weekly_chunks)}"
                if RECENT_GIT_WINDOW_LABEL == "today"
                else f"{RECENT_GIT_WINDOW_DAYS}-Day Commit Appendix {idx}/{len(weekly_chunks)}"
            ),
            RECENT_GIT_APPENDIX_SUBTITLE,
            THEME["copper"] if idx % 2 else THEME["jade"],
        )
        weekly_lines = [f"{iso_date(c.committed_at)} | {c.short_hash} | {c.subject}" for c in chunk]
        add_panel(slide, audit, 0.82, 1.68, 11.55, 4.86, "Commit timeline", weekly_lines, THEME["copper"] if idx % 2 else THEME["jade"], f"week-{idx}", 12)
        audit_layout(audit, len(prs.slides))

    tree_chunks = split_lines(context["tree_text"], 31)
    for idx, chunk in enumerate(tree_chunks, 1):
        slide, audit = add_slide(
            prs,
            f"Repository File Tree Appendix {idx}/{len(tree_chunks)}",
            "complete repository inventory tree, including git-unignored working additions",
            THEME["jade"] if idx % 2 else THEME["copper"],
        )
        add_text(slide, audit, 0.72, 1.56, 11.95, 5.42, chunk, 7, THEME["white"], False, name=f"tree-{idx}", font="Cascadia Mono")
        audit_layout(audit, len(prs.slides))

    slide, audit = add_slide(prs, "How To Keep Docs Excellent", "future refresh discipline", THEME["copper"])
    add_panel(slide, audit, 0.85, 1.75, 11.55, 4.6, "Recommended practice", [
        "Regenerate the PDF and deck whenever README, architecture, agent catalog, line inventory, or packaging behavior changes.",
        "Keep the PDF exhaustive and evidence-heavy; keep the PPT visual, split dense appendices, and audit geometry before delivery.",
        "Use the new Skills in `.codex/skills/` so future refreshes follow the same no-overlap and full-dossier rules.",
    ], THEME["copper"], "final", 17)
    audit_layout(audit, len(prs.slides))

    prs.save(PPT_OUTPUT)


def serialize_context(context: dict) -> dict:
    return {
        "generated_at": context["generated_at"],
        "head_short": context["head_short"],
        "head_full": context["head_full"],
        "head_subject": context["head_subject"],
        "head_date": context["head_date"],
        "inventory_files": context["inventory_files"],
        "tracked_files": context["tracked_files"],
        "untracked_files": context["untracked_files"],
        "total_effective_lines": context["total_effective_lines"],
        "total_lines": context["total_lines"],
        "workflow_agent_count": context["workflow_agent_count"],
        "agent_description_rows": context["agent_description_rows"],
        "wrapped_chat_agent_count": context["wrapped_chat_agent_count"],
        "core_python_tool_count": context["core_python_tool_count"],
        "acpx_tool_count": context["acpx_tool_count"],
        "total_multi_turn_tools": context["total_multi_turn_tools"],
        "skills_count": context["skills_count"],
        "requirements_count": context["requirements_count"],
        "js_modules": context["js_modules"],
        "css_files": context["css_files"],
        "html_templates": context["html_templates"],
        "migrations": context["migrations"],
        "binary_count": context["binary_count"],
        "skipped_count": context["skipped_count"],
        "version": context["version_info"]["version"],
        "version_source": context["version_info"]["source"],
        "version_info": context["version_info"],
        "language_rows": [row.__dict__ for row in context["language_rows"]],
        "largest_files": [row.__dict__ for row in context["file_rows"][:50]],
        "recent_commits": [row.__dict__ for row in context["recent_commits"]],
        "weekly_commits": [row.__dict__ for row in context["weekly_commits"]],
        "weekly_highlights": context["weekly_highlights"],
        "visual_doc_baseline": (
            None
            if context["visual_doc_baseline"] is None
            else context["visual_doc_baseline"].__dict__
        ),
        "visual_doc_commits": [row.__dict__ for row in context["visual_doc_commits"]],
        "visual_doc_highlights": context["visual_doc_highlights"],
        "workflow_agents": context["workflow_agents"],
    }


def main() -> None:
    BUILD_DIR.mkdir(parents=True, exist_ok=True)
    context = collect_context()
    build_pdf(context)
    build_ppt(context)
    CONTEXT_OUTPUT.write_text(json.dumps(serialize_context(context), indent=2), encoding="utf-8")
    print(f"Updated PDF: {PDF_OUTPUT}")
    print(f"Updated PPTX: {PPT_OUTPUT}")
    print(f"Wrote context: {CONTEXT_OUTPUT}")
    print(f"Wrote tree: {TREE_OUTPUT}")


if __name__ == "__main__":
    main()
